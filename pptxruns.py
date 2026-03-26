import streamlit as st
from concurrent.futures import ThreadPoolExecutor, as_completed
import json
import os
import tempfile
import io
import time
import hashlib
import hmac
import random
import base64
import requests
import pandas as pd
import pypdf
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ─────────────────────────────────────────────────────────────
# THEME & FILL HELPERS — prevent white-on-white text
# ─────────────────────────────────────────────────────────────

# Dark SEGA theme: dk1=white, lt1=dark-navy.
# This means any text that falls back to theme colour is white,
# which is correct on our dark backgrounds.  Replaces the default
# Office theme (dk1=windowText which can resolve to white on some
# systems/dark-mode setups, making text invisible).
_SEGA_THEME_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="SEGA Dark">'
    '<a:themeElements>'
    '<a:clrScheme name="SEGA Dark">'
    '<a:dk1><a:srgbClr val="FFFFFF"/></a:dk1>'
    '<a:lt1><a:srgbClr val="040A1C"/></a:lt1>'
    '<a:dk2><a:srgbClr val="0055AA"/></a:dk2>'
    '<a:lt2><a:srgbClr val="1A2A4A"/></a:lt2>'
    '<a:accent1><a:srgbClr val="00AADD"/></a:accent1>'
    '<a:accent2><a:srgbClr val="F5C218"/></a:accent2>'
    '<a:accent3><a:srgbClr val="00BB66"/></a:accent3>'
    '<a:accent4><a:srgbClr val="CC2244"/></a:accent4>'
    '<a:accent5><a:srgbClr val="8899BB"/></a:accent5>'
    '<a:accent6><a:srgbClr val="AABBCC"/></a:accent6>'
    '<a:hlink><a:srgbClr val="00AADD"/></a:hlink>'
    '<a:folHlink><a:srgbClr val="0055AA"/></a:folHlink>'
    '</a:clrScheme>'
    '<a:fontScheme name="SEGA">'
    '<a:majorFont><a:latin typeface="Calibri"/><a:ea typeface=""/><a:cs typeface=""/></a:majorFont>'
    '<a:minorFont><a:latin typeface="Calibri"/><a:ea typeface=""/><a:cs typeface=""/></a:minorFont>'
    '</a:fontScheme>'
    '<a:fmtScheme name="Office">'
    '<a:fillStyleLst>'
    '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
    '<a:gradFill rotWithShape="1"><a:gsLst>'
    '<a:gs pos="0"><a:schemeClr val="phClr"><a:tint val="50000"/></a:schemeClr></a:gs>'
    '<a:gs pos="100000"><a:schemeClr val="phClr"><a:tint val="15000"/></a:schemeClr></a:gs>'
    '</a:gsLst><a:lin ang="16200000" scaled="1"/></a:gradFill>'
    '<a:gradFill rotWithShape="1"><a:gsLst>'
    '<a:gs pos="0"><a:schemeClr val="phClr"><a:shade val="51000"/></a:schemeClr></a:gs>'
    '<a:gs pos="100000"><a:schemeClr val="phClr"><a:shade val="94000"/></a:schemeClr></a:gs>'
    '</a:gsLst><a:lin ang="16200000" scaled="0"/></a:gradFill>'
    '</a:fillStyleLst>'
    '<a:lnStyleLst>'
    '<a:ln w="9525"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill><a:prstDash val="solid"/></a:ln>'
    '<a:ln w="25400"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill><a:prstDash val="solid"/></a:ln>'
    '<a:ln w="38100"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill><a:prstDash val="solid"/></a:ln>'
    '</a:lnStyleLst>'
    '<a:effectStyleLst>'
    '<a:effectStyle><a:effectLst/></a:effectStyle>'
    '<a:effectStyle><a:effectLst/></a:effectStyle>'
    '<a:effectStyle><a:effectLst/></a:effectStyle>'
    '</a:effectStyleLst>'
    '<a:bgFillStyleLst>'
    '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
    '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
    '<a:gradFill rotWithShape="1"><a:gsLst>'
    '<a:gs pos="0"><a:schemeClr val="phClr"><a:tint val="93000"/></a:schemeClr></a:gs>'
    '<a:gs pos="100000"><a:schemeClr val="phClr"><a:shade val="98000"/></a:schemeClr></a:gs>'
    '</a:gsLst><a:lin ang="16200000" scaled="1"/></a:gradFill>'
    '</a:bgFillStyleLst>'
    '</a:fmtScheme>'
    '</a:themeElements>'
    '</a:theme>'
)


def _patch_theme(pptx_bytes: bytes) -> bytes:
    """Swap every theme file in the PPTX zip for the SEGA dark theme."""
    import zipfile, re as _re
    buf_in  = io.BytesIO(pptx_bytes)
    buf_out = io.BytesIO()
    theme_bytes = _SEGA_THEME_XML.encode("utf-8")
    with zipfile.ZipFile(buf_in, "r") as zin,          zipfile.ZipFile(buf_out, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if _re.match(r"ppt/theme/theme\d+\.xml", item.filename):
                data = theme_bytes
            zout.writestr(item, data)
    return buf_out.getvalue()


def _lock_txb(txb):
    """
    Add explicit noFill + noLine to a textbox shape element so it can never
    inherit a white background from the theme or slide layout.
    Must be called after add_textbox() but before saving.
    """
    sp   = txb._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    # Remove any existing fill child
    for tag in ("a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill", "a:grpFill", "a:noFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    # Insert <a:noFill/> after the <a:xfrm> and <a:prstGeom> elements
    no_fill = etree.SubElement(spPr, qn("a:noFill"))
    # Also kill any line (border) that might show as white
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    # Clear the line's fill children and set noFill
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, qn("a:noFill"))


# Safe defaults and constants used inside _render_plan_modal.
run_btn           = False
business_question = ""
SLIDE_TYPES = ["title","section","bullets","stats","comparison","recommendation","chart","closing"]



def _send_otp(email: str, code: str) -> bool:
    """Send OTP via AWS SES using credentials from st.secrets."""
    try:
        import boto3
        from botocore.exceptions import ClientError
        ses = boto3.client(
            "ses",
            region_name=st.secrets.get("AWS_SES_REGION", "us-east-1"),
            aws_access_key_id=st.secrets.get("AWS_ACCESS_KEY_ID", ""),
            aws_secret_access_key=st.secrets.get("AWS_SECRET_ACCESS_KEY", ""),
        )
        ses.send_email(
            Source=st.secrets.get("EMAIL_FROM", "noreply@segaamerica.com"),
            Destination={"ToAddresses": [email]},
            Message={
                "Subject": {
                    "Data": "SEGA PowerPoint Creator — Your verification code",
                    "Charset": "UTF-8",
                },
                "Body": {
                    "Text": {
                        "Data": (
                            f"Your SEGA PowerPoint Creator verification code is: {code}\n\n"
                            "This code expires in 10 minutes.\n"
                            "If you didn't request this, you can safely ignore this email."
                        ),
                        "Charset": "UTF-8",
                    },
                    "Html": {
                        "Data": f"""
                        <div style="font-family:Arial,sans-serif;max-width:480px;margin:0 auto;padding:32px 24px;">
                          <div style="font-size:22px;font-weight:900;letter-spacing:0.1em;color:#1A6BFF;margin-bottom:4px;">SEGA</div>
                          <div style="font-size:14px;color:#444;margin-bottom:28px;">PowerPoint Creator</div>
                          <div style="font-size:14px;color:#222;margin-bottom:16px;">Your verification code is:</div>
                          <div style="font-size:42px;font-weight:900;letter-spacing:0.18em;color:#1a1a2e;
                                      background:#f0f4ff;border-radius:8px;padding:18px 24px;
                                      display:inline-block;margin-bottom:24px;">{code}</div>
                          <div style="font-size:12px;color:#888;">
                            This code expires in 10 minutes.<br>
                            If you didn't request this, you can safely ignore this email.
                          </div>
                        </div>""",
                        "Charset": "UTF-8",
                    },
                },
            },
        )
        return True
    except ClientError as e:
        st.error(f"SES error: {e.response['Error']['Message']}")
        return False
    except Exception as e:
        st.error(f"Failed to send email: {e}")
        return False
def _sign_cookie(email: str) -> str:
    """Create an HMAC-signed token: email|expiry|signature."""
    secret = st.secrets.get("COOKIE_SIGNING_KEY", "fallback-change-this")
    expiry = int(time.time()) + (COOKIE_EXPIRY_DAYS * 86400)
    payload = f"{email}|{expiry}"
    sig = hmac.new(secret.encode(), payload.encode(), hashlib.sha256).hexdigest()
    return base64.urlsafe_b64encode(f"{payload}|{sig}".encode()).decode()
def _verify_cookie(token: str) -> str | None:
    """Verify signed cookie. Returns email if valid, None otherwise."""
    try:
        secret = st.secrets.get("COOKIE_SIGNING_KEY", "fallback-change-this")
        decoded = base64.urlsafe_b64decode(token.encode()).decode()
        email, expiry_str, sig = decoded.rsplit("|", 2)
        payload = f"{email}|{expiry_str}"
        expected = hmac.new(secret.encode(), payload.encode(), hashlib.sha256).hexdigest()
        if not hmac.compare_digest(sig, expected):
            return None
        if int(time.time()) > int(expiry_str):
            return None
        return email
    except Exception:
        return None
def extract_text_from_file(uploaded_file) -> str:
    name = uploaded_file.name.lower()
    content = ""
    try:
        if name.endswith(".pdf"):
            reader = pypdf.PdfReader(io.BytesIO(uploaded_file.read()))
            content = "\n\n".join(p.extract_text() or "" for p in reader.pages)
        elif name.endswith((".xlsx", ".xls")):
            sheets = pd.read_excel(io.BytesIO(uploaded_file.read()), sheet_name=None)
            content = "\n\n".join(f"=== Sheet: {s} ===\n{df.to_string(max_rows=200)}" for s, df in sheets.items())
        elif name.endswith(".csv"):
            content = pd.read_csv(io.BytesIO(uploaded_file.read())).to_string(max_rows=300)
        elif name.endswith(".txt"):
            content = uploaded_file.read().decode("utf-8", errors="replace")
        elif name.endswith(".docx"):
            import zipfile, xml.etree.ElementTree as ET
            zf   = zipfile.ZipFile(io.BytesIO(uploaded_file.read()))
            root = ET.fromstring(zf.read("word/document.xml"))
            ns   = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
            content = "\n".join(
                "".join(r.text for r in p.iter(f"{ns}t") if r.text)
                for p in root.iter(f"{ns}p")
            )
        else:
            content = f"[Unsupported type: {name}]"
    except Exception as e:
        content = f"[Extraction error for {uploaded_file.name}: {e}]"
    if len(content) > 15000:
        content = content[:15000] + "\n\n[... truncated ...]"
    return content
def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
def _in(v):  return Inches(v)
def _pt(v):  return Pt(v)
def _rect(slide, x, y, w, h, fill_hex, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        _in(x), _in(y), _in(w), _in(h)
    )
    shape.line.fill.background()
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_hex)
    if alpha is not None:
        # alpha 0-100 (0=opaque, 100=transparent in pptx land)
        shape.fill.fore_color._element.getparent().getparent()  # noop touch
    shape.line.fill.background()
    return shape
def _add_text(slide, text, x, y, w, h,
              size=12, bold=False, italic=False, color="FFFFFF",
              align=PP_ALIGN.LEFT, wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    _lock_txb(txb)  # prevent white-box inheritance
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size  = _pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    run.font.name  = font_name
    return txb
def _chrome(slide, idx, total, C):
    """Top stripe, bottom bar, page number, right edge accent."""
    _rect(slide, 0, 0,       W_IN, 0.1,  C["accent"])
    _rect(slide, 0, H_IN-0.36, W_IN, 0.36, C["header_bg"])
    _add_text(slide, "SEGA POWERPOINT CREATOR",
              0.3, H_IN-0.30, 5, 0.24,
              size=7, bold=True, color=C["midgray"],
              font_name=C["body_font"])
    _add_text(slide, f"{idx} / {total}",
              W_IN-1.4, H_IN-0.30, 1.1, 0.24,
              size=8, color=C["neutral"], align=PP_ALIGN.RIGHT,
              font_name=C["body_font"])
    _rect(slide, W_IN-0.1, 0, 0.1, H_IN, C["primary"])
def _set_bg(slide, hex_color):
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)
def _add_bullets(slide, bullets, x, y, w, h, bullet_color, text_color,
                 size=13, font_name=None):
    if not bullets:
        return
    font_name = font_name or "Calibri"   # resolved at call time, not definition time
    txb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    _lock_txb(txb)  # prevent white-box inheritance
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        # bullet marker
        r1 = p.add_run()
        r1.text = "▸ "
        r1.font.color.rgb = _rgb(bullet_color)
        r1.font.bold  = True
        r1.font.size  = _pt(size)
        r1.font.name  = font_name
        # bullet text
        r2 = p.add_run()
        r2.text = b
        r2.font.color.rgb = _rgb(text_color)
        r2.font.size  = _pt(size)
        r2.font.name  = font_name
def _slide_title(slide, s, idx, total, C):
    _set_bg(slide, C["bg"])
    _rect(slide, 7.4, 0, 5.9, H_IN, C["primary"])
    # SEGA wordmark
    _add_text(slide, "SEGA", 8.0, 0.55, 4.8, 0.9,
              size=52, bold=True, color=C["white"], font_name="Arial Black",
              align=PP_ALIGN.CENTER)
    _rect(slide, 8.3, 1.52, 4.2, 0.05, C["accent"])
    _add_text(slide, "POWERPOINT CREATOR", 8.0, 1.58, 4.8, 0.38,
              size=9, color=C["accent"], align=PP_ALIGN.CENTER, font_name=C["body_font"])
    # Main title
    _add_text(slide, s.get("title",""), 0.6, 1.8, 6.8, 2.2,
              size=34, bold=True, color=C["white"], font_name=C["body_font"])
    if s.get("subtitle"):
        _add_text(slide, s["subtitle"], 0.6, 4.2, 6.6, 0.6,
                  size=16, italic=True, color=C["accent"], font_name=C["body_font"])
    if s.get("body"):
        _add_text(slide, s["body"], 0.6, 4.95, 6.6, 1.4,
                  size=12, color=C["midgray"], font_name=C["body_font"])
    _chrome(slide, idx, total, C)
def _slide_section(slide, s, idx, total, C):
    _set_bg(slide, C["bg"])
    _rect(slide, 0, 0, 0.28, H_IN, C["accent"])
    _add_text(slide, s.get("title",""), 0.55, 2.3, 12.0, 1.6,
              size=40, bold=True, color=C["white"], font_name=C["heading_font"])
    if s.get("subtitle"):
        _add_text(slide, s["subtitle"], 0.55, 4.1, 11.0, 0.65,
                  size=19, italic=True, color=C["accent"], font_name=C["body_font"])
    _chrome(slide, idx, total, C)
def _slide_bullets(slide, s, idx, total, C):
    _set_bg(slide, C["bg"])
    _rect(slide, 0, 0.1, W_IN, 0.95, C["subtle"])
    _rect(slide, 0, 0.1, 0.17, 0.95, C["accent"])
    _add_text(slide, s.get("title",""), 0.38, 0.14, W_IN-0.6, 0.84,
              size=22, bold=True, color=C["white"], font_name=C["body_font"])
    _add_bullets(slide, s.get("bullets",[]),
                 0.45, 1.2, W_IN-1.0, H_IN-2.0,
                 bullet_color=C["accent"], text_color=C["white"],
                 size=13)
    if s.get("body"):
        _add_text(slide, s["body"], 0.45, H_IN-1.45, W_IN-1.0, 1.0,
                  size=11, italic=True, color=C["midgray"], font_name=C["body_font"])
    _chrome(slide, idx, total, C)
def _slide_stats(slide, s, idx, total, C):
    _set_bg(slide, C["bg"])
    _rect(slide, 0, 0.1, W_IN, 0.95, C["subtle"])
    _rect(slide, 0, 0.1, 0.17, 0.95, C["gold"])
    _add_text(slide, s.get("title",""), 0.38, 0.14, W_IN-0.6, 0.84,
              size=22, bold=True, color=C["white"], font_name=C["body_font"])
    stats = s.get("stats", [])[:4]
    if stats:
        box_w = (W_IN - 1.2) / len(stats)
        for i, stat in enumerate(stats):
            x = 0.6 + i * (box_w + 0.1)
            _rect(slide, x, 1.3, box_w, 2.9, C["subtle"])
            _add_text(slide, stat.get("value","—"),
                      x+0.1, 1.55, box_w-0.2, 1.2,
                      size=36, bold=True, color=C["accent"],
                      align=PP_ALIGN.CENTER, font_name=C["body_font"])
            _add_text(slide, stat.get("label",""),
                      x+0.1, 2.85, box_w-0.2, 0.7,
                      size=12, bold=True, color=C["white"],
                      align=PP_ALIGN.CENTER, font_name=C["body_font"])
            if stat.get("note"):
                _add_text(slide, stat["note"],
                          x+0.1, 3.6, box_w-0.2, 0.45,
                          size=9, color=C["midgray"],
                          align=PP_ALIGN.CENTER, font_name=C["body_font"])
    if s.get("body"):
        _add_text(slide, s["body"], 0.45, H_IN-1.5, W_IN-1.0, 1.0,
                  size=11, italic=True, color=C["midgray"], font_name=C["body_font"])
    _chrome(slide, idx, total, C)
def _slide_comparison(slide, s, idx, total, C):
    _set_bg(slide, C["bg"])
    _rect(slide, 0, 0.1, W_IN, 0.95, C["subtle"])
    _rect(slide, 0, 0.1, 0.17, 0.95, C["accent"])
    _add_text(slide, s.get("title",""), 0.38, 0.14, W_IN-0.6, 0.84,
              size=20, bold=True, color=C["white"], font_name=C["body_font"])
    cmp    = s.get("comparison", {})
    rows   = cmp.get("rows", [])
    label_w = 2.8
    col_w   = (W_IN - 1.0 - label_w) / 2
    left_x  = 0.5
    mid_x   = left_x + label_w + 0.05
    right_x = mid_x + col_w + 0.05
    start_y = 1.22
    row_h   = 0.42
    # Column headers
    _rect(slide, mid_x,   start_y, col_w, row_h, C["primary"])
    _rect(slide, right_x, start_y, col_w, row_h, C["gold"])
    _add_text(slide, cmp.get("left_title","Internal"),
              mid_x+0.05, start_y, col_w-0.1, row_h,
              size=10, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    _add_text(slide, cmp.get("right_title","Reference"),
              right_x+0.05, start_y, col_w-0.1, row_h,
              size=10, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows[:10]):
        y   = start_y + row_h + ri * row_h
        row_bg = "0D1530" if ri % 2 == 0 else "111E3A"
        total_w = label_w + col_w * 2 + 0.1
        _rect(slide, left_x, y, total_w, row_h - 0.04, row_bg)
        _add_text(slide, row.get("label",""),
                  left_x+0.08, y, label_w-0.12, row_h,
                  size=9, bold=True, color=C["midgray"], font_name=C["body_font"])
        _add_text(slide, row.get("left","—"),
                  mid_x+0.05, y, col_w-0.1, row_h,
                  size=9, color=C["light"], align=PP_ALIGN.CENTER, font_name=C["body_font"])
        delta = (row.get("delta","") or "").lower()
        dc = C["green"] if delta == "positive" else C["red"] if delta == "negative" else C["neutral"]
        _add_text(slide, row.get("right","—"),
                  right_x+0.05, y, col_w-0.1, row_h,
                  size=9, color=dc, align=PP_ALIGN.CENTER, font_name=C["body_font"])
    _chrome(slide, idx, total, C)
def _slide_recommendation(slide, s, idx, total, C):
    _set_bg(slide, C["bg"])
    _rect(slide, 0, 0.1, W_IN, 0.95, C["subtle"])
    _rect(slide, 0, 0.1, 0.17, 0.95, C["green"])
    _add_text(slide, s.get("title",""), 0.38, 0.14, W_IN-0.6, 0.84,
              size=22, bold=True, color=C["white"], font_name=C["body_font"])
    for i, b in enumerate((s.get("bullets") or [])[:6]):
        y = 1.3 + i * 0.88
        _rect(slide, 0.45, y, W_IN-1.0, 0.76, C["subtle"])
        # Numbered circle (just a coloured rect behind the number)
        _rect(slide, 0.55, y+0.1, 0.5, 0.5, C["green"])
        _add_text(slide, str(i+1), 0.55, y+0.1, 0.5, 0.5,
                  size=14, bold=True, color="000000",
                  align=PP_ALIGN.CENTER, font_name=C["body_font"])
        _add_text(slide, b, 1.18, y+0.1, W_IN-1.8, 0.56,
                  size=12, color=C["white"], font_name=C["body_font"])
    _chrome(slide, idx, total, C)
def _slide_closing(slide, s, idx, total, C):
    _set_bg(slide, C["bg"])
    _rect(slide, 0, H_IN/2 - 0.05, W_IN, 0.1,  C["accent"])
    _rect(slide, 0, 0,              0.28, H_IN,  C["accent"])
    _add_text(slide, s.get("title",""), 0.55, 1.8, 12.0, 1.6,
              size=38, bold=True, color=C["white"], font_name=C["body_font"])
    if s.get("subtitle"):
        _add_text(slide, s["subtitle"], 0.55, 3.6, 11.0, 0.65,
                  size=17, italic=True, color=C["accent"], font_name=C["body_font"])
    if s.get("body"):
        _add_text(slide, s["body"], 0.55, 4.4, 11.0, 1.5,
                  size=12, color=C["midgray"], font_name=C["body_font"])
    _add_text(slide, "SEGA  •  CONFIDENTIAL",
              0.55, H_IN-0.75, 8, 0.38,
              size=9, bold=True, color=C["midgray"], font_name=C["body_font"])
    _chrome(slide, idx, total, C)
def _darken(hex6: str, amount: float = 0.15) -> str:
    h = hex6.lstrip('#').upper().ljust(6,'0')[:6]
    r,g,b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
    return f"{max(0,int(r*(1-amount))):02X}{max(0,int(g*(1-amount))):02X}{max(0,int(b*(1-amount))):02X}"
def _lighten(hex6: str, amount: float = 0.1) -> str:
    h = hex6.lstrip('#').upper().ljust(6,'0')[:6]
    r,g,b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
    return f"{min(255,int(r+(255-r)*amount)):02X}{min(255,int(g+(255-g)*amount)):02X}{min(255,int(b+(255-b)*amount)):02X}"
def extract_template_palette(pptx_bytes: bytes) -> dict:
    """
    Read a .pptx template and return a palette dict compatible with
    the generate_pptx() C dict.  Pulls colours and fonts from the
    theme XML and slide master, then derives text colours based on
    whether the background is dark or light.
    """
    import zipfile as _zf, xml.etree.ElementTree as ET

    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

    def _lum(h):
        h = (h or '').lstrip('#').upper()
        if len(h) != 6: return 0.5
        r,g,b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
        return (0.299*r + 0.587*g + 0.114*b) / 255

    def _safe(h):
        if not h: return None
        h = h.lstrip('#').upper()
        return h if len(h) == 6 and all(c in '0123456789ABCDEF' for c in h) else None

    zf = _zf.ZipFile(io.BytesIO(pptx_bytes))
    colors, fonts = {}, {}

    # ── Theme XML ────────────────────────────────────────────────────────
    theme_files = [n for n in zf.namelist() if 'ppt/theme/theme' in n and n.endswith('.xml')]
    if theme_files:
        tree = ET.fromstring(zf.read(theme_files[0]))
        clr  = tree.find('.//a:clrScheme', ns)
        if clr is not None:
            for child in clr:
                tag = child.tag.split('}')[-1]
                for cel in child:
                    ctag = cel.tag.split('}')[-1]
                    val  = cel.get('val') or cel.get('lastClr')
                    if val:
                        if ctag == 'srgbClr':
                            colors[tag] = _safe(val)
                        elif ctag == 'sysClr':
                            colors[tag] = 'FFFFFF' if 'window' in val.lower() else '000000'

        fscheme = tree.find('.//a:fontScheme', ns)
        if fscheme is not None:
            for child in fscheme:
                tag   = child.tag.split('}')[-1]
                latin = child.find('a:latin', ns)
                if latin is not None:
                    tf = (latin.get('typeface') or '').strip()
                    if tf and not tf.startswith('+'):
                        fonts['major' if 'major' in tag.lower() else 'minor'] = tf

    # ── Slide master background ──────────────────────────────────────────
    master_bg = None
    master_files = [n for n in zf.namelist() if 'slideMasters/slideMaster' in n and n.endswith('.xml')]
    if master_files:
        mtree = ET.fromstring(zf.read(master_files[0]))
        for solidFill in mtree.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill'):
            for srgb in solidFill:
                val = _safe(srgb.get('val') or '')
                if val:
                    master_bg = val
                    break
            if master_bg:
                break

    # ── Derive palette ───────────────────────────────────────────────────
    bg      = _safe(master_bg or colors.get('lt1') or colors.get('dk1') or '') or '040A1C'
    is_dark = _lum(bg) < 0.4

    primary = _safe(colors.get('dk2') or colors.get('accent1') or '') or '0055AA'
    # Pick accent that differs from primary
    accent  = _safe(colors.get('accent1') if colors.get('accent1') != primary else
                    colors.get('accent5') or '') or '00AADD'

    if is_dark:
        white, light = 'FFFFFF', 'D0E8FF'
        midgray      = '8899BB'
        subtle       = _darken(bg, 0.12) if _lum(_darken(bg, 0.12)) > 0.05 else _lighten(bg, 0.12)
        header_bg    = _darken(primary, 0.3)
    else:
        white, light = '111122', '334466'
        midgray      = '556677'
        subtle       = _lighten(bg, 0.06)
        header_bg    = primary

    heading_font = fonts.get('major', 'Calibri')
    body_font    = fonts.get('minor', 'Calibri')

    return {
        'bg': bg, 'subtle': subtle, 'header_bg': header_bg,
        'white': white, 'light': light, 'midgray': midgray, 'neutral': midgray,
        'gold':  _safe(colors.get('accent2') or '') or 'F5C242',
        'green': '22DD88', 'red': 'EE3355', 'dark': '111122',
        'primary': primary, 'accent': accent,
        'heading_font': heading_font, 'body_font': body_font,
        'is_dark': is_dark,
    }
def _render_chart_png(chart_spec: dict, width_px=860, height_px=480) -> bytes:
    """
    Render a chart spec to PNG bytes using matplotlib.
    chart_spec keys:
      chart_type : bar | line | scatter | pie | horizontal_bar
      title      : str
      x_label    : str
      y_label    : str
      series     : [{"label": str, "values": [num, ...]}]
      categories : [str, ...]   (x-axis labels / pie slices)
      colors     : [hex, ...]   optional
    """
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.ticker as mticker
    import numpy as np

    spec        = chart_spec or {}
    ctype       = spec.get("chart_type", "bar").lower()
    title       = spec.get("title", "")
    x_label     = spec.get("x_label", "")
    y_label     = spec.get("y_label", "")
    categories  = spec.get("categories") or []
    series      = spec.get("series") or []
    colors      = spec.get("colors") or ["#003D7A","#ED7D31","#22AA66","#CC2244","#8B5CF6","#0EA5E9"]

    # Normalize colors — ensure they have # prefix for matplotlib
    colors = [(c if c.startswith('#') else f'#{c}') for c in colors]

    dpi = 120
    fw  = width_px  / dpi
    fh  = height_px / dpi

    # SEGA-branded style
    bg      = "#FFFFFF"
    grid_c  = "#E2E8F0"
    text_c  = "#1E293B"
    plt.rcParams.update({
        "font.family":     "DejaVu Sans",
        "font.size":       9,
        "axes.labelcolor": text_c,
        "xtick.color":     text_c,
        "ytick.color":     text_c,
        "text.color":      text_c,
        "axes.edgecolor":  grid_c,
        "axes.facecolor":  bg,
        "figure.facecolor":bg,
    })

    fig, ax = plt.subplots(figsize=(fw, fh), dpi=dpi)
    ax.set_facecolor(bg)
    ax.grid(axis="y", color=grid_c, linewidth=0.6, zorder=0)
    ax.set_axisbelow(True)
    for spine in ax.spines.values():
        spine.set_color(grid_c)

    if ctype == "pie":
        vals   = series[0]["values"] if series else []
        labels = categories or [s["label"] for s in series]
        wedge_colors = [colors[i % len(colors)] for i in range(len(vals))]
        ax.pie(vals, labels=labels, autopct="%1.1f%%", colors=wedge_colors,
               textprops={"fontsize": 8}, startangle=140)
        ax.axis("equal")

    elif ctype in ("line", "scatter"):
        x = list(range(len(categories))) if categories else []
        for si, s in enumerate(series):
            vals = [float(v) for v in (s.get("values") or [])]
            col  = colors[si % len(colors)]
            xs   = x[:len(vals)] if x else list(range(len(vals)))
            if ctype == "line":
                ax.plot(xs, vals, marker="o", color=col, linewidth=2,
                        markersize=5, label=s.get("label",""))
            else:
                ax.scatter(xs, vals, color=col, s=40, label=s.get("label",""))
        if categories:
            ax.set_xticks(x)
            ax.set_xticklabels(categories, rotation=25, ha="right", fontsize=8)
        if len(series) > 1:
            ax.legend(fontsize=8, framealpha=0.7)
        if x_label: ax.set_xlabel(x_label, fontsize=8)
        if y_label: ax.set_ylabel(y_label, fontsize=8)

    elif ctype == "horizontal_bar":
        cats = categories or [s["label"] for s in series]
        vals = series[0]["values"] if series else [s.get("values",[0])[0] for s in series]
        vals = [float(v) for v in vals]
        y    = np.arange(len(cats))
        bar_colors = [colors[i % len(colors)] for i in range(len(cats))]
        bars = ax.barh(y, vals, color=bar_colors, height=0.6, zorder=3)
        ax.set_yticks(y)
        ax.set_yticklabels(cats, fontsize=8)
        ax.bar_label(bars, fmt="%.1f", padding=4, fontsize=7.5, color=text_c)
        ax.invert_yaxis()
        ax.grid(axis="x", color=grid_c, linewidth=0.6, zorder=0)
        if x_label: ax.set_xlabel(x_label, fontsize=8)

    else:  # default: grouped bar
        cats  = categories or []
        n_ser = len(series)
        x     = np.arange(len(cats)) if cats else np.arange(
                    max(len(s.get("values", [])) for s in series) if series else 1)
        width = min(0.7 / max(n_ser, 1), 0.35)
        offsets = np.linspace(-(n_ser-1)*width/2, (n_ser-1)*width/2, n_ser) if n_ser > 1 else [0]
        for si, s in enumerate(series):
            vals = [float(v) for v in (s.get("values") or [])]
            # Single series: use one color per category for visual clarity
            if n_ser == 1:
                bar_cols = [colors[i % len(colors)] for i in range(len(vals))]
            else:
                bar_cols = colors[si % len(colors)]
            bars = ax.bar(x[:len(vals)] + offsets[si], vals, width,
                          label=s.get("label",""), color=bar_cols, zorder=3)
            ax.bar_label(bars, fmt="%.1f", padding=2, fontsize=7, color=text_c)
        if cats:
            ax.set_xticks(x)
            ax.set_xticklabels(cats, rotation=25, ha="right", fontsize=8)
        if n_ser > 1:
            ax.legend(fontsize=8, framealpha=0.7)
        if x_label: ax.set_xlabel(x_label, fontsize=8)
        if y_label: ax.set_ylabel(y_label, fontsize=8)

    if title:
        ax.set_title(title, fontsize=11, fontweight="bold", color=text_c, pad=10)

    fig.tight_layout(pad=1.2)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=bg, edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf.read()
def _slide_chart(prs, s: dict, C: dict):
    """
    Chart slide: title bar + full-width chart image + optional caption.
    Adds its own slide to prs and returns it.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    import io as _io

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C.get("bg", "040A1C"))

    title_text   = s.get("title", "Chart")
    caption_text = s.get("body") or s.get("subtitle") or ""
    chart_spec   = s.get("chart") or {}

    white  = RGBColor(0xFF, 0xFF, 0xFF)
    accent = _rgb(C.get("accent", "00AADD"))

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(W_IN-0.8), Inches(0.55))
    tf = title_box.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = title_text; r.font.bold = True; r.font.size = Pt(22)
    r.font.color.rgb = white
    _lock_txb(title_box)
    # Rule
    _rect(slide, 0.4, 0.67, W_IN-0.8, 0.03, C.get("accent", "00AADD"))

    chart_top = 0.78
    chart_h   = H_IN - chart_top - (0.55 if caption_text else 0.25)

    try:
        png_bytes  = _render_chart_png(chart_spec,
                                        width_px=int((W_IN-0.6)*96),
                                        height_px=int(chart_h*96))
        slide.shapes.add_picture(_io.BytesIO(png_bytes),
                                  Inches(0.3), Inches(chart_top),
                                  Inches(W_IN-0.6), Inches(chart_h))
    except Exception as e:
        err_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(W_IN-2), Inches(1.5))
        err_box.text_frame.paragraphs[0].add_run().text = f"Chart render error: {e}"

    if caption_text:
        cap_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(H_IN-0.62), Inches(W_IN-0.8), Inches(0.5))
        cap_tf = cap_box.text_frame; cap_tf.word_wrap = True
        cap_r = cap_tf.paragraphs[0].add_run()
        cap_r.text = caption_text; cap_r.font.size = Pt(9.5); cap_r.font.italic = True
        cap_r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

    return slide



    """
    Build a PPTX from slide_data.

    Two modes:
    - template_bytes provided → open the template as the Presentation base,
      map slide types to the template's own layouts, fill title/body
      placeholders, and draw extra shapes (stats, comparison tables) on top.
      All master backgrounds, logos, fonts, and decorative elements are
      preserved exactly as they appear in the template.
    - No template → scratch-build with the SEGA dark palette (unchanged).
    """
    from pptx.util import Inches, Emu

    if template_bytes:
        return _generate_from_template(slide_data, template_bytes)

    # ── SCRATCH-BUILD PATH (SEGA dark palette) ────────────────────────────

    theme = slide_data.get("theme", {})

    def _safe_dark_hex(val, fallback):
        if not val or not isinstance(val, str): return fallback
        h = val.lstrip("#").upper()
        if len(h) != 6: return fallback
        try: int(h, 16)
        except ValueError: return fallback
        r, g, b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
        lum = (0.299*r + 0.587*g + 0.114*b) / 255
        return h if 0.04 < lum < 0.82 else fallback

    C = {
        "bg":        "040A1C",
        "subtle":    "1A2A4A",
        "header_bg": "0033AA",
        "white":     "FFFFFF",
        "light":     "D0E4FF",
        "midgray":   "8899BB",
        "neutral":   "AABBCC",
        "gold":      "F5C242",
        "green":     "22DD88",
        "red":       "EE3355",
        "dark":      "040A1C",
        "primary":   _safe_dark_hex(theme.get("primary"), "0055AA"),
        "accent":    _safe_dark_hex(theme.get("accent"),  "00AADD"),
        "heading_font": "Calibri",
        "body_font":    "Calibri",
    }

    prs = Presentation()
    prs.slide_width  = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    blank_layout = prs.slide_layouts[6]

    slides_data = slide_data.get("slides", [])
    total       = len(slides_data)

    RENDERERS = {
        "title":          _slide_title,
        "section":        _slide_section,
        "bullets":        _slide_bullets,
        "stats":          _slide_stats,
        "comparison":     _slide_comparison,
        "recommendation": _slide_recommendation,
        "closing":        _slide_closing,
        "chart":          _slide_chart,
    }

    for idx, s in enumerate(slides_data, start=1):
        slide    = prs.slides.add_slide(blank_layout)
        stype    = (s.get("type") or "bullets").lower()
        renderer = RENDERERS.get(stype, _slide_bullets)
        # chart slides manage their own slide creation
        if stype == "chart":
            renderer(prs, s, C)
        else:
            renderer(slide, s, idx, total, C)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return _patch_theme(buf.read())
W_IN, H_IN = 13.3, 7.5  # slide dimensions in inches

def generate_pptx(slide_data: dict, template_bytes: bytes | None = None) -> bytes:
    """
    Build a PPTX from slide_data.

    Two modes:
    - template_bytes provided → open the template as the Presentation base,
      map slide types to the template's own layouts, fill title/body
      placeholders, and draw extra shapes (stats, comparison tables) on top.
      All master backgrounds, logos, fonts, and decorative elements are
      preserved exactly as they appear in the template.
    - No template → scratch-build with the SEGA dark palette (unchanged).
    """
    from pptx.util import Inches, Emu

    if template_bytes:
        return _generate_from_template(slide_data, template_bytes)

    # ── SCRATCH-BUILD PATH (SEGA dark palette) ────────────────────────────

    theme = slide_data.get("theme", {})

    def _safe_dark_hex(val, fallback):
        if not val or not isinstance(val, str): return fallback
        h = val.lstrip("#").upper()
        if len(h) != 6: return fallback
        try: int(h, 16)
        except ValueError: return fallback
        r, g, b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
        lum = (0.299*r + 0.587*g + 0.114*b) / 255
        return h if 0.04 < lum < 0.82 else fallback

    C = {
        "bg":        "040A1C",
        "subtle":    "1A2A4A",
        "header_bg": "0033AA",
        "white":     "FFFFFF",
        "light":     "D0E4FF",
        "midgray":   "8899BB",
        "neutral":   "AABBCC",
        "gold":      "F5C242",
        "green":     "22DD88",
        "red":       "EE3355",
        "dark":      "040A1C",
        "primary":   _safe_dark_hex(theme.get("primary"), "0055AA"),
        "accent":    _safe_dark_hex(theme.get("accent"),  "00AADD"),
        "heading_font": "Calibri",
        "body_font":    "Calibri",
    }

    prs = Presentation()
    prs.slide_width  = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    blank_layout = prs.slide_layouts[6]

    slides_data = slide_data.get("slides", [])
    total       = len(slides_data)

    RENDERERS = {
        "title":          _slide_title,
        "section":        _slide_section,
        "bullets":        _slide_bullets,
        "stats":          _slide_stats,
        "comparison":     _slide_comparison,
        "recommendation": _slide_recommendation,
        "closing":        _slide_closing,
        "chart":          _slide_chart,
    }

    for idx, s in enumerate(slides_data, start=1):
        slide    = prs.slides.add_slide(blank_layout)
        stype    = (s.get("type") or "bullets").lower()
        renderer = RENDERERS.get(stype, _slide_bullets)
        if stype == "chart":
            renderer(prs, s, C)
        else:
            renderer(slide, s, idx, total, C)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return _patch_theme(buf.read())
def _copy_slide(prs_dest, slide_src):
    """
    Copy a slide from slide_src into prs_dest and return the new slide.
    Copies all shapes AND migrates image/media relationships so that
    pictures (logos, artwork, EMF files) actually appear in the output.
    """
    import copy
    from pptx.oxml.ns import qn

    # ── 1. Add blank slide with matching layout ───────────────────────────
    try:
        layout = slide_src.slide_layout
        dest_layout = next(
            (l for l in prs_dest.slide_layouts if l.name == layout.name),
            prs_dest.slide_layouts[0]
        )
        new_slide = prs_dest.slides.add_slide(dest_layout)
    except Exception:
        new_slide = prs_dest.slides.add_slide(prs_dest.slide_layouts[0])

    # ── 2. Migrate image/media parts ─────────────────────────────────────
    # Build a map: old rId → new rId so we can fix up the XML references
    SLIDE_IMG_REL = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
    )
    rId_map = {}   # old_rId → new_rId

    src_part  = slide_src.part
    dest_part = new_slide.part

    for rId, rel in src_part.rels.items():
        if rel.reltype == SLIDE_IMG_REL:
            try:
                img_part   = rel.target_part
                new_rId    = dest_part.relate_to(img_part, rel.reltype)
                rId_map[rId] = new_rId
            except Exception:
                pass   # if the part can't be related, skip — image won't show but won't crash

    # ── 3. Deep-copy the shape tree ───────────────────────────────────────
    new_sp_tree = new_slide.shapes._spTree
    src_sp_tree = slide_src.shapes._spTree

    for child in list(new_sp_tree):
        new_sp_tree.remove(child)

    for el in src_sp_tree:
        new_el = copy.deepcopy(el)
        # Fix any rId references in blip/hlinkClick/etc. within this element
        for node in new_el.iter():
            for attr in list(node.attrib):
                if attr.endswith('}embed') or attr.endswith('}link') or attr.endswith('}id'):
                    old = node.get(attr)
                    if old in rId_map:
                        node.set(attr, rId_map[old])
        new_sp_tree.append(new_el)

    # ── 4. Copy slide-level background if explicitly set ─────────────────
    bg = slide_src._element.find(qn("p:bg"))
    if bg is not None:
        new_slide._element.insert(2, copy.deepcopy(bg))

    return new_slide
def _set_txb_text(slide, shape_name, text, size_pt=None, bold=None, color_hex=None):
    """Find a shape by name on a slide and replace its text."""
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            tf = shape.text_frame
            # Preserve first run's formatting, replace text
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            if tf.paragraphs:
                p = tf.paragraphs[0]
                if p.runs:
                    run = p.runs[0]
                else:
                    run = p.add_run()
                run.text = str(text)
                if size_pt:
                    from pptx.util import Pt
                    run.font.size = Pt(size_pt)
                if bold is not None:
                    run.font.bold = bold
                if color_hex:
                    from pptx.dml.color import RGBColor
                    h = color_hex.lstrip("#")
                    run.font.color.rgb = RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
            return True
    return False
def _find_slide_by_pattern(prs, patterns):
    """
    Return the first slide whose shape names/text match any of the given patterns.
    patterns is a list of strings to look for in shape names or text.
    Falls back to slide index 0 if nothing matches.
    """
    for slide in prs.slides:
        names = {s.name.lower() for s in slide.shapes}
        texts = set()
        for s in slide.shapes:
            if s.has_text_frame:
                texts.add(s.text_frame.text.lower()[:30])
        for pat in patterns:
            pat_l = pat.lower()
            if any(pat_l in n for n in names) or any(pat_l in t for t in texts):
                return slide
    return prs.slides[0]
def _analyse_template(prs):
    """
    Classify each slide in the template into one of our 7 content types
    and return a type → slide_index mapping.

    Uses a scored approach: each candidate heuristic contributes points,
    and the highest scorer wins.  This is much more robust than the
    previous if/elif chain which missed stats slides entirely.

    For "blank branding" templates (≤5 slides, no bullet/content slides),
    returns a special dict with key "__blank_template__" = True so the
    caller can switch to the overlay-on-content-slide strategy.
    """
    # Detect blank branding templates — just title + section + blank + closing
    n_slides = len(prs.slides)
    n_bullet_slides = sum(
        1 for slide in prs.slides
        if sum(1 for s in slide.shapes
               if s.has_text_frame and s.text_frame.text.strip().startswith("▸")) >= 3
    )
    if n_slides <= 6 and n_bullet_slides == 0:
        # Find the best "content" slide — the white one with branding sidebar
        # (typically slide 2 in SEGA templates)
        content_idx = 0
        for i, slide in enumerate(prs.slides):
            texts = [s.text_frame.text.strip().upper() for s in slide.shapes if s.has_text_frame]
            # Content slide: has SUBJECT but NOT HEADLINE or THANK YOU
            if any("SUBJECT" in t for t in texts) and not any(t in ("HEADLINE","THANK YOU!") for t in texts):
                content_idx = i
                break
        title_idx = 0
        closing_idx = max(range(n_slides), key=lambda i: sum(
            1 for s in prs.slides[i].shapes
            if s.has_text_frame and "thank" in s.text_frame.text.lower()
        ))
        return {
            "__blank_template__": True,
            "title":          0,
            "section":        content_idx,
            "bullets":        content_idx,
            "stats":          content_idx,
            "comparison":     content_idx,
            "recommendation": content_idx,
            "closing":        closing_idx,
            "__content_idx__": content_idx,
        }

    type_map = {}

    for i, slide in enumerate(prs.slides):
        shapes     = list(slide.shapes)
        n_shapes   = len(shapes)
        texts_raw  = [s.text_frame.text for s in shapes if s.has_text_frame]
        texts      = [t.lower().strip() for t in texts_raw]
        texts_full = " ".join(texts)

        bullet_shapes  = [t for t in texts_raw if t.strip().startswith("▸")]
        n_bullets      = len(bullet_shapes)
        n_text_shapes  = len([t for t in texts_raw if t.strip()])
        n_rects        = len([s for s in shapes
                               if s.shape_type == 1  # MSO_SHAPE_TYPE.AUTO_SHAPE
                               and not s.has_text_frame])

        # Column balance — computed here so both stats and comparison can use it
        mid_x       = (prs.slide_width / 914400) / 2
        left_count  = sum(1 for s in shapes
                          if s.has_text_frame and s.text_frame.text.strip()
                          and (s.left or 0)/914400 < mid_x * 0.8)
        right_count = sum(1 for s in shapes
                          if s.has_text_frame and s.text_frame.text.strip()
                          and (s.left or 0)/914400 > mid_x * 1.1)

        scores = {t: 0 for t in
                  ("title","section","bullets","stats","comparison","recommendation","closing")}

        # ── TITLE: has HEADLINE and SUBJECT placeholder text ──────────────
        if any("headline" in t for t in texts):   scores["title"] += 5
        if any("subject"  in t for t in texts):   scores["title"] += 3
        # Penalise if it also has lots of bullets (probably not a title)
        if n_bullets >= 3: scores["title"] -= 4

        # ── CLOSING: thank-you text or subject-only (no headline) ─────────
        if any("thank you" in t for t in texts):  scores["closing"] += 6
        if (any("subject" in t for t in texts) and
                not any("headline" in t for t in texts)):
            scores["closing"] += 3
        if n_text_shapes <= 3:                    scores["closing"] += 1

        # ── BULLETS: 4+ bullet (▸) shapes ────────────────────────────────
        if n_bullets >= 6:  scores["bullets"] += 6
        elif n_bullets >= 4: scores["bullets"] += 4
        elif n_bullets >= 2: scores["bullets"] += 1

        # ── RECOMMENDATION: second-priority bullets slide ─────────────────
        # scored same as bullets; de-duplication happens in fallback below

        # ── STATS: card layout — values & labels in X-grouped columns ────
        import re as _re2
        number_texts = [t for t in texts_raw if _re2.search(r'\d+\.?\d*[MmKk%$]?', t.strip()) and
                        len(t.strip()) < 25 and not t.strip().startswith("▸")]
        if len(number_texts) >= 3:       scores["stats"] += 4
        if len(number_texts) >= 6:       scores["stats"] += 2
        if n_bullets == 0:               scores["stats"] += 2
        if 10 <= n_shapes <= 24:         scores["stats"] += 1
        if n_rects >= 2:                 scores["stats"] += 2
        # Penalise if left/right columns are balanced (that's a comparison slide)
        if left_count >= 4 and right_count >= 4:
            balance = min(left_count, right_count) / max(left_count, right_count)
            if balance > 0.5:            scores["stats"] -= 5  # strong column symmetry = comparison
        # Bonus if shapes cluster in clear X columns with ~3-4 cols (card grid)
        if n_text_shapes >= 8:
            x_vals = sorted((s.left or 0)/914400 for s in shapes
                            if s.has_text_frame and s.text_frame.text.strip())
            if x_vals:
                x_range = x_vals[-1] - x_vals[0]
                # Cards spread across the full width
                if x_range > 7.0:        scores["stats"] += 3

        # ── COMPARISON: two column headers + rows of label/left/right ─────
        if n_text_shapes >= 6 and n_bullets == 0:
            if left_count >= 3 and right_count >= 3:
                scores["comparison"] += 5
            if n_text_shapes >= 10:  scores["comparison"] += 2
            if n_text_shapes >= 16:  scores["comparison"] += 2

        # ── SECTION: sparse, large text, no bullets ───────────────────────
        if n_text_shapes <= 2 and n_bullets == 0:  scores["section"] += 3
        if n_shapes <= 4:                          scores["section"] += 2

        # Determine best type for this slide
        best_type  = max(scores, key=lambda t: scores[t])
        best_score = scores[best_type]

        if best_score < 2:
            continue   # too ambiguous, skip

        # Store first occurrence of each type; bullets → recommendation if bullets taken
        if best_type == "bullets" and "bullets" in type_map and "recommendation" not in type_map:
            type_map["recommendation"] = i
        elif best_type not in type_map:
            type_map[best_type] = i

    # Fill missing types with best fallback
    fallback_order = ["bullets","title","comparison","stats","recommendation","section","closing"]
    first_available = next(
        (type_map[t] for t in fallback_order if t in type_map), 0
    )
    for t in ["title","section","bullets","stats","comparison","recommendation","closing"]:
        if t not in type_map:
            type_map[t] = first_available

    return type_map
def _generate_from_template(slide_data: dict, template_bytes: bytes) -> bytes:
    """
    Build a PPTX using the uploaded template as the actual presentation base.

    Approach:
    1. Load the template, analyse which slides match which content types
    2. For each output slide, copy the best-matching template slide
    3. Replace text box content in-place, preserving all styling/positioning

    This means ALL template branding — backgrounds, logos, fonts, colours,
    decorative shapes — is preserved exactly. We only change text values.
    """
    import copy
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs_template = Presentation(io.BytesIO(template_bytes))
    W = prs_template.slide_width.inches
    H = prs_template.slide_height.inches

    # Analyse the template
    type_map = _analyse_template(prs_template)

    # Quality check: if fewer than 4 distinct slide indices are mapped,
    # the template couldn't be classified properly — return a diagnostic error
    # so the caller can fall back to scratch-build.
    is_blank_template = type_map.get("__blank_template__", False)
    unique_indices = set(v for k, v in type_map.items() if not k.startswith("__"))

    if not is_blank_template and len(unique_indices) < 4:
        raise ValueError(
            f"Template classification failed — only {len(unique_indices)} distinct slide "
            f"types found (need ≥4). type_map={type_map}. "
            "Check that the template contains bullet, comparison, and stats slides."
        )

    # Strip existing slides at zip level (only reliable approach in python-pptx)
    # then open the cleaned template as our base.
    import zipfile as _zf2, re as _re2
    _src2 = _zf2.ZipFile(io.BytesIO(template_bytes))
    _out2 = io.BytesIO()
    _dst2 = _zf2.ZipFile(_out2, "w", _zf2.ZIP_DEFLATED)
    _spat = _re2.compile(r"^ppt/slides/slide\d+\.xml(\.rels)?$")
    for _item in _src2.infolist():
        if _spat.match(_item.filename): continue
        _data = _src2.read(_item.filename)
        if _item.filename == "[Content_Types].xml":
            _t = etree.fromstring(_data)
            _ct  = "http://schemas.openxmlformats.org/package/2006/content-types"
            _sct = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
            for _o in _t.findall(f"{{{_ct}}}Override"):
                if _o.get("ContentType") == _sct: _t.remove(_o)
            _data = etree.tostring(_t, xml_declaration=True, encoding="UTF-8", standalone=True)
        elif _item.filename == "ppt/presentation.xml":
            _t = etree.fromstring(_data)
            _pns = "http://schemas.openxmlformats.org/presentationml/2006/main"
            _lst = _t.find(f"{{{_pns}}}sldIdLst")
            if _lst is not None:
                for _el in list(_lst): _lst.remove(_el)
            _data = etree.tostring(_t, xml_declaration=True, encoding="UTF-8", standalone=True)
        elif _item.filename == "ppt/_rels/presentation.xml.rels":
            _t = etree.fromstring(_data)
            _sr = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
            for _el in list(_t):
                if _el.get("Type") == _sr: _t.remove(_el)
            _data = etree.tostring(_t, xml_declaration=True, encoding="UTF-8", standalone=True)
        _dst2.writestr(_item, _data)
    _dst2.close()
    prs_out = Presentation(io.BytesIO(_out2.getvalue()))

    slides_data = slide_data.get("slides", [])

    def _get_source_slide(stype):
        idx = type_map.get(stype, type_map.get("bullets", 0))
        return prs_template.slides[idx]

    def _txbs(slide):
        """Return all text-bearing shapes on a slide, sorted top-to-bottom."""
        return sorted(
            [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()],
            key=lambda s: (s.top or 0)
        )

    def _set_text(shape, text):
        """Replace text in a shape preserving run formatting. Auto-shrinks if needed."""
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        tf.auto_size = None   # let the shape constrain text
        # Clear all paragraphs except first, set first paragraph's first run
        for i, para in enumerate(tf.paragraphs):
            for j, run in enumerate(para.runs):
                if i == 0 and j == 0:
                    run.text = str(text)
                else:
                    run.text = ""
            if i > 0:
                # Remove extra paragraphs by clearing their text
                for run in para.runs:
                    run.text = ""
        if tf.paragraphs and not tf.paragraphs[0].runs:
            tf.paragraphs[0].add_run().text = str(text)

    # Helper: find the SUBJECT sidebar shape (rotated TextBox 54 / TextBox 23 etc.)
    def _subject_shape(slide):
        """Find the rotated subject/sidebar label shape."""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            nm = shape.name.lower()
            if t.upper() == "SUBJECT" or "subject" in nm:
                return shape
            xfrm = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
            if xfrm is not None and xfrm.get('rot'):
                rot = int(xfrm.get('rot', 0))
                if abs(rot) > 8000000:
                    return shape
        return None

    # Shared palette + layout helpers (used by both blank-template overlay and chart slides)
    from pptx.util import Inches as _In, Pt as _Pt, Emu as _Emu
    from pptx.dml.color import RGBColor as _RGB
    from pptx.enum.text import PP_ALIGN as _PPA
    import io as _ioC

    tp     = extract_template_palette(template_bytes)
    TEXT   = tp.get("white", "111122")
    MUTED  = tp.get("midgray", "556677")
    ACCENT = tp.get("accent",  "00337E")
    FONT   = tp.get("body_font", "Calibri")
    CX     = 0.95
    CW     = W - CX - 0.15
    CY     = 0.45

    def _rgb3(h):
        h = (h or "111122").lstrip('#').upper().ljust(6,'0')[:6]
        return _RGB(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    def _txb(slide_ref, x, y, w, h, text, size=12, bold=False, color=None,
             align=None, italic=False):
        """Add a text box to slide_ref."""
        tb = slide_ref.shapes.add_textbox(_In(x), _In(y), _In(w), _In(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align:
            p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = _Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = FONT
        run.font.color.rgb = _rgb3(color or TEXT)
        return tb

    def _rect(slide_ref, x, y, w, h, fill_hex):
        """Add a filled rectangle to slide_ref."""
        shp = slide_ref.shapes.add_shape(1, _In(x), _In(y), _In(w), _In(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb3(fill_hex)
        shp.line.fill.background()
        return shp

    def _add_slide_from_source(stype, s):
        """Copy the template slide for stype and populate with slide data s."""
        src = _get_source_slide(stype)
        new_slide = _copy_slide(prs_out, src)

        # ── BLANK TEMPLATE OVERLAY PATH ──────────────────────────────────
        if is_blank_template and stype not in ("title", "closing", "chart"):

            # Write title into SUBJECT sidebar — truncate to prevent overflow
            subj = _subject_shape(new_slide)
            if subj:
                sidebar_text = s.get("title", "")
                if len(sidebar_text) > 28:
                    sidebar_text = sidebar_text[:26] + "…"
                _set_text(subj, sidebar_text)

            if stype in ("bullets", "recommendation", "section"):
                # Section header + rule (add rule BEFORE bullets — it's a background element)
                _txb(new_slide, CX, CY, CW, 0.5, s.get("title",""),
                     size=22, bold=True, color=ACCENT)
                _rect(new_slide, CX, CY+0.58, CW, 0.02, ACCENT)  # thin rule

                bullets = (s.get("bullets") or [])[:6]
                for bi, bullet in enumerate(bullets):
                    y = CY + 0.75 + bi * 0.72
                    _txb(new_slide, CX, y, CW, 0.65, f"▸  {bullet}", size=13, color=TEXT)
                if s.get("body"):
                    # Place body text as a subtitle line under the title, not at the bottom
                    _txb(new_slide, CX, CY + 0.65, CW, 0.55, s["body"], size=11, italic=True, color=MUTED)

            elif stype == "stats":
                _txb(new_slide, CX, CY, CW, 0.5, s.get("title",""),
                     size=20, bold=True, color=ACCENT)

                stats = (s.get("stats") or [])[:4]
                n = max(len(stats), 1)
                card_w = (CW - 0.12 * (n-1)) / n
                for si, stat in enumerate(stats):
                    cx = CX + si * (card_w + 0.12)
                    # Background card first
                    _rect(new_slide, cx, CY+0.62, card_w, H-CY-0.82, "EEF2FF")
                    # Then text on top
                    _txb(new_slide, cx+0.1, CY+0.78, card_w-0.2, 1.2,
                         stat.get("value",""), size=36, bold=True,
                         color=ACCENT, align=_PPA.CENTER)
                    _txb(new_slide, cx+0.1, CY+2.05, card_w-0.2, 0.5,
                         stat.get("label",""), size=12, bold=True,
                         color=TEXT, align=_PPA.CENTER)
                    _txb(new_slide, cx+0.1, CY+2.62, card_w-0.2, 0.4,
                         stat.get("note",""), size=10,
                         color=MUTED, align=_PPA.CENTER)

            elif stype == "comparison":
                _txb(new_slide, CX, CY, CW, 0.5, s.get("title",""),
                     size=20, bold=True, color=ACCENT)

                cmp   = s.get("comparison") or {}
                rows  = (cmp.get("rows") or [])[:8]
                lbl_w = CW * 0.28
                col_w = (CW - lbl_w) / 2

                # Draw ALL backgrounds first, then ALL text on top
                # Header backgrounds
                _rect(new_slide, CX + lbl_w,           CY+0.57, col_w-0.04, 0.38, ACCENT)
                _rect(new_slide, CX + lbl_w + col_w,   CY+0.57, col_w-0.04, 0.38, "ED7D31")

                # Row backgrounds
                for ri in range(len(rows)):
                    ry = CY + 1.02 + ri * 0.46
                    _rect(new_slide, CX, ry, CW, 0.40, "EEF2FF" if ri % 2 == 0 else "F5F7FF")

                # Now header text (on top of header backgrounds)
                _txb(new_slide, CX + lbl_w,           CY+0.60, col_w-0.04, 0.32,
                     cmp.get("left_title","Left"), size=10, bold=True,
                     color="FFFFFF", align=_PPA.CENTER)
                _txb(new_slide, CX + lbl_w + col_w,   CY+0.60, col_w-0.04, 0.32,
                     cmp.get("right_title","Right"), size=10, bold=True,
                     color="FFFFFF", align=_PPA.CENTER)

                # Row text (on top of row backgrounds)
                for ri, row in enumerate(rows):
                    ry = CY + 1.02 + ri * 0.46
                    _txb(new_slide, CX+0.05,               ry+0.04, lbl_w-0.1, 0.34,
                         row.get("label",""), size=9, bold=True, color=MUTED)
                    _txb(new_slide, CX+lbl_w,              ry+0.04, col_w-0.08, 0.34,
                         row.get("left",""), size=9, color=TEXT, align=_PPA.CENTER)
                    delta = row.get("delta","neutral")
                    dc = "22AA66" if delta=="positive" else "CC2244" if delta=="negative" else MUTED
                    _txb(new_slide, CX+lbl_w+col_w,        ry+0.04, col_w-0.08, 0.34,
                         row.get("right",""), size=9, color=dc, align=_PPA.CENTER)

            return new_slide


        # ── CHART (both template paths) ───────────────────────────────────
        if stype == "chart":
            chart_spec   = s.get("chart") or {}
            title_text   = s.get("title", "Chart")
            caption_text = s.get("body") or s.get("subtitle") or ""

            subj = _subject_shape(new_slide)
            if subj:
                _set_text(subj, title_text[:26]+"…" if len(title_text)>28 else title_text)

            # Title + rule — pass slide_ref=new_slide since closure would grab a stale ref
            _txb(new_slide, CX, 0.1, CW+0.1, 0.55, title_text, size=22, bold=True,
                 color=ACCENT)
            _rect(new_slide, CX, 0.67, CW, 0.02, ACCENT)

            chart_top = 0.78
            chart_h   = H - chart_top - (0.60 if caption_text else 0.25)
            try:
                png_bytes = _render_chart_png(chart_spec,
                                               width_px=int(CW*96),
                                               height_px=int(chart_h*96))
                new_slide.shapes.add_picture(
                    _ioC.BytesIO(png_bytes),
                    _In(CX), _In(chart_top), _In(CW), _In(chart_h)
                )
            except Exception as _ce:
                eb = new_slide.shapes.add_textbox(_In(CX), _In(2), _In(CW), _In(1.5))
                eb.text_frame.paragraphs[0].add_run().text = f"Chart render error: {_ce}"

            if caption_text:
                _txb(new_slide, CX, H-0.65, CW, 0.5, caption_text, size=9, italic=True,
                     color=MUTED)
            return new_slide

        # ── TITLE ────────────────────────────────────────────────────────
        if stype == "title":
            for shape in new_slide.shapes:
                if not shape.has_text_frame:
                    continue
                t = shape.text_frame.text.strip()
                nm = shape.name.lower()
                if t.upper() == "HEADLINE" or "headline" in nm:
                    _set_text(shape, s.get("title", ""))
                elif t.upper() == "SUBJECT" or "subject" in nm:
                    _set_text(shape, s.get("subtitle") or s.get("body") or "")

        # ── SECTION ──────────────────────────────────────────────────────
        elif stype == "section":
            # This template's section slide is mostly blank — just a rotated
            # SUBJECT sidebar. Write title into SUBJECT and add a large
            # title overlay in the center of the white content area.
            title_text = s.get("title", "")
            subj = _subject_shape(new_slide)
            if subj:
                _set_text(subj, title_text)

            # Add a large centered title text box over the white area
            txb = new_slide.shapes.add_textbox(
                _In(1.5), _In(2.5), _In(W - 2.5), _In(2.5)
            )
            tf = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = _PPA.LEFT
            run = p.add_run()
            run.text = title_text
            run.font.bold = True
            run.font.size = _Pt(40)
            run.font.color.rgb = _RGB(0x11, 0x11, 0x22)   # dark navy — matches template text
            if s.get("subtitle"):
                p2 = tf.add_paragraph()
                p2.alignment = _PPA.LEFT
                r2 = p2.add_run()
                r2.text = s["subtitle"]
                r2.font.size = _Pt(18)
                r2.font.color.rgb = _RGB(0x44, 0x54, 0x6A)

        # ── BULLETS ──────────────────────────────────────────────────────
        elif stype == "bullets":
            bullets = (s.get("bullets") or [])[:6]
            bullet_shapes = sorted(
                [sh for sh in new_slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip().startswith("▸")],
                key=lambda sh: (sh.top or 0)
            )
            # This template's bullet slides have no title — write title into
            # the SUBJECT sidebar (rotated label), clear it if no title
            subj = _subject_shape(new_slide)
            if subj:
                _set_text(subj, s.get("title","")[:26]+"…" if len(s.get("title",""))>28 else s.get("title",""))
            # Also check for explicit HEADLINE shape (other templates may have it)
            for shape in new_slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    nm = shape.name.lower()
                    if t.upper() == "HEADLINE" or "headline" in nm:
                        _set_text(shape, s.get("title", ""))
            # Fill bullet slots
            for i, shape in enumerate(bullet_shapes[:len(bullets)]):
                _set_text(shape, "▸  " + bullets[i])
            # Clear unused bullet slots
            for shape in bullet_shapes[len(bullets):]:
                _set_text(shape, "")

        elif stype == "comparison":
            cmp   = s.get("comparison") or {}
            rows  = (cmp.get("rows") or [])[:8]
            # Find column header boxes — first two non-label text boxes near top
            # and the row label/value boxes pattern
            all_txbs = sorted(
                [sh for sh in new_slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()],
                key=lambda sh: ((sh.top or 0), (sh.left or 0))
            )
            # Detect column positions from existing shapes
            # The template has 3 columns: label, left-value, right-value
            # Find by x-position clustering
            x_positions = sorted(set(
                round((sh.left or 0) / 914400 / 0.5) * 0.5   # round to 0.5"
                for sh in all_txbs if (sh.left or 0) > 0
            ))
            # Column x groups: label ~0.5", left ~3.3", right ~8.2"
            label_x = min(x_positions) if x_positions else 0.5
            remaining_x = [x for x in x_positions if x > label_x + 0.5]
            left_x  = remaining_x[0] if len(remaining_x) > 0 else 3.3
            right_x = remaining_x[1] if len(remaining_x) > 1 else 8.2

            def x_group(shape, target, tol=0.8):
                return abs((shape.left or 0)/914400 - target) < tol

            header_row = sorted(
                [sh for sh in all_txbs
                 if x_group(sh, left_x) or x_group(sh, right_x)],
                key=lambda sh: (sh.top or 0)
            )
            # First shapes at left_x and right_x are column headers
            left_headers  = [sh for sh in header_row if x_group(sh, left_x)]
            right_headers = [sh for sh in header_row if x_group(sh, right_x)]
            if left_headers:  _set_text(left_headers[0],  cmp.get("left_title",  "Internal"))
            if right_headers: _set_text(right_headers[0], cmp.get("right_title", "Reference"))

            # Row data: label, left_val, right_val in y order
            label_shapes = sorted(
                [sh for sh in all_txbs if x_group(sh, label_x) and sh not in left_headers[:1]],
                key=lambda sh: sh.top or 0
            )
            left_shapes = left_headers[1:] if len(left_headers) > 1 else []
            right_shapes = right_headers[1:] if len(right_headers) > 1 else []

            for i, row in enumerate(rows):
                if i < len(label_shapes): _set_text(label_shapes[i], row.get("label", ""))
                if i < len(left_shapes):  _set_text(left_shapes[i],  row.get("left", ""))
                if i < len(right_shapes): _set_text(right_shapes[i], row.get("right", ""))
            # Clear unused rows
            for i in range(len(rows), max(len(label_shapes), len(left_shapes), len(right_shapes))):
                if i < len(label_shapes): _set_text(label_shapes[i], "")
                if i < len(left_shapes):  _set_text(left_shapes[i],  "")
                if i < len(right_shapes): _set_text(right_shapes[i], "")

        elif stype == "stats":
            stats = (s.get("stats") or [])[:4]
            # Stats cards: groups of 3 stacked text boxes (value, label, note)
            # Each group is above a filled rectangle
            # Sort all non-empty txbs by (left, top) to find card groups
            all_txbs_l = sorted(
                [sh for sh in new_slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()],
                key=lambda sh: ((sh.left or 0), (sh.top or 0))
            )
            # Group by x-position (card columns)
            from itertools import groupby
            def x_col(sh):
                return round((sh.left or 0) / 914400 / 3.5)  # group into ~3.5" wide buckets
            cards = []
            for _, grp in groupby(all_txbs_l, key=x_col):
                card_shapes = sorted(list(grp), key=lambda sh: sh.top or 0)
                if len(card_shapes) >= 2:
                    cards.append(card_shapes)

            for i, stat in enumerate(stats):
                if i < len(cards):
                    card = cards[i]
                    # First shape (top) = big value
                    if len(card) > 0: _set_text(card[0], stat.get("value", ""))
                    # Second = label
                    if len(card) > 1: _set_text(card[1], stat.get("label", ""))
                    # Third = note
                    if len(card) > 2: _set_text(card[2], stat.get("note", ""))

        elif stype == "recommendation":
            # Same layout as bullets — 6 bullet slots + SUBJECT sidebar for title
            bullets = (s.get("bullets") or [])[:6]
            bullet_shapes = sorted(
                [sh for sh in new_slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip().startswith("▸")],
                key=lambda sh: sh.top or 0
            )
            subj = _subject_shape(new_slide)
            if subj:
                _set_text(subj, s.get("title","")[:26]+"…" if len(s.get("title",""))>28 else s.get("title",""))
            for shape in new_slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t.upper() == "HEADLINE" or "headline" in shape.name.lower():
                        _set_text(shape, s.get("title", ""))
            for i, shape in enumerate(bullet_shapes[:len(bullets)]):
                _set_text(shape, "▸  " + bullets[i])
            for shape in bullet_shapes[len(bullets):]:
                _set_text(shape, "")

        elif stype == "closing":
            # Closing slide — THANK YOU! box is 72pt in a 9"×2.5" area.
            # Scale font down for longer titles to prevent overflow.
            title_text    = s.get("title", "")
            subtitle_text = s.get("subtitle") or s.get("body") or ""
            subj = _subject_shape(new_slide)
            if subj:
                _set_text(subj, subtitle_text)
            for shape in new_slide.shapes:
                if not shape.has_text_frame:
                    continue
                t  = shape.text_frame.text.strip()
                nm = shape.name.lower()
                if t.upper() in ("THANK YOU!", "HEADLINE") or "headline" in nm:
                    _set_text(shape, title_text)
                    # Scale font size down for longer titles
                    from pptx.util import Pt as _Pt2
                    nchars = len(title_text)
                    pt = 72 if nchars <= 15 else 54 if nchars <= 25 else 40 if nchars <= 40 else 30
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = _Pt2(pt)
                elif t.upper() == "SUBJECT" or "subject" in nm:
                    _set_text(shape, subtitle_text)
                elif t.isdigit() or "presentation template" in t.lower():
                    _set_text(shape, "")
            # Add subtitle as a visible overlay on the blue content area
            if subtitle_text:
                stb = new_slide.shapes.add_textbox(
                    _In(1.3), _In(H * 0.62), _In(W - 2.0), _In(0.6)
                )
                stf = stb.text_frame
                stf.word_wrap = True
                sp = stf.paragraphs[0]
                sp.alignment = _PPA.CENTER
                sr = sp.add_run()
                sr.text = subtitle_text
                sr.font.size = _Pt(18)
                sr.font.color.rgb = _RGB(0xFF, 0xFF, 0xFF)
                sr.font.bold = False

    for s in slides_data:
        stype = (s.get("type") or "bullets").lower()
        _add_slide_from_source(stype, s)

    buf = io.BytesIO()
    prs_out.save(buf)
    return buf.getvalue()
def _api_post(headers: dict, payload: dict, timeout: int = 90,
              on_wait=None) -> dict:
    """
    POST to /v1/messages with exponential back-off on 429.
    on_wait(seconds, attempt) is called before each sleep so callers can
    surface a countdown message to the UI.
    """
    for attempt in range(_RL_MAX_TRIES):
        resp = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers=headers, json=payload, timeout=timeout,
        )
        if resp.status_code == 429:
            if attempt == _RL_MAX_TRIES - 1:
                raise RuntimeError(
                    f"Rate limited after {_RL_MAX_TRIES} retries. "
                    "Try switching to Haiku in the sidebar (much higher rate limit) "
                    "or wait a minute and run again."
                )
            # Honour retry-after header when present; otherwise double each time
            try:
                wait = max(float(
                    resp.headers.get("retry-after") or
                    resp.headers.get("x-ratelimit-reset-requests") or
                    _RL_WAIT_BASE * (2 ** attempt)
                ), 1.0)
            except (TypeError, ValueError):
                wait = _RL_WAIT_BASE * (2 ** attempt)
            if on_wait:
                on_wait(wait, attempt + 1)
            time.sleep(wait)
            continue
        resp.raise_for_status()
        return resp.json()
    raise RuntimeError("Unexpected exit from retry loop.")
def _api_stream(headers: dict, payload: dict, timeout: int = 240,
                on_rate_limit=None):
    """
    POST with stream=True using httpx for a true wall-clock timeout.
    Yields text delta chunks as they arrive via SSE.
    Retries on 429; raises RuntimeError on other errors.
    """
    import httpx

    for attempt in range(_RL_MAX_TRIES):
        stream_payload = {**payload, "stream": True}
        stream_headers = {**headers, "Accept": "text/event-stream"}

        _retry_after = None  # set inside the with-block, acted on outside

        try:
            with httpx.stream(
                "POST",
                "https://api.anthropic.com/v1/messages",
                headers=stream_headers,
                json=stream_payload,
                timeout=httpx.Timeout(
                    timeout,          # total wall-clock limit
                    connect=10,       # connection timeout
                    read=timeout,     # per-read timeout
                    write=30,         # write timeout
                    pool=10,          # pool acquisition timeout
                ),
            ) as resp:
                if resp.status_code == 429:
                    if attempt == _RL_MAX_TRIES - 1:
                        raise RuntimeError(
                            "Rate limited after max retries. "
                            "Switch to Haiku in the sidebar or wait a minute."
                        )
                    try:
                        _retry_after = max(float(
                            resp.headers.get("retry-after") or
                            _RL_WAIT_BASE * (2 ** attempt)
                        ), 1.0)
                    except (TypeError, ValueError):
                        _retry_after = _RL_WAIT_BASE * (2 ** attempt)
                    # Must exit the with-block before sleeping and retrying
                elif resp.status_code != 200:
                    # Must call read() before accessing .text in a stream context
                    resp.read()
                    raise RuntimeError(
                        f"API error {resp.status_code}: {resp.text[:400]}"
                    )
                else:
                    for raw_line in resp.iter_lines():
                        line = raw_line.strip()
                        if not line or not line.startswith("data: "):
                            continue
                        payload_str = line[6:].strip()
                        if payload_str == "[DONE]":
                            return
                        try:
                            evt = json.loads(payload_str)
                        except json.JSONDecodeError:
                            continue
                        if evt.get("type") == "content_block_delta":
                            delta = evt.get("delta", {})
                            if delta.get("type") == "text_delta":
                                yield delta.get("text", "")
                    return  # clean exit

        except httpx.TimeoutException:
            raise RuntimeError(
                f"Analysis API timed out after {timeout}s. "
                "Try reducing the slide count or switching to Haiku in the sidebar."
            )

        # Handle 429 retry outside the with-block
        if _retry_after is not None:
            if on_rate_limit:
                on_rate_limit(_retry_after, attempt + 1)
            time.sleep(_retry_after)
def _render_plan_modal(template_bytes_ref):
    """
    Renders the Plan Mode outline editor as a full-page modal overlay.
    Reads/writes st.session_state["plan_slide_data"].
    """
    sd     = st.session_state.get("plan_slide_data", {})
    slides = sd.get("slides", [])

    st.markdown(
        "<style>"
        ".pm-card{background:#1e293b;border:1px solid #334155;border-radius:8px;"
        "padding:.8rem 1rem;margin-bottom:.5rem}"
        ".pm-num{color:#60a5fa;font-size:.68rem;font-weight:700;"
        "text-transform:uppercase;letter-spacing:.07em;margin-bottom:.25rem}"
        "button[data-testid='stBaseButton-secondary'] p{"
        "width:100%;text-align:center !important;margin:0 !important;"
        "display:block !important;font-size:1.05rem !important;line-height:1 !important}"
        "button[data-testid='stBaseButton-secondary']{"
        "padding:0.2rem 0 !important}"
        "</style>",
        unsafe_allow_html=True,
    )

    hcol, xcol = st.columns([4, 1])
    with hcol:
        st.markdown(
            f"<h3 style='color:#e2e8f0;margin:0'>✏️ Plan Mode &nbsp;·&nbsp; "
            f"<span style='color:#60a5fa'>{sd.get('title','Presentation')}</span></h3>",
            unsafe_allow_html=True,
        )
        st.caption(f"{len(slides)} slides · Edit titles & content · ⬆⬇ reorder · ➕ add · 🗑 delete")
    with xcol:
        if st.button("✕ Close", key="pm_close", use_container_width=True):
            st.session_state.pop("plan_slide_data", None)
            st.session_state.pop("plan_mode_active", None)
            st.rerun()

    st.divider()

    updated_slides = list(slides)
    move_up = move_down = delete_idx = insert_after = None

    for i, slide in enumerate(slides):
        stype = slide.get("type", "bullets")
        with st.expander(
            f"**{i+1}.** `{stype.upper()}` — {slide.get('title','(untitled)')}",
            expanded=(i == 0),
        ):
            rc1, rc2 = st.columns([2, 1])
            with rc1:
                new_type = st.selectbox(
                    "Type", SLIDE_TYPES,
                    index=SLIDE_TYPES.index(stype) if stype in SLIDE_TYPES else 1,
                    key=f"pm_type_{i}", label_visibility="collapsed",
                )
            with rc2:
                b1, b2, b3, b4 = st.columns(4)
                if b1.button("⬆", key=f"pu_{i}", help="Move up",      use_container_width=True): move_up      = i
                if b2.button("⬇", key=f"pd_{i}", help="Move down",    use_container_width=True): move_down    = i
                if b3.button("➕", key=f"pa_{i}", help="Insert after", use_container_width=True): insert_after = i
                if b4.button("🗑", key=f"px_{i}", help="Delete",       use_container_width=True): delete_idx   = i

            new_title    = st.text_input("Title",    slide.get("title",""),    key=f"pm_ti_{i}")
            new_subtitle = st.text_input("Subtitle", slide.get("subtitle") or slide.get("body",""), key=f"pm_su_{i}")

            new_slide = {**slide, "type": new_type, "title": new_title, "subtitle": new_subtitle}

            # ── Type-specific editors ──────────────────────────────────────
            if new_type in ("bullets", "recommendation"):
                raw = st.text_area(
                    "Bullets (one per line, max 6)",
                    value="\n".join(slide.get("bullets") or []),
                    height=150, key=f"pm_bu_{i}",
                )
                new_slide["bullets"] = [b.strip() for b in raw.split("\n") if b.strip()][:6]

            elif new_type == "stats":
                st.caption("Format: value | label | note")
                raw = st.text_area(
                    "Stats", height=110, key=f"pm_st_{i}",
                    label_visibility="collapsed",
                    value="\n".join(
                        f"{s.get('value','')} | {s.get('label','')} | {s.get('note','')}"
                        for s in (slide.get("stats") or [])
                    ),
                )
                new_stats = []
                for line in raw.split("\n"):
                    p = [x.strip() for x in line.split("|")]
                    if any(p):
                        new_stats.append({"value": p[0] if p else "", "label": p[1] if len(p)>1 else "", "note": p[2] if len(p)>2 else ""})
                new_slide["stats"] = new_stats[:4]

            elif new_type == "comparison":
                cmp  = slide.get("comparison") or {}
                rows = cmp.get("rows") or []
                cl, cr = st.columns(2)
                lt = cl.text_input("Left col title",  cmp.get("left_title",""),  key=f"pm_lt_{i}")
                rt = cr.text_input("Right col title", cmp.get("right_title",""), key=f"pm_rt_{i}")
                st.caption("Rows: label | left | right | delta")
                raw = st.text_area("Rows", height=140, key=f"pm_ro_{i}", label_visibility="collapsed",
                    value="\n".join(
                        f"{r.get('label','')} | {r.get('left','')} | {r.get('right','')} | {r.get('delta','neutral')}"
                        for r in rows
                    ),
                )
                new_rows = []
                for line in raw.split("\n"):
                    p = [x.strip() for x in line.split("|")]
                    if len(p) >= 3 and any(p):
                        new_rows.append({"label":p[0],"left":p[1],"right":p[2],"delta":p[3] if len(p)>3 else "neutral"})
                new_slide["comparison"] = {"left_title":lt,"right_title":rt,"rows":new_rows[:8]}

            elif new_type == "chart":
                chart = slide.get("chart") or {}
                ct_opts = ["bar","line","scatter","pie","horizontal_bar"]
                ct  = st.selectbox("Chart type", ct_opts,
                                    index=ct_opts.index(chart.get("chart_type","bar")),
                                    key=f"pm_ct_{i}")
                ca, cb = st.columns(2)
                xl  = ca.text_input("X label", chart.get("x_label",""), key=f"pm_xl_{i}")
                yl  = cb.text_input("Y label", chart.get("y_label",""), key=f"pm_yl_{i}")
                cats = st.text_input(
                    "Categories (comma-separated)",
                    ", ".join(chart.get("categories") or []), key=f"pm_ca_{i}"
                )
                st.caption("Series: label | val1, val2, …  (one per line)")
                raw = st.text_area("Series", height=100, key=f"pm_se_{i}", label_visibility="collapsed",
                    value="\n".join(
                        f"{s.get('label','')} | {', '.join(str(v) for v in s.get('values',[]))}"
                        for s in (chart.get("series") or [])
                    ),
                )
                new_series = []
                for line in raw.split("\n"):
                    p = [x.strip() for x in line.split("|")]
                    if len(p) >= 2:
                        try:
                            vals = [float(v.strip()) for v in p[1].split(",") if v.strip()]
                        except ValueError:
                            vals = []
                        new_series.append({"label": p[0], "values": vals})
                new_slide["chart"] = {
                    "chart_type": ct, "x_label": xl, "y_label": yl,
                    "categories": [c.strip() for c in cats.split(",") if c.strip()],
                    "series": new_series,
                }

            updated_slides[i] = new_slide

    # ── Reorder / delete / insert ────────────────────────────────────────────
    if move_up is not None and move_up > 0:
        updated_slides[move_up-1], updated_slides[move_up] = updated_slides[move_up], updated_slides[move_up-1]
        sd["slides"] = updated_slides; st.session_state["plan_slide_data"] = sd; st.rerun()
    if move_down is not None and move_down < len(updated_slides)-1:
        updated_slides[move_down], updated_slides[move_down+1] = updated_slides[move_down+1], updated_slides[move_down]
        sd["slides"] = updated_slides; st.session_state["plan_slide_data"] = sd; st.rerun()
    if delete_idx is not None and len(updated_slides) > 1:
        updated_slides.pop(delete_idx)
        sd["slides"] = updated_slides; st.session_state["plan_slide_data"] = sd; st.rerun()
    if insert_after is not None:
        updated_slides.insert(insert_after+1, {"type":"bullets","title":"New Slide","bullets":[]})
        sd["slides"] = updated_slides; st.session_state["plan_slide_data"] = sd; st.rerun()

    # Always persist live edits
    sd["slides"] = updated_slides
    st.session_state["plan_slide_data"] = sd

    st.divider()

    # ── Export ────────────────────────────────────────────────────────────────
    ec1, ec2, _ = st.columns([1, 1, 1])
    with ec1:
        if st.button("🚀 Export to PPTX", key="pm_export", type="primary", use_container_width=True):
            _tb = None
            _tf = st.session_state.get("template_upload")
            if _tf is not None:
                try:
                    _tf.seek(0); _tb = _tf.read()
                except Exception:
                    pass
            with st.spinner("Building PPTX…"):
                try:
                    pptx_out = generate_pptx(st.session_state["plan_slide_data"], template_bytes=_tb)
                    _title   = st.session_state["plan_slide_data"].get("title","Plan")
                    fname    = f"SEGA_Plan_{_title.replace(' ','_')[:40]}.pptx"
                    st.session_state["pptx_bytes"]    = pptx_out
                    st.session_state["pptx_filename"] = fname
                    st.session_state.pop("plan_mode_active", None)
                    st.success("PPTX ready — close plan to download.")
                    st.rerun()
                except Exception as _ex:
                    st.error(f"Export failed: {_ex}")
    with ec2:
        if "pptx_bytes" in st.session_state:
            st.download_button(
                "⬇️ Download PPTX",
                data=st.session_state["pptx_bytes"],
                file_name=st.session_state.get("pptx_filename","SEGA_Plan.pptx"),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="pm_dl",
                use_container_width=True,
            )
def run_pipeline(model, uploaded_files, game_title, business_question, audience,
                 theme_preset, web_search_en, slide_count, template_bytes=None,
                 data_files=None, plan_mode=False):
    """
    Generator yielding (event_type, payload) tuples consumed by the run-button
    handler.  Produces rich narrative log messages so users understand exactly
    what is happening at each stage.
    """
    api_key = st.secrets.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        yield ("error", "ANTHROPIC_API_KEY not found in st.secrets. "
                        "Add it to .streamlit/secrets.toml.")
        return

    headers = {
        "Content-Type": "application/json",
        "x-api-key": api_key,
        "anthropic-version": "2023-06-01",
    }

    # ── STAGE 1 + 2: document extraction & web research in parallel ──────────
    yield ("spinner", (
        "📂 <b>Stage 1 of 4 — Loading your documents</b><br>"
        "<span class='log-detail'>Reading uploaded files and extracting their text content. "
        "PDFs are parsed page-by-page; Excel/CSV sheets are converted to plain text tables; "
        "Word documents are unzipped and their paragraph text pulled out. "
        "Content is capped at {:,} characters to stay within API token limits.</span>"
    ).format(_MAX_DOC_CHARS))

    # Parse comma-separated game titles into a clean list
    game_titles = [g.strip() for g in (game_title or "").split(",") if g.strip()]
    game_title_display = ", ".join(game_titles) if game_titles else ""

    if web_search_en and game_titles:
        _n_games = len(game_titles)
        yield ("spinner", (
            "🔍 <b>Stage 2 of 4 — Web research running in parallel</b><br>"
            "<span class='log-detail'>While your documents are being extracted, Claude is "
            "simultaneously searching the web for <i>{}</i>. "
            "{}"
            "Each search has a 90s wall-clock deadline — falls back to model knowledge if it times out.</span>"
        ).format(
            game_title_display,
            f"{_n_games} parallel searches running. " if _n_games > 1 else ""
        ))

    combined_docs  = "[No documents uploaded]"
    research_text  = "[No reference game specified]"

    def _extract_docs():
        if not uploaded_files:
            return "[No documents uploaded]"
        texts = []
        for f in uploaded_files:
            f.seek(0)
            texts.append(f"=== {f.name} ===\n{extract_text_from_file(f)}")
        full = "\n\n".join(texts)
        # Trim to token-safe length
        if len(full) > _MAX_DOC_CHARS:
            full = full[:_MAX_DOC_CHARS] + "\n\n[... document content trimmed to fit token budget ...]"
        return full

    def _web_research_one(title: str) -> str:
        """
        Fetch competitive intel for a single game title.
        Runs the blocking httpx.post() in a daemon thread and enforces a hard
        wall-clock deadline via future.result(timeout=HARD_TIMEOUT). This fires
        even if the server is dribbling bytes and the socket never goes idle.
        """
        if not web_search_en:
            return f"[Web search disabled — using model knowledge for '{title}']"

        import httpx
        from concurrent.futures import ThreadPoolExecutor as _TPE, TimeoutError as _FTE

        HARD_TIMEOUT = 60   # true wall-clock seconds — future.result() fires this hard

        prompt = (
            f"Research the video game \"{title}\" for an executive competitive analysis presentation. "
            "Search the web for current information and write a structured report covering:\n\n"
            "OVERVIEW: Developer, publisher, release date, platforms, genre, launch price.\n\n"
            "CRITICAL RECEPTION: Metacritic score (critic + user), scores from 3-4 named outlets "
            "(IGN, GameSpot, Eurogamer etc). 4 specific things reviewers praised and 4 they criticised.\n\n"
            "COMMERCIAL PERFORMANCE: Launch window sales, lifetime sales if available, "
            "publisher statements, chart positions.\n\n"
            "GAMEPLAY & FEATURES: 5-6 core mechanics in 1-2 sentences each. "
            "Story length estimate. Any multiplayer or co-op features.\n\n"
            "POST-LAUNCH: Notable DLC or updates released.\n\n"
            "MARKET CONTEXT: 2-3 biggest competitor titles in the same window. "
            "How it compares to the previous franchise entry.\n\n"
            "Use real numbers throughout. Aim for 650-750 words total."
        )

        payload = {
            "model": "claude-haiku-4-5-20251001",  # Haiku is faster for web search
            "max_tokens": 2000,
            "tools": [{"type": "web_search_20250305", "name": "web_search"}],
            "messages": [{"role": "user", "content": prompt}],
        }

        def _do_fetch():
            """Single blocking HTTP call — no retries. Wall-clock enforced by future.result() outside."""
            try:
                resp = httpx.post(
                    "https://api.anthropic.com/v1/messages",
                    headers=headers,
                    json=payload,
                    timeout=httpx.Timeout(300, connect=10),
                )
                if resp.status_code == 429:
                    return "[Web search rate-limited — using training knowledge.]"
                resp.raise_for_status()
                blocks = resp.json().get("content", [])
                text = "\n".join(b.get("text", "") for b in blocks if b.get("type") == "text")
                return text.strip() or f"[No web results for '{title}']"
            except Exception as e:
                return f"[Web search error: {e}]"

        # Run in a single-thread executor so future.result(timeout=N) gives a true
        # wall-clock deadline regardless of socket activity.
        _ex = _TPE(max_workers=1)
        _fut = _ex.submit(_do_fetch)
        try:
            return _fut.result(timeout=HARD_TIMEOUT)
        except _FTE:
            _ex.shutdown(wait=False)
            return (
                f"[Web search timed out after {HARD_TIMEOUT}s. "
                f"Using Claude's training knowledge for '{title}' instead.]"
            )
        except Exception as e:
            return f"[Web search error: {e}]"
        finally:
            _ex.shutdown(wait=False)


    # ── Poll futures with heartbeat so Streamlit never blocks ────────────────
    # One future per game title (parallel searches) + one for doc extraction.
    # Poll every 3 s to keep Streamlit alive with heartbeat spinner ticks.
    # Hard wall-clock cap: if any game future exceeds RESEARCH_CAP seconds,
    # cancel it and move on — future.result(timeout) enforces this per-game,
    # but we also track total elapsed here as a belt-and-suspenders guard.
    RESEARCH_CAP = 100  # seconds — bail out of a stalled future
    _n_workers = 1 + max(len(game_titles), 1)
    with ThreadPoolExecutor(max_workers=_n_workers) as pool:
        fut_docs   = pool.submit(_extract_docs)
        game_futs  = {pool.submit(_web_research_one, t): t for t in game_titles}

        pending       = {fut_docs: "__docs__"}
        pending.update(game_futs)
        docs_done     = False
        research_done = not bool(game_titles)
        game_results  = {}   # title -> research text
        elapsed       = 0
        TICK          = 2   # poll every 2s for more responsive UI

        while pending:
            resolved = [(f, l) for f, l in list(pending.items()) if f.done()]

            # Hard cap: if any game future has been running longer than RESEARCH_CAP,
            # treat it as timed out and remove it from pending so we don't hang forever
            if elapsed >= RESEARCH_CAP:
                for fut, label in list(pending.items()):
                    if label != "__docs__" and not fut.done():
                        title_stuck = label
                        del pending[fut]
                        game_results[title_stuck] = (
                            f"[Web search exceeded {RESEARCH_CAP}s wall-clock limit for "
                            f"'{title_stuck}' — using Claude's training knowledge instead.]"
                        )
                        yield ("log",
                            f"⚠️ <b>Web search timed out for <i>{title_stuck}</i></b> — "
                            f"exceeded {RESEARCH_CAP}s hard limit<br>"
                            "<span class='log-detail'>Moving on with training knowledge.</span>")
                remaining_games = [l for l in pending.values() if l != "__docs__"]
                if not remaining_games:
                    research_done = True
                    if game_results:
                        parts = []
                        for t, r in game_results.items():
                            parts.append("=== RESEARCH: " + t + " ===\n" + r)
                        research_text = "\n\n".join(parts)
                    yield ("step_done", "research")

            for fut, label in resolved:
                del pending[fut]
                if label == "__docs__":
                    docs_done     = True
                    combined_docs = fut.result()
                    n_files = len(uploaded_files)
                    n_chars = len(combined_docs)
                    yield ("log", (
                        "✅ <b>Documents extracted</b> — {} file{}, {:,} characters of content<br>"
                        "<span class='log-detail'>Text successfully pulled from your uploads. "
                        "This content will be the primary source of facts about your internal game.</span>"
                    ).format(n_files, "s" if n_files != 1 else "", n_chars))
                    yield ("step_done", "extract")
                else:
                    # label is the game title string
                    res = fut.result()
                    game_results[label] = res
                    is_error    = "[Web search error" in res
                    is_timeout  = "timed out after" in res
                    is_fallback = any(x in res for x in (
                        "[Web search disabled", "[No reference", "[No web results",
                        "timed out after", "[No results"
                    ))
                    if is_error and not is_timeout:
                        yield ("log",
                            f"⚠️ <b>Web research error for <i>{label}</i></b><br>"
                            "<span class='log-detail'>Will use Claude's training knowledge instead.</span>")
                    elif is_fallback:
                        yield ("log",
                            f"⚠️ <b>Web search timed out for <i>{label}</i></b> — "
                            "falling back to model training knowledge<br>"
                            "<span class='log-detail'>Claude has extensive knowledge of most "
                            "released titles. The analysis will still be data-rich.</span>")
                    else:
                        word_count = len(res.split())
                        yield ("log",
                            f"✅ <b>Web research complete — <i>{label}</i></b> "
                            f"— ~{word_count} words<br>"
                            "<span class='log-detail'>Review scores, sales data, mechanics, "
                            "and market context compiled.</span>")
                    # Mark research step done once all game futures resolved
                    remaining_games = [l for l in pending.values() if l != "__docs__"]
                    if not remaining_games:
                        research_done = True
                        # Build combined research text
                        if game_results:
                            parts = []
                            for t, r in game_results.items():
                                parts.append("=== RESEARCH: " + t + " ===\n" + r)
                            research_text = "\n\n".join(parts)
                        yield ("step_done", "research")

            if pending:
                time.sleep(TICK)
                elapsed += TICK
                still_doing = []
                if not docs_done:
                    still_doing.append("extracting documents")
                pending_games = [l for l in pending.values() if l != "__docs__"]
                if pending_games:
                    still_doing.append(
                        "searching web for <i>" + ", ".join(pending_games) +
                        f"</i> ({elapsed}s elapsed)"
                    )
                if still_doing:
                    yield ("spinner",
                        "⏳ <b>Still working…</b> " + " &amp; ".join(still_doing) + "<br>"
                        "<span class='log-detail'>Web search has a 90s wall-clock deadline — "
                        "if it times out the pipeline falls back to training knowledge automatically."
                        "</span>")


    # ── STAGE 3: streaming analysis ───────────────────────────────────────────

    # Estimate input token count roughly (1 token ≈ 4 chars)
    prompt_chars = len(combined_docs) + len(research_text) + len(business_question) + 800
    est_input_tokens = prompt_chars // 4
    yield ("spinner", (
        "🤖 <b>Stage 3 of 4 — Claude is writing your presentation</b><br>"
        "<span class='log-detail'>Sending ~{:,} input tokens to {}. Claude will read all your "
        "document content, cross-reference the research on <i>{}</i>, interpret your business "
        "question, and produce a structured {}-slide JSON outline with titles, bullet points, "
        "comparison tables, stat callouts, and speaker notes. "
        "You will see progress updates as each section is written.</span>"
    ).format(est_input_tokens, model, game_title_display or "the reference game", slide_count))

    theme_desc = {
        "SEGA Blue — Corporate Executive": "Professional SEGA corporate blue (#0055AA), boardroom-ready.",
        "SEGA Dark — Game Reveal Style":   "Dark dramatic (#040A1C) with electric blue accents.",
        "SEGA Sonic — High Energy":        "Vibrant SEGA blue with gold accents, high energy.",
    }.get(theme_preset, "SEGA corporate blue")

    # ── DATA FILES: extract full data for chart context ──────────────────────
    # Send the complete dataset (all rows, all sheets) up to a token-safe cap.
    # Claude uses these exact numbers to populate chart.series values.
    _DATA_CHAR_CAP = 40_000  # ~10k tokens — enough for large sheets
    data_summary = ""
    if data_files:
        import pandas as _pd
        data_parts = []
        total_chars = 0
        for _df_file in (data_files or []):
            if total_chars >= _DATA_CHAR_CAP:
                data_parts.append(f"FILE: {_df_file.name} — skipped (data cap reached)")
                continue
            try:
                _df_file.seek(0)
                _dname = _df_file.name.lower()
                file_parts = []
                if _dname.endswith(".csv"):
                    _df = _pd.read_csv(_df_file)
                    _full_csv = _df.to_csv(index=False)
                    # Cap per-file if huge
                    if len(_full_csv) > _DATA_CHAR_CAP - total_chars:
                        _full_csv = _df.head(200).to_csv(index=False)
                        _full_csv += f"\n[...truncated — showing first 200 of {len(_df)} rows]"
                    file_parts.append(
                        f"Sheet: (csv)\n"
                        f"Rows: {len(_df)} | Columns: {list(_df.columns)}\n"
                        f"Numeric columns: {_df.select_dtypes(include='number').columns.tolist()}\n"
                        f"Full data:\n{_full_csv}"
                    )
                else:
                    # Excel — read every sheet
                    _sheets = _pd.read_excel(_df_file, sheet_name=None)
                    for _sheet_name, _df in _sheets.items():
                        _full_csv = _df.to_csv(index=False)
                        if len(_full_csv) > (_DATA_CHAR_CAP - total_chars) // max(len(_sheets), 1):
                            _full_csv = _df.head(200).to_csv(index=False)
                            _full_csv += f"\n[...truncated — showing first 200 of {len(_df)} rows]"
                        file_parts.append(
                            f"Sheet: {_sheet_name}\n"
                            f"Rows: {len(_df)} | Columns: {list(_df.columns)}\n"
                            f"Numeric columns: {_df.select_dtypes(include='number').columns.tolist()}\n"
                            f"Full data:\n{_full_csv}"
                        )
                sheet_block = "\n\n".join(file_parts)
                total_chars += len(sheet_block)
                data_parts.append(f"FILE: {_df_file.name}\n{sheet_block}")
            except Exception as _de:
                data_parts.append(f"FILE: {_df_file.name} — could not read: {_de}")
        if data_parts:
            data_summary = "\n\n".join(data_parts)
            total_rows = data_summary.count("\n")
            yield ("log", (
                f"📊 <b>Data loaded</b> — {len(data_files)} file(s), "
                f"{len(data_summary):,} characters of exact data sent to Claude for chart generation"
            ))


        analysis_prompt = f"""You are a senior game industry analyst at SEGA.
Analyse the following and produce a JSON object for a {slide_count}-slide executive presentation.

## INTERNAL GAME DOCUMENTS:
{combined_docs}

## REFERENCE GAME RESEARCH — {game_title_display}:
{research_text}

## BUSINESS QUESTION:
{business_question}

## AUDIENCE: {audience}

## UPLOADED DATA FOR CHARTS:
{data_summary if data_summary else "(no data files uploaded)"}

Output a single JSON object. Schema:
{{
  "title":"...", "subtitle":"...",
  "theme":{{"primary":"hex (dark-to-mid blue for SEGA branding, e.g. 0055AA)","accent":"hex (vivid cyan or teal, e.g. 00AADD)"}},
  "slides":[
    {{
      "type":"title|section|comparison|bullets|stats|recommendation|closing|chart",
      "title":"...","subtitle":"...","body":"...",
      "bullets":["..."],
      "stats":[{{"label":"...","value":"...","note":"..."}}],
      "comparison":{{
        "left_title":"Internal Game","right_title":"<Game Title>",
        "rows":[{{"label":"...","left":"...","right":"...","delta":"positive|negative|neutral"}}]
      }},
      "chart":{{
        "chart_type":"bar|line|scatter|pie|horizontal_bar",
        "title":"...","x_label":"...","y_label":"...",
        "categories":["..."],
        "series":[{{"label":"...","values":[0.0]}}],
        "colors":["hex"]
      }},
      "speaker_notes":"..."
    }}
  ]
}}

Rules:
- Use REAL data from the documents and research — no generic placeholders
- Be specific and data-driven for {audience}
- theme.primary and theme.accent must be dark-to-mid vivid hex colours (6 digits, no #).
  Never use white, near-white, or light pastels (no values above DDDDDD).
  Good: "0055AA", "003380", "00AADD". Bad: "FFFFFF", "F0F0F0", "C8D8EE".
- Keep speaker_notes to 1-2 sentences maximum — they are brief presenter cues, not essays
- Bullets: max 6 per slide, each under 15 words
- Comparison rows: max 8 per slide
- If multiple reference games are provided, you may produce one comparison slide per game, or a multi-column comparison slide covering all of them. Use the game title(s) as the right_title in comparison slides.
- Use "chart" type when the business question or uploaded data suggests a chart would be clearer than a table. Populate chart.series with real numeric data extracted from documents or research.
- Return ONLY valid JSON — no markdown fences, no explanation"""

    raw_chunks   = []
    char_count   = 0
    last_tick_at = 0

    # Slide detection: watch for "type": patterns in the accumulating JSON
    slides_announced = 0

    def _count_slides_so_far(text):
        """Count how many slide title fields have appeared so far in the stream."""
        import re
        return len(re.findall(r'"type"\s*:\s*"(?:title|section|bullets|stats|comparison|recommendation|closing)"', text))

    # Run the API stream in a background thread, feeding chunks into a queue.
    # The main generator polls that queue every 2s, yielding heartbeat ticks
    # between chunks — this way the UI updates even before the first token arrives.
    import queue as _queue, threading as _threading

    _chunk_q   = _queue.Queue()
    _SENTINEL  = object()   # signals stream is done
    _stream_err = [None]    # mutable container so thread can write to it

    def _stream_worker():
        """
        Makes a regular blocking POST (no SSE streaming) in a background thread.
        Puts the full response text as a single chunk, then the sentinel.
        This avoids httpx.stream() context manager issues in background threads.
        """
        import httpx
        try:
            # Use regular non-streaming POST — simpler and works reliably from threads
            resp = httpx.post(
                "https://api.anthropic.com/v1/messages",
                headers={**headers, "Content-Type": "application/json"},
                json={
                    "model": model,
                    "max_tokens": 8000,
                    "stream": False,
                    "system": "You are a precise game industry analyst. Return valid JSON only.",
                    "messages": [{"role": "user", "content": analysis_prompt}],
                },
                timeout=httpx.Timeout(300, connect=15),
            )
            if resp.status_code == 429:
                _stream_err[0] = RuntimeError(
                    "Rate limited. Switch to Haiku in the sidebar or wait a minute.")
                return
            if resp.status_code != 200:
                _stream_err[0] = RuntimeError(
                    f"API error {resp.status_code}: {resp.text[:400]}")
                return
            data = resp.json()
            text = ""
            for block in data.get("content", []):
                if block.get("type") == "text":
                    text += block.get("text", "")
            # Feed in small chunks so the progress counter updates as it "arrives"
            chunk_size = 120
            for i in range(0, len(text), chunk_size):
                _chunk_q.put(text[i:i+chunk_size])
        except Exception as exc:
            _stream_err[0] = exc
        finally:
            _chunk_q.put(_SENTINEL)

    _stream_thread = _threading.Thread(target=_stream_worker, daemon=True)
    _stream_thread.start()

    _stream_elapsed = 0
    _first_chunk    = True
    _done           = False
    _STREAM_TIMEOUT = 300  # hard wall-clock cap on the entire generation phase

    try:
        while not _done:
            # Poll the queue with a 2-second timeout so we can yield heartbeats
            try:
                item = _chunk_q.get(timeout=2)
            except _queue.Empty:
                # Nothing arrived in 2s — emit a heartbeat tick
                _stream_elapsed += 2
                if _stream_elapsed >= _STREAM_TIMEOUT:
                    yield ("error",
                        f"⏱️ <b>Generation timed out after {_STREAM_TIMEOUT}s.</b><br>"
                        "The API connection stalled. Try again — if it keeps happening, "
                        "switch to <b>Haiku</b> in the sidebar or reduce uploaded documents.")
                    return
                if _first_chunk:
                    yield ("spinner",
                        f"⏳ <b>Waiting for Claude…</b> {_stream_elapsed}s elapsed — "
                        f"generating ~{slide_count} slides from {est_input_tokens:,} tokens. "
                        "Typically takes 20–60s depending on model and load.")
                else:
                    pct = min(int(char_count / (slide_count * 220) * 100), 95)
                    yield ("spinner",
                        f"  📝 Generating… {char_count:,} chars written (~{pct}% complete) "
                        f"— {_stream_elapsed}s elapsed")
                continue

            if item is _SENTINEL:
                # Check if the worker raised an error
                if _stream_err[0] is not None:
                    raise _stream_err[0]
                _done = True
                break

            # Got a real chunk
            _first_chunk = False
            raw_chunks.append(item)
            char_count += len(item)

            # Announce new slides as they appear in the stream
            current_text  = "".join(raw_chunks)
            slides_so_far = _count_slides_so_far(current_text)
            if slides_so_far > slides_announced:
                for _ in range(slides_so_far - slides_announced):
                    slides_announced += 1
                    yield ("spinner",
                        f"  ✏️ Writing slide {slides_announced} of {slide_count}…")

            # Periodic byte-count tick every ~600 chars
            if char_count - last_tick_at >= 600:
                last_tick_at = char_count
                pct = min(int(char_count / (slide_count * 220) * 100), 95)
                yield ("spinner",
                    f"  📝 Generating JSON… {char_count:,} chars (~{pct}% complete)")

    except RuntimeError as e:
        err = str(e)
        if "rate limit" in err.lower() or "429" in err:
            yield ("error",
                "⏱️ <b>Rate limit hit</b> — your organisation has exceeded 30,000 input tokens/min.<br><br>"
                "Quick fixes:<br>"
                "• Switch to <b>Haiku</b> in the sidebar (much higher rate limits)<br>"
                "• Wait 60 seconds then click Run again<br>"
                "• Reduce the number of uploaded documents to lower input token usage"
            )
        else:
            yield ("error", err)
        return
    except Exception as e:
        yield ("error", f"Streaming error: {e}")
        return

    raw_json = "".join(raw_chunks).strip()
    if raw_json.startswith("```"):
        raw_json = raw_json.split("\n", 1)[1]
        raw_json = raw_json.rsplit("```", 1)[0]

    # ── Truncation recovery ───────────────────────────────────────────────
    # If the stream was cut off mid-JSON (max_tokens hit), try to salvage
    # whatever complete slides were already written before the truncation.
    def _recover_truncated_json(s: str) -> dict | None:
        """
        Try progressively shorter substrings until we get valid JSON,
        or attempt to close open brackets to recover partial output.
        Returns parsed dict or None.
        """
        import re
        # First: try as-is
        try: return json.loads(s)
        except json.JSONDecodeError: pass

        # Second: strip everything after the last complete slide object
        # by finding the last `}` before the slides array closes
        # Look for the last complete slide: ends with either },\n  { or }]\n}
        last_close = s.rfind('}\n    ]')  # end of slides array
        if last_close == -1:
            last_close = s.rfind('}, {')  # between slides
        if last_close == -1:
            last_close = s.rfind('"}')    # end of any string field

        if last_close > 100:
            # Try to close the JSON properly
            truncated = s[:last_close + 1]
            # Count unclosed brackets
            opens  = truncated.count('{') - truncated.count('}')
            opens2 = truncated.count('[') - truncated.count(']')
            candidate = truncated + (']' * opens2) + ('}' * opens)
            try: return json.loads(candidate)
            except json.JSONDecodeError: pass

        return None

    try:
        slide_data = json.loads(raw_json)
    except json.JSONDecodeError as e:
        # Try recovery before giving up
        recovered = _recover_truncated_json(raw_json)
        if recovered and recovered.get("slides"):
            n_recovered = len(recovered["slides"])
            yield ("log", (
                f"⚠️ <b>Output was truncated</b> — recovered {n_recovered} of {slide_count} slides.<br>"
                f"<span class='log-detail'>The model hit the token limit mid-stream. "
                f"Recovered {n_recovered} complete slides. Consider reducing slide count or switching "
                f"to a higher-capacity model.</span>"
            ))
            slide_data = recovered
        else:
            yield ("error",
                f"JSON parse error: {e}<br>"
                f"The model returned malformed JSON and recovery failed. Try running again, "
                f"or reduce the slide count in the sidebar.<br><br>"
                f"First 400 chars received:<br><code>{raw_json[:400]}</code>"
            )
            return

    n_slides = len(slide_data.get("slides", []))
    yield ("step_done", "analyze")
    yield ("log", (
        "✅ <b>Analysis complete</b> — {} slides generated, {:,} characters of content<br>"
        "<span class='log-detail'>Claude has finished writing all slide content including "
        "titles, bullet points, comparison rows, stat callouts, and speaker notes. "
        "Slide types used: {}.</span>"
    ).format(
        n_slides,
        char_count,
        ", ".join(sorted(set(s.get("type","?") for s in slide_data.get("slides",[]))))
    ))
    yield ("slide_data", slide_data)

    # ── PLAN MODE: stop here and let user edit the outline ────────────────────
    if plan_mode:
        yield ("plan_ready", slide_data)
        return

    # ── STAGE 4: PPTX rendering ───────────────────────────────────────────────
    _template_slide_count = ""
    if template_bytes:
        try:
            from pptx import Presentation as _PrsCheck
            _t = _PrsCheck(io.BytesIO(template_bytes))
            _template_slide_count = f" ({len(_t.slides)}-slide template)"
        except Exception:
            pass

    yield ("spinner", (
        "🖥️ <b>Stage 4 of 4 — Building the PowerPoint file</b><br>"
        "<span class='log-detail'>Rendering {} slides with python-pptx{}. "
        "Setting backgrounds, drawing shape layers, placing text with correct fonts, "
        "adding chrome footer with page numbers, writing to memory.</span>"
    ).format(
        n_slides,
        f" using your uploaded .pptx{_template_slide_count} as the base" if template_bytes else " with SEGA dark theme",
    ))

    try:
        pptx_bytes_out = generate_pptx(slide_data, template_bytes=template_bytes)
    except ValueError as e:
        # Template classification failed — fall back to SEGA dark scratch-build
        err_str = str(e)
        yield ("log", (
            f"⚠️ <b>Template could not be used</b> — falling back to SEGA dark theme.<br>"
            f"<span class='log-detail'>{err_str}</span>"
        ))
        try:
            pptx_bytes_out = generate_pptx(slide_data, template_bytes=None)
        except Exception as e2:
            yield ("error", f"PPTX generation error (fallback): {e2}")
            return
    except Exception as e:
        yield ("error", f"PPTX generation error: {e}")
        return

    yield ("step_done", "generate")
    yield ("log", (
        "🎉 <b>All done!</b> — PPTX is {:.0f} KB across {} slides<br>"
        "<span class='log-detail'>Your presentation is ready to download. "
        "Open it in PowerPoint or Google Slides. Speaker notes are included on each slide. "
        "The file uses standard OOXML format and is compatible with all modern presentation apps.</span>"
    ).format(len(pptx_bytes_out) / 1024, n_slides))
    yield ("pptx_bytes_out", pptx_bytes_out)

# ─── STREAMLIT UI ───

st.set_page_config(
    page_title="SEGA PowerPoint Creator",
    page_icon="🎮",
    layout="wide",
)


# ─────────────────────────────────────────────────────────────
# OTP AUTHENTICATION  (mirrors BIreport.py exactly)
# ─────────────────────────────────────────────────────────────

ALLOWED_DOMAIN     = "@segaamerica.com"
OTP_EXPIRY_SECS    = 600   # 10 minutes
COOKIE_EXPIRY_DAYS = 7
COOKIE_NAME        = "sega_analyzer_auth"








# ── Check for existing valid cookie ──────────────────────────
try:
    import extra_streamlit_components as stx
    _cookie_manager = stx.CookieManager(key="auth_cookies")
    _existing_cookie = _cookie_manager.get(COOKIE_NAME)
except Exception:
    _cookie_manager = None
    _existing_cookie = None

_cookie_email = _verify_cookie(_existing_cookie) if _existing_cookie else None

# ── Auth state init ───────────────────────────────────────────
for _k, _v in [
    ("auth_verified",   False),
    ("auth_email",      ""),
    ("otp_code",        ""),
    ("otp_email",       ""),
    ("otp_expiry",      0),
    ("otp_sent",        False),
    ("otp_attempts",    0),
]:
    if _k not in st.session_state:
        st.session_state[_k] = _v

# If valid cookie found, mark as verified
if _cookie_email and not st.session_state.auth_verified:
    st.session_state.auth_verified = True
    st.session_state.auth_email    = _cookie_email

# ── Render login gate if not verified ────────────────────────
if not st.session_state.auth_verified:
    st.markdown("""
    <style>
    .auth-wrap {
        max-width: 420px; margin: 5rem auto; padding: 2.5rem 2.5rem 2rem;
        background: var(--surface); border: 1px solid var(--border);
        border-top: 3px solid var(--blue); border-radius: 0 0 10px 10px;
    }
    .auth-logo  { font-family:'Arial Black',sans-serif; font-size:1.6rem; font-weight:900;
                  letter-spacing:0.12em; color:var(--blue) !important; margin-bottom:0.2rem; }
    .auth-title { font-size:1rem; font-weight:700; color:var(--text) !important; margin-bottom:0.25rem; }
    .auth-sub   { font-size:0.8rem; color:var(--muted) !important; margin-bottom:1.5rem; }
    .auth-note  { font-size:0.72rem; color:var(--muted) !important; margin-top:1rem;
                  text-align:center; line-height:1.5; }
    :root { --bg:#0a0c1a; --surface:#0f1120; --border:#232640; --blue:#4080ff; --text:#eef0fa; --muted:#5a5f82; }
    html, body, .stApp { background: var(--bg) !important; }
    </style>
    """, unsafe_allow_html=True)

    _lc, _mc, _rc = st.columns([1, 2, 1])
    with _mc:
        st.markdown("""
        <div class="auth-wrap">
          <div class="auth-logo">SEGA</div>
          <div class="auth-title">PowerPoint Creator</div>
          <div class="auth-sub">Sign in with your SEGA America email</div>
        </div>
        """, unsafe_allow_html=True)

        if not st.session_state.otp_sent:
            _email_input = st.text_input(
                "Email address",
                placeholder="you@segaamerica.com",
                label_visibility="hidden",
                key="auth_email_input",
            )
            _send_btn = st.button("Send verification code", use_container_width=True)

            if _send_btn and _email_input:
                if not _email_input.strip().lower().endswith(ALLOWED_DOMAIN):
                    st.error(f"Access restricted to {ALLOWED_DOMAIN} addresses.")
                else:
                    _code = str(random.randint(100000, 999999))
                    if _send_otp(_email_input.strip().lower(), _code):
                        st.session_state.otp_code     = _code
                        st.session_state.otp_email    = _email_input.strip().lower()
                        st.session_state.otp_expiry   = time.time() + OTP_EXPIRY_SECS
                        st.session_state.otp_sent     = True
                        st.session_state.otp_attempts = 0
                        st.rerun()

        else:
            st.info(f"Code sent to **{st.session_state.otp_email}** — check your inbox.")
            _code_input = st.text_input(
                "6-digit code",
                placeholder="123456",
                label_visibility="hidden",
                max_chars=6,
                key="auth_code_input",
            )
            _verify_btn = st.button("Verify code", use_container_width=True)

            if _verify_btn and _code_input:
                if st.session_state.otp_attempts >= 5:
                    st.error("Too many attempts. Please request a new code.")
                    st.session_state.otp_sent = False
                elif time.time() > st.session_state.otp_expiry:
                    st.error("Code has expired. Please request a new one.")
                    st.session_state.otp_sent = False
                elif _code_input.strip() != st.session_state.otp_code:
                    st.session_state.otp_attempts += 1
                    _remaining = 5 - st.session_state.otp_attempts
                    st.error(f"Incorrect code. {_remaining} attempt{'s' if _remaining != 1 else ''} remaining.")
                else:
                    st.session_state.auth_verified = True
                    st.session_state.auth_email    = st.session_state.otp_email
                    st.session_state.otp_code      = ""
                    if _cookie_manager:
                        _token = _sign_cookie(st.session_state.auth_email)
                        _cookie_manager.set(COOKIE_NAME, _token, expires_at=None, key="set_auth_cookie")
                    st.rerun()

            if st.button("← Use a different email", key="auth_back"):
                st.session_state.otp_sent = False
                st.session_state.otp_code = ""
                st.rerun()

        st.markdown(
            '<div class="auth-note">Restricted to @segaamerica.com addresses only.<br>'
            'Codes expire after 10 minutes.</div>',
            unsafe_allow_html=True,
        )

    st.stop()

# ── Signed-in user + sign out (sidebar) ──────────────────────
with st.sidebar:
    st.markdown(
        f'<div style="font-size:.6rem;font-weight:700;letter-spacing:.12em;text-transform:uppercase;'
        f'color:#5a5f82;margin-bottom:.3rem;">SEGA America</div>'
        f'<div style="font-size:.7rem;font-weight:600;color:#b8bcd4;margin-bottom:.75rem;">'
        f'{st.session_state.auth_email}</div>',
        unsafe_allow_html=True,
    )
    if st.button("Sign out", key="sign_out_btn"):
        if _cookie_manager:
            _cookie_manager.delete(COOKIE_NAME, key="delete_auth_cookie")
        for _k in ["auth_verified", "auth_email", "otp_sent", "otp_code",
                   "otp_email", "otp_expiry", "otp_attempts"]:
            st.session_state[_k] = False if _k == "auth_verified" else ""
        st.rerun()
    st.markdown("<hr style='border:none;border-top:1px solid #232640;margin:.5rem 0 1rem;'>", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────
# GLOBAL STYLES
# ─────────────────────────────────────────────────────────────
st.markdown("""
<style>
[data-testid="stToolbar"]      { display: none !important; }
[data-testid="stDecoration"]   { display: none !important; }
header[data-testid="stHeader"] { display: none !important; }

html, body, [class*="css"], .stApp {
    background-color: #0a0c1a !important;
    color: #e2e8f0 !important;
    font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif !important;
}
[data-testid="stSidebar"] { background: #0f172a !important; border-right: 1px solid #1e293b !important; }
[data-testid="stSidebar"] * { color: #cbd5e1 !important; }
[data-testid="stSidebar"] .block-container { padding: 1.25rem 1rem !important; max-width: 100% !important; }
.block-container { max-width: 1400px !important; padding: 1.5rem 2rem 3rem !important; margin: 0 auto !important; }

h1 { font-size: 2rem !important; font-weight: 700 !important; color: #f1f5f9 !important;
     letter-spacing: -.02em !important; margin-bottom: .15rem !important; }

.sidebar-section { font-size:.7rem; font-weight:600; text-transform:uppercase;
                   letter-spacing:.1em; color:#64748b; margin:1.2rem 0 .4rem; display:block; }

[data-testid="stFileUploader"] {
    border: 1px dashed #334155 !important;
    border-radius: 8px !important;
    background: transparent !important;
}

.stTextInput input, .stTextArea textarea {
    background: #1e293b !important;
    border: 1px solid #334155 !important;
    border-radius: 6px !important;
    color: #e2e8f0 !important;
    font-size: .85rem !important;
}

.section-label {
    font-size: .68rem; font-weight: 600; text-transform: uppercase;
    letter-spacing: .1em; color: #475569;
    margin: 2rem 0 .75rem;
    border-bottom: 1px solid #1e293b;
    padding-bottom: .4rem;
}

.status-card {
    background: #1e293b; border: 1px solid #334155; border-radius: 8px; padding: 1rem 1.25rem;
}
.status-card-label { font-size:.65rem; text-transform:uppercase; letter-spacing:.1em; color:#64748b; margin-bottom:.3rem; }
.status-card-value { font-size:.95rem; font-weight:600; color:#f1f5f9; line-height:1.4; }

.step-row { display:flex; align-items:center; gap:.6rem; padding:.45rem .75rem;
            margin:.2rem 0; border-radius:5px; font-size:.8rem; }
.step-done    { background:rgba(52,211,153,.1); color:#34d399; }
.step-active  { background:rgba(96,165,250,.12); color:#60a5fa; }
.step-pending { color:#475569; }

.result-log {
    background: #0f172a; border: 1px solid #1e293b; border-radius: 8px;
    padding: 1rem 1.25rem; font-size: .82rem;
    font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
    color: #94a3b8; max-height: 480px; overflow-y: auto; line-height: 1.7;
}
.result-log b  { color: #e2e8f0; font-weight: 600; }
.result-log i  { color: #60a5fa; }
.log-detail    { color: #64748b; font-size: .76rem; line-height: 1.55; display: block;
                 margin-top: .15rem; margin-bottom: .35rem; }
.log-entry     { border-left: 2px solid #1e3a5f; padding-left: .6rem;
                 margin-bottom: .5rem; }

/* ── Active spinner entry ── */
@keyframes spin { to { transform: rotate(360deg); } }
.log-active {
    border-left: 2px solid #00AADD;
    padding-left: .6rem;
    margin-bottom: .5rem;
    display: flex;
    align-items: flex-start;
    gap: .55rem;
}
.log-spinner {
    flex-shrink: 0;
    width: 14px; height: 14px;
    margin-top: 3px;
    border: 2px solid rgba(0,170,221,0.25);
    border-top-color: #00AADD;
    border-radius: 50%;
    animation: spin 0.8s linear infinite;
    display: inline-block;
}
.log-active-text { flex: 1; }
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────
# SIDEBAR — settings (post-auth)
# ─────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown("<div style='font-size:1rem;font-weight:700;color:#f1f5f9;margin-bottom:.25rem;'>PowerPoint Creator</div>", unsafe_allow_html=True)
    st.markdown("<div style='font-size:.65rem;color:#475569;text-transform:uppercase;letter-spacing:.1em;margin-bottom:1rem;'>Game PowerPoint Generator</div>", unsafe_allow_html=True)

    st.markdown('<span class="sidebar-section">Model</span>', unsafe_allow_html=True)
    model = st.selectbox(
        "Model", ["claude-sonnet-4-5", "claude-opus-4-5", "claude-haiku-4-5-20251001"],
        label_visibility="hidden",
    )

    st.markdown('<span class="sidebar-section">Options</span>', unsafe_allow_html=True)
    web_search_enabled = st.checkbox("Web search for reference game", value=True)
    plan_mode          = True  # always on — pipeline stops at outline for review
    slide_count        = st.slider("Target slides", 6, 20, 10)

    st.markdown('<span class="sidebar-section">Theme</span>', unsafe_allow_html=True)
    theme_preset = st.selectbox(
        "Theme",
        ["SEGA Blue — Corporate Executive", "SEGA Dark — Game Reveal Style", "SEGA Sonic — High Energy"],
        label_visibility="hidden",
    )

    template_file = st.file_uploader(
        "Upload .pptx template (optional)",
        type=["pptx"],
        label_visibility="hidden",
        help="Upload a branded .pptx file to extract its colour palette and fonts. "
             "Overrides the theme preset above.",
        key="template_upload",
    )
    if template_file:
        st.caption(f"✓ Template: {template_file.name}")

    st.markdown('<span class="sidebar-section">Pipeline</span>', unsafe_allow_html=True)
    pipeline_steps = st.session_state.get("pipeline_steps", {
        "upload": False, "extract": False, "research": False, "analyze": False, "generate": False,
    })
    for key, label in {"upload":"Document upload","extract":"Content extraction",
                        "research":"Web research","analyze":"AI analysis","generate":"PPTX generation"}.items():
        done = pipeline_steps.get(key, False)
        st.markdown(
            f'<div class="step-row {"step-done" if done else "step-pending"}">'
            f'{"✓" if done else "○"}&nbsp; {label}</div>',
            unsafe_allow_html=True,
        )

# ─────────────────────────────────────────────────────────────
# PAGE HEADER
# ─────────────────────────────────────────────────────────────
st.markdown(
    "<h1>Game PowerPoint Creator</h1>"
    "<div style='font-size:.78rem;color:#475569;margin-bottom:1.5rem;'>"
    "Upload internal documents &nbsp;·&nbsp; Benchmark against a released title &nbsp;·&nbsp; Generate a SEGA-branded PPTX"
    "</div>",
    unsafe_allow_html=True,
)

# ─────────────────────────────────────────────────────────────
# MAIN LAYOUT
# ─────────────────────────────────────────────────────────────
col_left, col_right = st.columns([1.1, 0.9], gap="large")

with col_left:
    st.markdown('<div class="section-label">Documents</div>', unsafe_allow_html=True)
    uploaded_files = st.file_uploader(
        "Upload internal game documents",
        type=["pdf", "xlsx", "xls", "csv", "txt", "docx"],
        accept_multiple_files=True,
        label_visibility="hidden",
    )
    if uploaded_files:
        st.caption(f"{len(uploaded_files)} file(s): " + ", ".join(f.name for f in uploaded_files))

    st.markdown('<div class="section-label" style="margin-top:.75rem;">Data for charts (optional)</div>', unsafe_allow_html=True)
    data_files = st.file_uploader(
        "Upload data files for chart generation",
        type=["xlsx", "xls", "csv"],
        accept_multiple_files=True,
        label_visibility="hidden",
        key="data_upload",
        help="Upload Excel or CSV files. Mention specific charts in the Business Question field.",
    )
    if data_files:
        st.caption(f"📊 {len(data_files)} data file(s): " + ", ".join(f.name for f in data_files))

    st.markdown('<div class="section-label">Analysis inputs</div>', unsafe_allow_html=True)

    game_title = st.text_input(
        "Reference / competitor game titles (comma-separated for multiple)",
        placeholder="e.g. Mario Odyssey, Astro Bot, Hollow Knight — separate multiple with commas",
    )
    business_question = st.text_area(
        "Business question",
        placeholder=(
            "e.g. Compare our internal game's combat mechanics and scope against Sonic Frontiers, "
            "highlighting gaps and opportunities for the executive team…"
        ),
        height=130,
    )
    audience = st.text_input(
        "Presentation audience", value="Executive team",
        placeholder="e.g. Executive team, Product leads, Marketing…",
    )

    run_btn = st.button("⚡ Run analysis", use_container_width=True, type="primary")

with col_right:
    st.markdown('<div class="section-label">Output</div>', unsafe_allow_html=True)
    output_area   = st.empty()
    download_area = st.empty()

    if st.session_state.get("plan_mode_active") and not run_btn:
        with output_area.container():
            _render_plan_modal(st.session_state.get("template_upload"))
    elif not run_btn and "pptx_bytes" not in st.session_state:
        output_area.markdown("""
<div class="status-card">
<div class="status-card-label">Ready</div>
<div class="status-card-value" style="color:#475569;font-size:.82rem;line-height:1.8;">
Fill in the inputs on the left and click <strong style="color:#e2e8f0;">Run analysis</strong>.<br><br>
The pipeline will:<br>
&nbsp;1. Extract your uploaded documents<br>
&nbsp;2. Search the web for the reference game<br>
&nbsp;3. Run a Claude-powered comparative analysis<br>
&nbsp;4. Show you the outline to review and edit<br>
&nbsp;5. Export to a SEGA-branded PPTX on demand
</div>
</div>
""", unsafe_allow_html=True)

    if "pptx_bytes" in st.session_state and not run_btn:
        download_area.download_button(
            label="⬇️ Download previous PPTX",
            data=st.session_state["pptx_bytes"],
            file_name=st.session_state.get("pptx_filename", "SEGA_Analysis.pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )

# ─────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────




# ─────────────────────────────────────────────────────────────
# PPTX GENERATION — pure python-pptx (no Node.js required)
# ─────────────────────────────────────────────────────────────


# Slide canvas: 13.3 × 7.5 inches (LAYOUT_WIDE)







# ── Slide type renderers ──────────────────────────────────────





























# ─────────────────────────────────────────────────────────────
# API HELPERS  (with rate-limit retry + back-off)
# ─────────────────────────────────────────────────────────────

_RL_WAIT_BASE = 20   # seconds to wait on first 429
_RL_MAX_TRIES = 5    # max retries before giving up






# ─────────────────────────────────────────────────────────────
# PIPELINE  (parallel extraction + research, streaming analysis)
# ─────────────────────────────────────────────────────────────

# Characters of document context to send.  Keeping this under ~8 000 chars
# (~2 000 tokens) leaves plenty of headroom on the 30 000 input-token/min limit.
_MAX_DOC_CHARS = 8_000




# ─────────────────────────────────────────────────────────────
# RUN BUTTON HANDLER
# ─────────────────────────────────────────────────────────────



if run_btn:
    if not business_question.strip():
        st.error("Please enter a business question.")
    else:
        # Collect data files
        data_files = st.session_state.get("data_upload") or []
        if hasattr(data_files, "read"):  # single file — wrap
            data_files = [data_files]

        # Extract template palette if a file was uploaded
        _template_bytes = None
        _template_file = st.session_state.get("template_upload")
        if _template_file is not None:
            try:
                _template_file.seek(0)
                _template_bytes = _template_file.read()
            except Exception as _e:
                st.warning(f"Could not read template file: {_e}. Using default theme.")

        pipeline_steps = {
            "upload": bool(uploaded_files), "extract": False,
            "research": False, "analyze": False, "generate": False,
        }
        st.session_state["pipeline_steps"] = pipeline_steps
        log_lines = []

        with output_area.container():
            st.markdown('<div class="section-label">Pipeline log</div>', unsafe_allow_html=True)
            log_area = st.empty()

            try:
                for event in run_pipeline(
                    model, uploaded_files, game_title, business_question,
                    audience, theme_preset, web_search_enabled, slide_count,
                    template_bytes=_template_bytes,
                    data_files=data_files,
                    plan_mode=plan_mode,
                ):
                    etype = event[0]

                    if etype in ("log", "spinner"):
                        # "spinner" = active working entry (animated ring at bottom)
                        # "log"     = completed entry (static, left-border style)
                        # When a new event arrives, any previous spinner is promoted
                        # to a plain log entry so only the latest one animates.
                        if etype == "spinner":
                            # Promote the last spinner (if any) to a plain log entry
                            if log_lines and log_lines[-1][0] == "spinner":
                                log_lines[-1] = ("log", log_lines[-1][1])
                            log_lines.append(("spinner", event[1]))
                        else:
                            # Promote any trailing spinner before adding the completed msg
                            if log_lines and log_lines[-1][0] == "spinner":
                                log_lines[-1] = ("log", log_lines[-1][1])
                            log_lines.append(("log", event[1]))

                        # Build self-contained HTML with inline CSS.
                        # Streamlit does NOT share CSS between separate markdown() calls,
                        # so all styles must be embedded in every render.
                        _CSS = (
                            "<style>"
                            "@keyframes _sp{to{transform:rotate(360deg)}}"
                            "._lw{background:#0f172a;border:1px solid #1e293b;border-radius:8px;"
                            "padding:1rem 1.25rem;font-size:.82rem;"
                            "font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;"
                            "color:#94a3b8;max-height:480px;overflow-y:auto;line-height:1.7}"
                            "._lw b{color:#e2e8f0;font-weight:600}"
                            "._lw i{color:#60a5fa}"
                            "._ld{color:#64748b;font-size:.76rem;line-height:1.55;display:block;"
                            "margin-top:.15rem;margin-bottom:.35rem}"
                            "._le{border-left:2px solid #1e3a5f;padding-left:.6rem;margin-bottom:.5rem}"
                            "._la{border-left:2px solid #00AADD;padding-left:.6rem;"
                            "margin-bottom:.5rem;display:flex;align-items:flex-start;gap:.55rem}"
                            "._lr{flex-shrink:0;width:14px;height:14px;margin-top:3px;"
                            "border:2px solid rgba(0,170,221,.25);border-top-color:#00AADD;"
                            "border-radius:50%;animation:_sp .8s linear infinite}"
                            "._lt{flex:1}"
                            "</style>"
                        )
                        html_parts = [_CSS, '<div class="_lw">']
                        for _kind, _text in log_lines[-14:]:
                            _t = _text.replace("class='log-detail'", "class='_ld'")
                            if _kind == "spinner":
                                html_parts.append(
                                    f'<div class="_la"><div class="_lr"></div>'
                                    f'<div class="_lt">{_t}</div></div>'
                                )
                            else:
                                html_parts.append(f'<div class="_le">{_t}</div>')
                        html_parts.append("</div>")
                        log_area.markdown("".join(html_parts), unsafe_allow_html=True)
                    elif etype == "step_done":
                        pipeline_steps[event[1]] = True
                        st.session_state["pipeline_steps"] = pipeline_steps

                    elif etype == "slide_data":
                        st.session_state["slide_data"] = event[1]

                    elif etype == "plan_ready":
                        st.session_state["plan_slide_data"] = event[1]
                        st.session_state["plan_mode_active"] = True

                    elif etype == "pptx_bytes_out":
                        fname = f"SEGA_Analysis_{(game_title_display or 'Report').replace(' ','_').replace(',','_')[:50]}.pptx"
                        st.session_state["pptx_bytes"]    = event[1]
                        st.session_state["pptx_filename"] = fname

                    elif etype == "error":
                        st.error(event[1])
                        break

            except Exception as ex:
                st.error(f"Unexpected error: {ex}")
                import traceback
                st.code(traceback.format_exc())

        if st.session_state.get("plan_mode_active"):
            with output_area.container():
                _render_plan_modal(st.session_state.get("template_upload"))

        if "pptx_bytes" in st.session_state:
            with download_area.container():
                st.success("Analysis complete.")
                st.download_button(
                    label="⬇️ Download PPTX",
                    data=st.session_state["pptx_bytes"],
                    file_name=st.session_state["pptx_filename"],
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )
                if "slide_data" in st.session_state:
                    with st.expander("Slide outline", expanded=False):
                        for i, sl in enumerate(st.session_state["slide_data"].get("slides", []), 1):
                            st.markdown(f"**{i}.** `{sl.get('type','?').upper()}` — {sl.get('title','')}")




# ─────────────────────────────────────────────────────────────
# PLAN MODE — OUTLINE EDITOR (modal dialog)
# ─────────────────────────────────────────────────────────────