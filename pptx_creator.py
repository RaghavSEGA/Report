"""
pptx_creator.py — SEGA PowerPoint Creator
A standalone Streamlit app for generating polished presentations on any topic.

Reuses the PPTX rendering engine from pptxruns.py (generate_pptx,
_render_plan_modal, extract_text_from_file etc.) — no duplication.

Auth: configurable domain via st.secrets["ALLOWED_DOMAIN"] or open to all.
"""

import streamlit as st
import io, time, hashlib, hmac, base64, json, re, os
import pypdf
import pandas as pd
from datetime import datetime, timedelta, timezone
from concurrent.futures import ThreadPoolExecutor, as_completed

# ── Reuse rendering engine from pptxruns ─────────────────────────────────────
import sys, importlib
sys.path.insert(0, os.path.dirname(__file__))

# We import just the pure functions — no Streamlit side-effects at import time
from pptxruns import (
    generate_pptx,
    extract_text_from_file,
    _make_anthropic_client,
)

# ── Storage ───────────────────────────────────────────────────────────────────
from storage_pptx import (
    init_db as _init_db,
    get_projects, project_exists, create_project,
    rename_project, delete_project, load_project, save_project,
)
_init_db()

# ── Constants ─────────────────────────────────────────────────────────────────
_MAX_DOC_CHARS     = 60_000
COOKIE_EXPIRY_DAYS = 1  # kept for storage_pptx compatibility

# ── Default template: SOA-HD_Template_Blue.pptx ──────────────────────────────
# Load from disk at startup so no upload is needed for the standard SEGA Blue look.
_DEFAULT_TEMPLATE_BYTES: bytes | None = None
_DEFAULT_TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "SOA-HD_Template_Blue.pptx")
if os.path.exists(_DEFAULT_TEMPLATE_PATH):
    with open(_DEFAULT_TEMPLATE_PATH, "rb") as _f:
        _DEFAULT_TEMPLATE_BYTES = _f.read()

# Auth: set ALLOWED_DOMAIN in secrets to restrict (e.g. "@mycompany.com")
# Leave blank or omit to allow any email address.
ALLOWED_DOMAIN = st.secrets.get("ALLOWED_DOMAIN", "")

# ── Purpose presets ───────────────────────────────────────────────────────────
PURPOSE_PRESETS = {
    "Executive briefing":   "Concise, data-driven. Lead with key takeaways. Use stats and comparison slides.",
    "Market analysis":      "Comprehensive research. Include market sizing, competitors, trends, opportunities.",
    "Project proposal":     "Problem → solution → plan → ask. Include timeline, budget, risks, success metrics.",
    "Sales deck":           "Hook early. Focus on value, outcomes, social proof. Close with clear CTA.",
    "Training material":    "Step-by-step. Clear objectives, examples, exercises, summaries per section.",
    "Research summary":     "Findings-first. Methodology, evidence, implications, limitations.",
    "Board update":         "High-level. KPIs, risks, decisions needed. Minimal detail.",
    "Product roadmap":      "Vision → now/next/later. Features, timeline, dependencies, success metrics.",
    "Investor pitch":       "Problem, solution, market, traction, team, ask. Narrative arc.",
    "General / Other":      "Balanced professional presentation suited to the audience and topic.",
}

# ── SEGA brand theme presets ──────────────────────────────────────────────────
THEME_PRESETS = {
    "SEGA Classic":      {"primary": "0033AA", "accent": "00CCFF"},   # deep blue + cyan
    "SEGA Midnight":     {"primary": "040A1C", "accent": "F5C242"},   # near-black + gold
    "SEGA Electric":     {"primary": "0055AA", "accent": "FFD700"},   # royal blue + gold
    "SEGA Stealth":      {"primary": "0A0F1E", "accent": "00CCFF"},   # dark navy + cyan
    "SEGA Prestige":     {"primary": "1A2A4A", "accent": "D0E4FF"},   # slate blue + ice
    "SEGA Arcade":       {"primary": "0033AA", "accent": "EE3355"},   # blue + red pop
    "SEGA Light":        {"primary": "0055AA", "accent": "0033AA"},   # blue on white
    "SEGA Monochrome":   {"primary": "1A1A2E", "accent": "8899BB"},   # near-black + muted
}

# ─────────────────────────────────────────────────────────────────────────────
# Auth — URL query-param token (no cookie library needed)
# Same pattern as documentcompare.py
# ─────────────────────────────────────────────────────────────────────────────

OTP_EXPIRY_SECS    = 600   # 10 minutes
TOKEN_EXPIRY_DAYS  = 1     # token survives 1 day in the URL

def _send_otp(email: str, code: str) -> bool:
    """Send OTP via AWS SES. Returns True on success."""
    try:
        import boto3
        from botocore.exceptions import ClientError
        ses = boto3.client(
            "ses",
            region_name=st.secrets.get("AWS_SES_REGION", "us-east-1"),
            aws_access_key_id=st.secrets.get("AWS_ACCESS_KEY_ID", ""),
            aws_secret_access_key=st.secrets.get("AWS_SECRET_ACCESS_KEY", ""),
        )
        ses.send_email(
            Source=st.secrets.get("EMAIL_FROM", "noreply@example.com"),
            Destination={"ToAddresses": [email]},
            Message={
                "Subject": {"Data": "SEGA PowerPoint Creator — Your verification code", "Charset": "UTF-8"},
                "Body": {
                    "Text": {
                        "Data": f"Your SEGA PowerPoint Creator verification code is: {code}\n\nExpires in 10 minutes.\nIf you didn't request this, you can safely ignore this email.",
                        "Charset": "UTF-8",
                    },
                    "Html": {
                        "Data": f"""
                        <div style="font-family:Arial,sans-serif;max-width:480px;margin:0 auto;
                                    background:#040A1C;border:1px solid #0033AA;border-top:3px solid #0055AA;
                                    border-radius:10px;overflow:hidden;">
                          <div style="padding:28px 32px 16px;border-bottom:1px solid rgba(0,102,204,0.3);">
                            <div style="font-family:'Arial Black',Arial,sans-serif;font-size:26px;
                                        font-weight:900;letter-spacing:.15em;color:#FFFFFF;">SEGA</div>
                            <div style="font-size:11px;letter-spacing:.3em;color:#00CCFF;
                                        text-transform:uppercase;margin-top:2px;">PowerPoint Creator</div>
                          </div>
                          <div style="padding:28px 32px 32px;">
                            <div style="font-size:14px;color:#D0E4FF;margin-bottom:20px;">
                              Your verification code is:
                            </div>
                            <div style="font-size:44px;font-weight:900;letter-spacing:.2em;color:#FFFFFF;
                                        background:rgba(0,85,170,0.35);border:1px solid rgba(0,204,255,0.3);
                                        border-radius:8px;padding:16px 24px;display:inline-block;
                                        margin-bottom:24px;font-family:'Arial Black',Arial,sans-serif;">
                              {code}
                            </div>
                            <div style="font-size:12px;color:#8899BB;line-height:1.6;">
                              Expires in 10 minutes.<br>
                              If you didn't request this, you can safely ignore this email.
                            </div>
                          </div>
                        </div>""",
                        "Charset": "UTF-8",
                    },
                },
            },
        )
        return True
    except Exception as e:
        st.error(f"Failed to send email: {e}")
        return False


def _make_token(email: str) -> str:
    """Create a signed URL token valid for TOKEN_EXPIRY_DAYS: base64(email|expiry|hmac)."""
    secret  = st.secrets.get("COOKIE_SIGNING_KEY", "fallback-change-this")
    expiry  = int(time.time()) + (TOKEN_EXPIRY_DAYS * 86400)
    payload = f"{email}|{expiry}"
    sig     = hmac.new(secret.encode(), payload.encode(), hashlib.sha256).hexdigest()
    return base64.urlsafe_b64encode(f"{payload}|{sig}".encode()).decode()


def _verify_token(token: str) -> str | None:
    """Verify token. Returns email if valid and unexpired, else None."""
    try:
        secret   = st.secrets.get("COOKIE_SIGNING_KEY", "fallback-change-this")
        decoded  = base64.urlsafe_b64decode(token.encode()).decode()
        email, expiry_str, sig = decoded.rsplit("|", 2)
        payload  = f"{email}|{expiry_str}"
        expected = hmac.new(secret.encode(), payload.encode(), hashlib.sha256).hexdigest()
        if not hmac.compare_digest(sig, expected):
            return None
        if int(time.time()) > int(expiry_str):
            return None
        return email
    except Exception:
        return None


# ── Read token from URL on every page load ────────────────────────────────────
_url_token   = st.query_params.get("t", "")
_token_email = _verify_token(_url_token) if _url_token else None

# ── Session state init ────────────────────────────────────────────────────────
for _k, _v in [
    ("auth_verified", False), ("auth_email", ""),
    ("auth_token", ""),
    ("otp_code", ""), ("otp_email", ""), ("otp_expiry", 0),
    ("otp_sent", False), ("otp_attempts", 0),
]:
    if _k not in st.session_state:
        st.session_state[_k] = _v

# If valid token in URL, auto-authenticate
if _token_email and not st.session_state.auth_verified:
    st.session_state.auth_verified = True
    st.session_state.auth_email    = _token_email
    st.session_state.auth_token    = _url_token

# ── Login gate ────────────────────────────────────────────────────────────────
if not st.session_state.auth_verified:
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Orbitron:wght@700;900&family=Rajdhani:wght@400;600&display=swap');
    .login-outer{
        min-height:80vh;display:flex;align-items:center;justify-content:center;
        background:linear-gradient(135deg,#040A1C 0%,#0D1B3E 50%,#040A1C 100%);
    }
    .login-wrap{
        max-width:420px;width:100%;padding:2.5rem 2rem;
        background:rgba(0,30,80,0.6);
        border:1px solid rgba(0,102,204,0.45);
        border-top:3px solid #0055AA;
        border-radius:12px;
        backdrop-filter:blur(12px);
        box-shadow:0 0 40px rgba(0,85,170,0.25);
    }
    .login-logo{
        font-family:'Orbitron',monospace;font-size:2.4rem;font-weight:900;
        letter-spacing:.15em;color:#FFFFFF;
        text-shadow:0 0 20px #0055AA,0 0 40px #00CCFF;
        margin-bottom:.2rem;text-align:center;
    }
    .login-sub{
        font-family:'Orbitron',monospace;font-size:.65rem;font-weight:700;
        letter-spacing:.35em;color:#00CCFF;text-transform:uppercase;
        text-align:center;margin-bottom:.35rem;
    }
    .login-divider{
        border:none;border-top:1px solid rgba(0,102,204,0.35);margin:1.2rem 0;
    }
    .login-title{
        font-family:'Rajdhani',sans-serif;font-size:1rem;font-weight:600;
        color:#D0E4FF;text-align:center;margin-bottom:1.5rem;
        letter-spacing:.04em;
    }
    </style>
    <div class="login-outer">
      <div class="login-wrap">
        <div class="login-logo">SEGA</div>
        <div class="login-sub">PowerPoint Creator</div>
        <hr class="login-divider"/>
        <div class="login-title">Sign in with your SEGA America email</div>
      </div>
    </div>
    """, unsafe_allow_html=True)

    if not st.session_state.otp_sent:
        with st.form("login_form"):
            _email_input = st.text_input("Email address", placeholder="you@yourcompany.com")
            _send_btn    = st.form_submit_button("Send verification code", use_container_width=True)

        if _send_btn and _email_input:
            if ALLOWED_DOMAIN and not _email_input.lower().endswith(ALLOWED_DOMAIN.lower()):
                st.error(f"Only {ALLOWED_DOMAIN} addresses are allowed.")
            else:
                _code = str(hash(time.time()) % 900000 + 100000)
                if _send_otp(_email_input.strip().lower(), _code):
                    st.session_state.otp_code     = _code
                    st.session_state.otp_email    = _email_input.strip().lower()
                    st.session_state.otp_expiry   = time.time() + OTP_EXPIRY_SECS
                    st.session_state.otp_sent     = True
                    st.session_state.otp_attempts = 0
                    st.rerun()
    else:
        st.info(f"Code sent to **{st.session_state.otp_email}** — check your inbox.")
        with st.form("otp_form"):
            _code_input = st.text_input("Enter 6-digit code", max_chars=6, key="auth_code_input")
            _verify_btn = st.form_submit_button("Verify code", use_container_width=True, type="primary")

        if _verify_btn and _code_input:
            if st.session_state.otp_attempts >= 5:
                st.error("Too many attempts. Please request a new code.")
                st.session_state.otp_sent = False
            elif time.time() > st.session_state.otp_expiry:
                st.error("Code expired. Please request a new one.")
                st.session_state.otp_sent = False
            elif _code_input.strip() != st.session_state.otp_code:
                st.session_state.otp_attempts += 1
                _rem = 5 - st.session_state.otp_attempts
                st.error(f"Incorrect code. {_rem} attempt{'s' if _rem != 1 else ''} remaining.")
            else:
                # Success — generate token and inject into URL so it survives refreshes
                st.session_state.auth_verified = True
                st.session_state.auth_email    = st.session_state.otp_email
                st.session_state.otp_code      = ""
                _token = _make_token(st.session_state.auth_email)
                st.session_state.auth_token    = _token
                st.query_params["t"] = _token
                st.rerun()

        _back_col, _ = st.columns([1, 1])
        with _back_col:
            if st.button("← Use a different email", key="auth_back"):
                st.session_state.otp_sent = False
                st.session_state.otp_code = ""
                st.rerun()

    if ALLOWED_DOMAIN:
        st.markdown(
            f"<div style='text-align:center;font-size:.72rem;color:#6080A8;margin-top:1rem'>"
            f"Restricted to {ALLOWED_DOMAIN} addresses · Codes expire after 10 minutes</div>",
            unsafe_allow_html=True,
        )
    st.stop()


# ─────────────────────────────────────────────────────────────────────────────
# Web research — generic, topic-aware
# ─────────────────────────────────────────────────────────────────────────────

def _web_research(topic: str, purpose: str, industry: str, question: str) -> tuple[str, list[dict]]:
    """Run a web search and return (research_text, sources).
    sources is a list of {"title": ..., "url": ...} dicts.

    The Anthropic web_search tool does not expose result URLs as structured data —
    they live only in the text Claude writes. So we instruct Claude to append a
    machine-parseable SOURCES block at the end of its response, then strip and
    parse it out here.
    """
    if not topic.strip():
        return "[No research topic specified]", []

    purpose_hint = PURPOSE_PRESETS.get(purpose, purpose)
    industry_hint = f" in the {industry} industry" if industry.strip() else ""

    prompt = (
        f'Research "{topic}"{industry_hint} for a presentation. '
        f'Presentation purpose: {purpose}. {purpose_hint}\n\n'
        f'Business question to address: {question or "(general overview)"}\n\n'
        "Write a structured research brief covering:\n"
        "OVERVIEW: What is this topic? Current state, key facts, relevant numbers.\n"
        "KEY PLAYERS / STAKEHOLDERS: Main organisations, people, or forces involved.\n"
        "DATA & EVIDENCE: Specific statistics, survey results, market figures, benchmarks.\n"
        "TRENDS: 3-4 notable recent developments or directions.\n"
        "CHALLENGES & OPPORTUNITIES: Main risks and areas for growth or improvement.\n"
        "CONTEXT: How does this compare to alternatives, competitors, or prior state?\n\n"
        "Use real, specific numbers wherever possible. Aim for 600-800 words.\n\n"
        "IMPORTANT: After the research brief, append a sources section in this exact format "
        "(do not skip this — it is required):\n\n"
        "SOURCES:\n"
        "- Title of page 1 | https://example.com/page1\n"
        "- Title of page 2 | https://example.com/page2\n"
        "(list every URL you actually retrieved content from, one per line)"
    )

    import concurrent.futures as _cf
    TIMEOUT = 90

    def _fetch():
        try:
            client = _make_anthropic_client()
            msg = client.messages.create(
                model="claude-haiku-4-5-20251001",
                max_tokens=2500,
                tools=[{"type": "web_search_20250305", "name": "web_search"}],
                messages=[{"role": "user", "content": prompt}],
            )

            full_text = "\n".join(
                b.text for b in msg.content
                if hasattr(b, "type") and b.type == "text" and b.text
            ) or "[No results returned]"

            # ── Parse out the SOURCES block ───────────────────────────────────
            sources = []
            seen_urls = set()
            research_text = full_text

            if "SOURCES:" in full_text:
                parts = full_text.split("SOURCES:", 1)
                research_text = parts[0].rstrip()
                sources_block = parts[1].strip()

                for line in sources_block.splitlines():
                    line = line.strip().lstrip("- ").strip()
                    if "|" in line:
                        title_part, url_part = line.rsplit("|", 1)
                        url   = url_part.strip()
                        title = title_part.strip() or url
                    elif line.startswith("http"):
                        url   = line
                        title = url
                    else:
                        continue

                    if url.startswith("http") and url not in seen_urls:
                        seen_urls.add(url)
                        sources.append({"title": title, "url": url})

            return research_text, sources

        except Exception as e:
            return f"[Web research error: {e}]", []

    with _cf.ThreadPoolExecutor(max_workers=1) as ex:
        fut = ex.submit(_fetch)
        try:
            return fut.result(timeout=TIMEOUT)
        except _cf.TimeoutError:
            return f"[Web research timed out after {TIMEOUT}s — using model knowledge]", []


# ─────────────────────────────────────────────────────────────────────────────
# Image fetching — uses Claude web_search to find an official image URL
# ─────────────────────────────────────────────────────────────────────────────

def _fetch_image_for_query(query: str) -> bytes | None:
    """
    Ask Claude (with web_search) to find a direct image URL from an official
    source for the query, then download and return the raw bytes.
    Explicitly avoids fan art, user-generated content, and social media.
    Returns None on any failure.
    """
    if not query or not query.strip():
        return None

    try:
        import urllib.request, re as _re

        client = _make_anthropic_client()

        msg = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=400,
            tools=[{"type": "web_search_20250305", "name": "web_search"}],
            messages=[{
                "role": "user",
                "content": (
                    f"Find an official image for: {query}\n\n"
                    "Requirements:\n"
                    "- Must be from an OFFICIAL source only: the game/product's official website, "
                    "publisher press kit, Wikipedia, official wiki, or official social media account.\n"
                    "- Absolutely NO fan art, fan sites, DeviantArt, ArtStation, Pinterest, Reddit, "
                    "Tumblr, Twitter/X fan accounts, Instagram fan accounts, or any user-generated content.\n"
                    "- Must be a direct URL ending in .jpg, .jpeg, .png, .webp, or .gif\n\n"
                    "Return ONLY the raw image URL on a single line. No markdown, no explanation."
                )
            }],
        )

        text = "\n".join(
            b.text for b in msg.content
            if hasattr(b, "type") and b.type == "text" and b.text
        ).strip()

        # Extract direct image URL — prefer explicit image extensions
        url_match = _re.search(
            r'https://[^\s\)\]"\'<>]+\.(?:jpg|jpeg|png|webp|gif)(?:\?[^\s\)\]"\'<>]*)?',
            text, _re.IGNORECASE
        )
        if not url_match:
            url_match = _re.search(r'https://[^\s\)\]"\'<>]+', text)
        if not url_match:
            return None

        img_url = url_match.group(0).rstrip('.,;)')

        # Block known fan art / UGC domains
        _BLOCKED = (
            "deviantart.com", "artstation.com", "pinterest.com", "reddit.com",
            "tumblr.com", "twitter.com", "x.com", "instagram.com",
            "pixiv.net", "furaffinity.net", "fandom.com/wiki",
        )
        if any(b in img_url.lower() for b in _BLOCKED):
            return None

        req = urllib.request.Request(
            img_url,
            headers={
                "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                              "AppleWebKit/537.36 (KHTML, like Gecko) "
                              "Chrome/120.0.0.0 Safari/537.36",
                "Accept": "image/webp,image/apng,image/*,*/*;q=0.8",
                "Referer": "https://www.google.com/",
            }
        )
        with urllib.request.urlopen(req, timeout=10) as resp:
            img_bytes = resp.read()

        # Validate: must be a real image
        if len(img_bytes) < 1000:
            return None
        magic = img_bytes[:4]
        if magic == b'\x89PNG':
            return img_bytes
        if magic[:3] == b'\xff\xd8\xff':
            return img_bytes
        if magic[:6] in (b'GIF87a', b'GIF89a'):
            return img_bytes
        if img_bytes[8:12] == b'WEBP':
            return img_bytes
        # JPEG2000 or other formats that still render fine
        if len(img_bytes) > 5000:
            return img_bytes

    except Exception as _e:
        import sys
        print(f"[image fetch] {query}: {_e}", file=sys.stderr)

    return None


# ─────────────────────────────────────────────────────────────────────────────
# Analysis prompt — fully generic
# ─────────────────────────────────────────────────────────────────────────────

def _build_analysis_prompt(
    topic: str,
    purpose: str,
    industry: str,
    audience: str,
    question: str,
    slide_count: int,
    doc_text: str,
    research_text: str,
    data_summary: str,
    theme: dict,
    sources: list | None = None,
) -> str:
    purpose_hint = PURPOSE_PRESETS.get(purpose, purpose)
    industry_ctx = f" in the {industry} industry" if industry.strip() else ""
    primary  = theme.get("primary", "1A3A6B")
    accent   = theme.get("accent",  "0099CC")

    sources_block = ""
    if sources:
        lines = "\n".join(f'  - {s["title"]}: {s["url"]}' for s in sources[:20])
        sources_block = f"\n## VERIFIED WEB SOURCES (include as final Sources slide):\n{lines}\n"

    return f"""You are an expert presentation strategist and analyst.
Create a JSON outline for a {slide_count}-slide professional presentation.

## TOPIC: {topic}{industry_ctx}
## PURPOSE: {purpose} — {purpose_hint}
## AUDIENCE: {audience}
## BUSINESS QUESTION / GOAL:
{question or "(No specific question — create a comprehensive overview)"}

## UPLOADED DOCUMENTS:
{doc_text}

## RESEARCH BRIEF:
{research_text}

## DATA FOR CHARTS:
{data_summary if data_summary else "(none uploaded)"}
{sources_block}
Output a single JSON object. Schema:
{{
  "title":"...", "subtitle":"...",
  "theme":{{"primary":"{primary}","accent":"{accent}"}},
  "slides":[
    {{
      "type":"title|section|bullets|stats|comparison|recommendation|chart|closing",
      "title":"...","subtitle":"...","body":"...",
      "bullets":["..."],
      "stats":[{{"label":"...","value":"...","note":"..."}}],
      "comparison":{{
        "left_title":"...","right_title":"...",
        "rows":[{{"label":"...","left":"...","right":"...","delta":"positive|negative|neutral"}}]
      }},
      "chart":{{
        "chart_type":"bar|line|scatter|pie|horizontal_bar",
        "title":"...","x_label":"...","y_label":"...",
        "categories":["..."],
        "series":[{{"label":"...","values":[0.0]}}],
        "colors":["hex"]
      }},
      "image_search_query":"concise image search query to find a real photo/logo for this slide, or empty string if no image needed",
      "source":"Source name — URL (the single most relevant source for the data on this slide, or empty string)",
      "speaker_notes":"..."
    }}
  ]
}}

Rules:
- Use REAL data from documents and research — no placeholders
- Tailor depth, tone, and structure to: {purpose} for {audience}
- Purpose hint: {purpose_hint}
- theme.primary and theme.accent must be vivid hex (6 digits, no #)
- speaker_notes: 1-2 sentences max — brief presenter cues only
- Bullets: max 6 per slide, each under 15 words
- Comparison rows: max 8 per slide
- Use "chart" type when data supports it
- image_search_query: provide a short, specific query (e.g. "PEAK game logo", "Sonic Frontiers screenshot") for slides where a real image would add value. Leave empty for data-heavy or abstract slides.
- source: cite the single most relevant source for each slide as "Publication Name — https://url". Use ONLY URLs from the VERIFIED WEB SOURCES list. Leave empty if no web source applies.
- STRICT TOPIC DISCIPLINE: Write ONLY about the stated topic. Do NOT introduce AI, machine learning, automation, or any technology trend unless it appears word-for-word in the uploaded documents or business question. If the topic is a game, genre, market, or business subject unrelated to AI, you MUST NOT mention AI anywhere in the slides. Violating this rule makes the deck unusable.
{f'- SOURCES SLIDE: The final slide MUST be type "bullets" titled "Sources", listing each source from the VERIFIED WEB SOURCES as a bullet: "Title — URL".' if sources else ''}
- Return ONLY valid JSON — no markdown fences, no explanation
"""


# ─────────────────────────────────────────────────────────────────────────────
# Conversational plan modal — chat-first outline editor
# ─────────────────────────────────────────────────────────────────────────────

_SLIDE_TYPES = ["title","section","bullets","stats","comparison","recommendation","chart","closing"]

def _render_plan_modal(template_bytes_ref=None, ns: str = "main"):
    """
    Chat-first plan editor.
    ns — namespace prefix for all widget keys, prevents duplicate-key errors
         when the modal is rendered in more than one tab.
    """
    sd     = st.session_state.get("plan_slide_data", {})
    slides = sd.get("slides", [])
    _active = st.session_state.get("active_project", "")

    # ── Header ────────────────────────────────────────────────────────────────
    hcol, xcol = st.columns([5, 1])
    with hcol:
        st.markdown(
            f"<h3 style='color:#D0E4FF;margin:0'>✏️ Review outline &nbsp;·&nbsp; "
            f"<span style='color:#00CCFF'>{sd.get('title','Presentation')}</span>"
            f"<span style='color:#4A6A9A;font-size:.8rem;font-weight:400'>"
            f" &nbsp;{len(slides)} slides</span></h3>",
            unsafe_allow_html=True,
        )
    with xcol:
        if st.button("✕ Close", key=f"{ns}_pm_close", use_container_width=True):
            st.session_state.pop("plan_slide_data", None)
            st.session_state.pop("plan_mode_active", None)
            st.rerun()

    st.divider()

    # ── Two-column layout: chat left, outline right ───────────────────────────
    chat_col, outline_col = st.columns([1, 1], gap="large")

    # ── LEFT: Conversational editor ───────────────────────────────────────────
    with chat_col:
        st.markdown(
            "<div style='font-size:.72rem;font-weight:700;letter-spacing:.1em;"
            "text-transform:uppercase;color:#00CCFF;margin-bottom:.5rem'>"
            "Chat to refine</div>",
            unsafe_allow_html=True,
        )

        if "plan_chat" not in st.session_state:
            st.session_state["plan_chat"] = []
        if "plan_slide_history" not in st.session_state:
            st.session_state["plan_slide_history"] = []

        # Welcome hint if no chat yet
        if not st.session_state["plan_chat"]:
            st.markdown(
                "<div style='color:#4A6A9A;font-size:.82rem;line-height:1.7;"
                "background:#0A1832;border-radius:6px;padding:.75rem 1rem;"
                "border-left:3px solid #0055AA;margin-bottom:.75rem'>"
                "Describe changes in plain language.<br>"
                "<b style='color:#8899BB'>Examples:</b><br>"
                "· \"Add a market size slide after slide 2\"<br>"
                "· \"Make slide 4 a comparison table vs competitors\"<br>"
                "· \"Swap slides 3 and 5, then tighten the bullets\"<br>"
                "· \"Remove the AI mentions from every slide\""
                "</div>",
                unsafe_allow_html=True,
            )

        # Chat history
        _chat_container = st.container(height=320)
        with _chat_container:
            for _msg in st.session_state["plan_chat"]:
                with st.chat_message(_msg["role"]):
                    st.markdown(_msg["content"])

        # Undo button
        _history = st.session_state["plan_slide_history"]
        undo_col, _ = st.columns([1, 3])
        with undo_col:
            if _history:
                if st.button(f"↩ Undo ({len(_history)})", key=f"{ns}_pm_undo", use_container_width=True):
                    st.session_state["plan_slide_data"]["slides"] = _history.pop()
                    st.session_state["plan_slide_history"] = _history
                    st.rerun()

        chat_input = st.chat_input("Describe what to change…", key=f"{ns}_plan_chat_input")

        if chat_input and chat_input.strip():
            current_slides = list(st.session_state["plan_slide_data"].get("slides", []))

            _system = f"""You are editing a PowerPoint presentation outline for SEGA America.
Respond conversationally — acknowledge what you're doing, explain any judgment calls, ask if anything is unclear.
Be direct and professional. No filler phrases.

SLIDE SCHEMA — every slide needs "type" and "title". Additional by type:
  title/closing/section: subtitle, body
  bullets/recommendation: bullets (list of strings, max 6)
  stats: stats (list of 4 x {{value, label, note}})
  comparison: comparison {{left_title, right_title, rows: [{{label,left,right,delta}}]}}
  chart: chart {{chart_type, x_label, y_label, categories, series:[{{label,values:[]}}], colors:[]}}
All slides may also have: image_search_query, source, speaker_notes

RESPONSE FORMAT — return exactly one JSON object (no markdown fences):
  Changes made: {{"action":"update","slides":[...complete array...],"message":"conversational explanation of what you did"}}
  Need clarification: {{"action":"message","text":"your question"}}

Rules:
- Always return the COMPLETE slides array when making changes.
- Do not introduce AI/ML/automation topics unless they are already in the outline.
- Keep slide count between 6 and 18.
- Return ONLY valid JSON."""

            _messages = []
            for _m in st.session_state["plan_chat"]:
                if _m["role"] == "user":
                    _messages.append({"role": "user", "content": _m["content"]})
                else:
                    _messages.append({"role": "assistant", "content": _m.get("_raw", _m["content"])})

            _messages.append({
                "role": "user",
                "content": (
                    f"Current outline ({len(current_slides)} slides):\n"
                    f"```json\n{json.dumps(current_slides, indent=2)}\n```\n\n"
                    f"Request: {chat_input.strip()}"
                )
            })

            with st.spinner("Updating outline…"):
                try:
                    client = _make_anthropic_client()
                    resp = client.messages.create(
                        model="claude-sonnet-4-6",
                        max_tokens=6000,
                        system=_system,
                        messages=_messages,
                    )
                    raw = "".join(
                        b.text for b in resp.content
                        if hasattr(b, "type") and b.type == "text"
                    ).strip()
                    if raw.startswith("```"):
                        raw = raw.split("\n", 1)[1].rsplit("```", 1)[0].strip()

                    parsed = json.loads(raw)

                    if parsed.get("action") == "update":
                        new_slides = parsed.get("slides", [])
                        message    = parsed.get("message", f"Updated to {len(new_slides)} slides.")
                        _history = st.session_state["plan_slide_history"]
                        _history.append(current_slides)
                        st.session_state["plan_slide_history"] = _history[-10:]
                        _sd = dict(st.session_state["plan_slide_data"])
                        _sd["slides"] = new_slides
                        st.session_state["plan_slide_data"] = _sd
                        display_reply = message
                    else:
                        display_reply = parsed.get("text", "No changes made.")
                        raw = json.dumps({"action": "message", "text": display_reply})

                    st.session_state["plan_chat"].append({"role": "user", "content": chat_input.strip()})
                    st.session_state["plan_chat"].append({"role": "assistant", "content": display_reply, "_raw": raw})

                    if _active:
                        _save_project(OWNER, _active)

                    st.rerun()
                except json.JSONDecodeError:
                    st.session_state["plan_chat"].append({"role": "user", "content": chat_input.strip()})
                    st.session_state["plan_chat"].append({"role": "assistant", "content": raw or "Could not parse response.", "_raw": raw or ""})
                    st.rerun()
                except Exception as e:
                    st.error(f"Chat error: {e}")

        if st.session_state["plan_chat"]:
            if st.button("🗑 Clear chat", key=f"{ns}_pm_clear_chat"):
                st.session_state["plan_chat"] = []
                st.rerun()

    # ── RIGHT: Live outline + expander editor ─────────────────────────────────
    with outline_col:
        st.markdown(
            "<div style='font-size:.72rem;font-weight:700;letter-spacing:.1em;"
            "text-transform:uppercase;color:#00CCFF;margin-bottom:.5rem'>"
            "Outline</div>",
            unsafe_allow_html=True,
        )

        # Quick slide list — always visible
        slides = st.session_state.get("plan_slide_data", {}).get("slides", [])
        _TYPE_ICON = {
            "title":"🎯","section":"📌","bullets":"📝","recommendation":"💡",
            "stats":"📊","comparison":"⚖️","chart":"📈","closing":"🏁",
        }
        for i, s in enumerate(slides):
            stype = s.get("type","bullets")
            st.markdown(
                f"<div style='font-size:.8rem;color:#8899BB;padding:.1rem 0'>"
                f"<span style='color:#4A6A9A;font-size:.72rem'>{i+1}.</span> "
                f"{_TYPE_ICON.get(stype,'📝')} "
                f"<span style='color:#D0E4FF'>{s.get('title','(untitled)')}</span>"
                f"</div>",
                unsafe_allow_html=True,
            )

        st.divider()

        # Expander: manual slide editor (secondary)
        with st.expander("✏️ Edit slides manually", expanded=False):
            updated_slides = list(slides)
            move_up = move_down = delete_idx = insert_after = None

            for i, slide in enumerate(slides):
                stype = slide.get("type", "bullets")
                with st.expander(f"**{i+1}.** {slide.get('title','(untitled)')}", expanded=False):
                    c1, c2 = st.columns([3, 2])
                    with c1:
                        new_type = st.selectbox("Type", _SLIDE_TYPES,
                            index=_SLIDE_TYPES.index(stype) if stype in _SLIDE_TYPES else 1,
                            key=f"{ns}_pm_type_{i}")
                    with c2:
                        b1,b2,b3,b4 = st.columns(4)
                        if b1.button("⬆", key=f"{ns}_pu_{i}", use_container_width=True): move_up = i
                        if b2.button("⬇", key=f"{ns}_pd_{i}", use_container_width=True): move_down = i
                        if b3.button("➕", key=f"{ns}_pa_{i}", use_container_width=True): insert_after = i
                        if b4.button("🗑", key=f"{ns}_px_{i}", use_container_width=True): delete_idx = i

                    new_title = st.text_input("Title", value=slide.get("title",""), key=f"{ns}_pm_ti_{i}")
                    new_sub   = st.text_input("Subtitle/body", value=slide.get("subtitle") or slide.get("body",""), key=f"{ns}_pm_su_{i}")
                    new_notes = st.text_input("Speaker notes", value=slide.get("speaker_notes",""), key=f"{ns}_pm_no_{i}")
                    new_src   = st.text_input("Source", value=slide.get("source",""), key=f"{ns}_pm_sr_{i}", placeholder="Publication — https://url")
                    new_img   = st.text_input("Image search query", value=slide.get("image_search_query",""), key=f"{ns}_pm_img_{i}")

                    new_slide = {**slide, "type": new_type, "title": new_title,
                                 "subtitle": new_sub, "speaker_notes": new_notes,
                                 "source": new_src, "image_search_query": new_img}

                    if new_type in ("bullets","recommendation"):
                        raw_b = st.text_area("Bullets (one per line)", value="\n".join(slide.get("bullets",[])), height=130, key=f"{ns}_pm_bu_{i}")
                        new_slide["bullets"] = [b.strip() for b in raw_b.split("\n") if b.strip()][:6]
                    elif new_type == "stats":
                        raw_s = st.text_area("Stats: value | label | note (4 rows)", height=100, key=f"{ns}_pm_st_{i}",
                            value="\n".join(f"{s.get('value','')} | {s.get('label','')} | {s.get('note','')}" for s in (slide.get("stats") or [{"value":"","label":"","note":""}]*4)[:4]))
                        new_slide["stats"] = [{"value":p[0],"label":p[1] if len(p)>1 else "","note":p[2] if len(p)>2 else ""} for line in raw_s.split("\n") if (p:=[x.strip() for x in line.split("|")]) and any(p)][:4]
                    elif new_type == "comparison":
                        cmp = slide.get("comparison") or {}
                        cl2,cr2 = st.columns(2)
                        lt = cl2.text_input("Left title", cmp.get("left_title",""), key=f"{ns}_pm_lt_{i}")
                        rt = cr2.text_input("Right title", cmp.get("right_title",""), key=f"{ns}_pm_rt_{i}")
                        raw_r = st.text_area("Rows: label | left | right | delta", height=130, key=f"{ns}_pm_ro_{i}",
                            value="\n".join(f"{r.get('label','')} | {r.get('left','')} | {r.get('right','')} | {r.get('delta','neutral')}" for r in (cmp.get("rows") or [])))
                        new_slide["comparison"] = {"left_title":lt,"right_title":rt,"rows":[{"label":p[0],"left":p[1] if len(p)>1 else "","right":p[2] if len(p)>2 else "","delta":p[3] if len(p)>3 else "neutral"} for line in raw_r.split("\n") if (p:=[x.strip() for x in line.split("|")]) and any(p)][:8]}

                    updated_slides[i] = new_slide

            # Apply manual reorder/insert/delete
            sd2 = dict(st.session_state["plan_slide_data"])
            if move_up is not None and move_up > 0:
                updated_slides[move_up-1], updated_slides[move_up] = updated_slides[move_up], updated_slides[move_up-1]
                sd2["slides"] = updated_slides; st.session_state["plan_slide_data"] = sd2; st.rerun()
            if move_down is not None and move_down < len(updated_slides)-1:
                updated_slides[move_down], updated_slides[move_down+1] = updated_slides[move_down+1], updated_slides[move_down]
                sd2["slides"] = updated_slides; st.session_state["plan_slide_data"] = sd2; st.rerun()
            if delete_idx is not None and len(updated_slides) > 1:
                updated_slides.pop(delete_idx)
                sd2["slides"] = updated_slides; st.session_state["plan_slide_data"] = sd2; st.rerun()
            if insert_after is not None:
                updated_slides.insert(insert_after+1, {"type":"bullets","title":"New Slide","bullets":[]})
                sd2["slides"] = updated_slides; st.session_state["plan_slide_data"] = sd2; st.rerun()
            # Persist live text edits
            sd2["slides"] = updated_slides
            st.session_state["plan_slide_data"] = sd2

    # ── Export bar ────────────────────────────────────────────────────────────
    st.divider()
    ex1, ex2, _ = st.columns([1, 1, 2])
    with ex1:
        if st.button("🚀 Export to PPTX", key=f"{ns}_pm_export", type="primary", use_container_width=True):
            _tb = st.session_state.get("saved_template_bytes") or _DEFAULT_TEMPLATE_BYTES
            _tf = template_bytes_ref
            if _tf is not None:
                try:
                    _tf.seek(0); _tb = _tf.read()
                except Exception:
                    pass
            with st.spinner("Building PPTX…"):
                try:
                    pptx_out = generate_pptx(st.session_state["plan_slide_data"], template_bytes=_tb)
                    pptx_out = _postprocess_pptx(pptx_out, st.session_state["plan_slide_data"].get("slides", []))
                    title = st.session_state["plan_slide_data"].get("title", "Plan")
                    fname = f"SEGA_{title.replace(' ','_')[:40]}.pptx"
                    st.session_state["pptx_bytes"]    = pptx_out
                    st.session_state["pptx_filename"] = fname
                    if _active:
                        _save_project(OWNER, _active)
                        st.toast(f'Saved to "{_active}"', icon="💾")
                    st.rerun()
                except Exception as ex:
                    st.error(f"Export failed: {ex}")
    with ex2:
        if st.session_state.get("pptx_bytes"):
            st.download_button(
                "⬇️ Download PPTX",
                data=st.session_state["pptx_bytes"],
                file_name=st.session_state.get("pptx_filename","presentation.pptx"),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key=f"{ns}_pm_dl", use_container_width=True,
            )


def _postprocess_pptx(pptx_bytes: bytes, slides_list: list) -> bytes:
    """
    Post-process a rendered PPTX:
    - Add a 'Source: ...' line at the bottom of slides that have one
    - Embed a real image in the right portion of eligible slides,
      preserving aspect ratio and not overwriting text content
    """
    has_images  = any(s.get("image_search_query","").strip() for s in slides_list)
    has_sources = any(s.get("source","").strip() for s in slides_list)
    if not has_images and not has_sources:
        return pptx_bytes

    try:
        from pptx import Presentation as _Prs
        from pptx.util import Inches as _In, Pt as _Pt, Emu as _Emu
        from pptx.dml.color import RGBColor as _RGB
        import io as _io

        prs2 = _Prs(_io.BytesIO(pptx_bytes))
        W = prs2.slide_width  / 914400   # slide width in inches  (e.g. 13.3)
        H = prs2.slide_height / 914400   # slide height in inches (e.g. 7.5)

        # SOA-HD template layout constants (match _generate_from_template)
        # Content area starts after the left sidebar and ends before source line
        CONTENT_X   = 0.95   # left edge of content area (after sidebar)
        CONTENT_TOP = 0.50   # top of content area (below title bar)
        CONTENT_BOT = H - 0.40  # bottom of content area (above chrome footer)

        # Image occupies the RIGHT portion of the content area
        IMG_FRAC   = 0.30    # fraction of content width used for image
        IMG_MARGIN = 0.15    # gap between text and image

        content_w = W - CONTENT_X - 0.15         # total content width
        img_w_max = content_w * IMG_FRAC          # max image width in inches
        img_h_max = CONTENT_BOT - CONTENT_TOP - 0.1  # max image height

        for slide, sdata in zip(prs2.slides, slides_list):

            # ── Source line ───────────────────────────────────────────────────
            src_text = (sdata.get("source") or "").strip()
            if src_text:
                try:
                    txb = slide.shapes.add_textbox(
                        _In(CONTENT_X), _In(H - 0.32), _In(W - CONTENT_X - 0.15), _In(0.25)
                    )
                    tf  = txb.text_frame
                    tf.word_wrap = False
                    run = tf.paragraphs[0].add_run()
                    run.text = f"Source: {src_text}"
                    run.font.size  = _Pt(7)
                    run.font.italic = True
                    run.font.color.rgb = _RGB(0x88, 0x99, 0xBB)
                    run.font.name  = "Calibri"
                except Exception:
                    pass

            # ── Image ─────────────────────────────────────────────────────────
            img_query = (sdata.get("image_search_query") or "").strip()
            stype = (sdata.get("type") or "").lower()
            if not img_query or stype in ("title", "closing", "chart"):
                continue

            img_bytes = _fetch_image_for_query(img_query)
            if not img_bytes:
                continue

            try:
                # Detect image dimensions to preserve aspect ratio
                from PIL import Image as _PIL_Image
                import io as _pil_io
                with _PIL_Image.open(_pil_io.BytesIO(img_bytes)) as _im:
                    px_w, px_h = _im.size
                src_ratio = px_w / px_h if px_h else 1.0
            except Exception:
                # PIL not available or can't decode — use 16:9 as safe default
                src_ratio = 16 / 9

            # Fit image into the allowed box preserving aspect ratio
            if src_ratio >= img_w_max / img_h_max:
                # Image is wider relative to box — constrain by width
                img_w = img_w_max
                img_h = img_w / src_ratio
            else:
                # Image is taller relative to box — constrain by height
                img_h = img_h_max
                img_w = img_h * src_ratio

            # Position: right side of content area, vertically centred in content band
            img_x = W - img_w - 0.15
            img_y = CONTENT_TOP + (img_h_max - img_h) / 2

            # Clip any text shapes whose right edge overlaps the image zone
            # Only clip shapes that are clearly in the content area (left > sidebar)
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                sh_left_in  = sh.left  / 914400
                sh_right_in = (sh.left + sh.width) / 914400
                # Only touch shapes that are in the content area and overlap
                if sh_left_in >= CONTENT_X - 0.1 and sh_right_in > img_x - IMG_MARGIN:
                    new_right = img_x - IMG_MARGIN
                    new_w_emu = max(int(0.5 * 914400),
                                   int((new_right - sh_left_in) * 914400))
                    sh.width = new_w_emu

            # Embed the image
            slide.shapes.add_picture(
                _io.BytesIO(img_bytes),
                _In(img_x), _In(img_y),
                _In(img_w), _In(img_h),
            )

        out = _io.BytesIO()
        prs2.save(out)
        return out.getvalue()

    except Exception as _e:
        import sys
        print(f"[postprocess error] {_e}", file=sys.stderr)
        return pptx_bytes


# ─────────────────────────────────────────────────────────────────────────────
# Pipeline — yields (event_type, payload) events
# ─────────────────────────────────────────────────────────────────────────────

def run_pipeline(
    model: str,
    uploaded_files: list,
    topic: str,
    purpose: str,
    industry: str,
    audience: str,
    question: str,
    web_search_en: bool,
    slide_count: int,
    theme: dict,
    template_bytes: bytes | None = None,
    data_files: list | None = None,
    plan_mode: bool = True,
):
    if not st.secrets.get("ANTHROPIC_API_KEY", ""):
        yield ("error", "ANTHROPIC_API_KEY not found in st.secrets.")
        return

    yield ("spinner", (
        "📂 <b>Stage 1 of 4 — Extracting your documents</b><br>"
        "<span class='log-detail'>Reading uploaded files and pulling out text content. "
        f"Content capped at {_MAX_DOC_CHARS:,} chars to stay within token limits.</span>"
    ))

    if web_search_en and topic.strip():
        yield ("spinner", (
            f"🔍 <b>Stage 2 of 4 — Web research on \"{topic}\" running in parallel</b><br>"
            "<span class='log-detail'>Searching for current data, stats, and context. "
            "90s deadline — falls back to model knowledge if it times out.</span>"
        ))

    doc_text     = "[No documents uploaded]"
    research_txt     = "[No web research — using model knowledge]"
    research_sources = []
    _file_stats  = []

    def _extract_docs():
        if not uploaded_files:
            return "[No documents uploaded]"
        parts = []
        for f in uploaded_files:
            f.seek(0)
            txt = extract_text_from_file(f)
            _file_stats.append((f.name, len(txt)))
            parts.append(f"=== {f.name} ===\n{txt}")
        full = "\n\n".join(parts)
        if len(full) > _MAX_DOC_CHARS:
            full = full[:_MAX_DOC_CHARS] + "\n\n[... trimmed ...]"
        return full

    def _do_research():
        if not web_search_en or not topic.strip():
            return "[Web search disabled]", []
        return _web_research(topic, purpose, industry, question)

    # ── Parallel: extract docs + web research ─────────────────────────────────
    with ThreadPoolExecutor(max_workers=2) as pool:
        fut_docs     = pool.submit(_extract_docs)
        fut_research = pool.submit(_do_research)
        pending = {fut_docs: "docs", fut_research: "research"}
        elapsed = 0
        TICK    = 2
        CAP     = 110

        while pending:
            done = [f for f in list(pending) if f.done()]
            for f in done:
                label = pending.pop(f)
                if label == "docs":
                    doc_text = f.result()
                    nf = len(uploaded_files) if uploaded_files else 0
                    yield ("log", (
                        f"✅ <b>Documents extracted</b> — {nf} file(s), "
                        f"{len(doc_text):,} chars<br>"
                        "<span class='log-detail'>" +
                        "".join(f"<br>&nbsp;&nbsp;· {n}: {c:,} chars"
                                for n, c in _file_stats) +
                        "</span>"
                    ))
                    yield ("step_done", "extract")
                else:
                    result = f.result()
                    if isinstance(result, tuple):
                        research_txt, research_sources = result
                    else:
                        research_txt, research_sources = result, []
                    ok = not research_txt.startswith("[")
                    yield ("log", (
                        f"{'✅' if ok else '⚠️'} <b>Web research</b> — "
                        f"{'~' + str(len(research_txt.split())) + ' words' if ok else 'used model knowledge'}<br>"
                        f"<span class='log-detail'>Topic: {topic}"
                        + (f" · {len(research_sources)} source(s) found" if research_sources else "")
                        + "</span>"
                    ))
                    yield ("step_done", "research")
                    if research_sources:
                        yield ("sources", research_sources)

            if pending:
                elapsed += TICK
                if elapsed >= CAP:
                    for f2 in list(pending):
                        pending.pop(f2)
                    research_txt = "[Research timed out]"
                    yield ("step_done", "research")
                    break
                still = list(pending.values())
                yield ("spinner", f"⏳ Still working: {', '.join(still)} ({elapsed}s)…")
                time.sleep(TICK)

    # ── Data files ─────────────────────────────────────────────────────────────
    _DATA_CAP   = 40_000
    data_summary = ""
    if data_files:
        parts = []
        total = 0
        for df_file in data_files:
            if total >= _DATA_CAP:
                parts.append(f"{df_file.name}: skipped (cap reached)")
                continue
            try:
                df_file.seek(0)
                name = df_file.name.lower()
                if name.endswith(".csv"):
                    df = pd.read_csv(df_file)
                    csv = df.to_csv(index=False)
                    if len(csv) > _DATA_CAP - total:
                        csv = df.head(200).to_csv(index=False) + "\n[truncated]"
                    parts.append(f"FILE: {df_file.name}\n{csv}")
                    total += len(csv)
                else:
                    sheets = pd.read_excel(df_file, sheet_name=None)
                    for sname, sdf in sheets.items():
                        csv = sdf.to_csv(index=False)
                        if len(csv) > (_DATA_CAP - total) // max(len(sheets), 1):
                            csv = sdf.head(200).to_csv(index=False) + "\n[truncated]"
                        parts.append(f"FILE: {df_file.name} / Sheet: {sname}\n{csv}")
                        total += len(csv)
            except Exception as e:
                parts.append(f"{df_file.name}: error — {e}")
        data_summary = "\n\n".join(parts)
        yield ("log", f"📊 <b>Data loaded</b> — {len(data_files)} file(s), {len(data_summary):,} chars")

    # ── Stage 3: streaming analysis ────────────────────────────────────────────
    prompt = _build_analysis_prompt(
        topic, purpose, industry, audience, question,
        slide_count, doc_text, research_txt, data_summary, theme,
        sources=research_sources,
    )
    est_tokens = len(prompt) // 4

    yield ("spinner", (
        f"🤖 <b>Stage 3 of 4 — Generating {slide_count}-slide outline</b><br>"
        f"<span class='log-detail'>Sending ~{est_tokens:,} tokens to {model}. "
        "Claude will synthesise your documents, research, and goal into a "
        "structured JSON slide plan with titles, bullets, stats, and speaker notes.</span>"
    ))

    import queue as _q, threading as _th

    chunk_q   = _q.Queue()
    SENTINEL  = object()
    err_box   = [None]

    def _stream():
        try:
            client = _make_anthropic_client()
            with client.messages.stream(
                model=model, max_tokens=8000,
                system="You are a precise presentation strategist. Return valid JSON only.",
                messages=[{"role": "user", "content": prompt}],
            ) as s:
                for chunk in s.text_stream:
                    chunk_q.put(chunk)
        except Exception as exc:
            err_box[0] = exc
        finally:
            chunk_q.put(SENTINEL)

    th = _th.Thread(target=_stream, daemon=True)
    th.start()

    raw_chunks   = []
    char_count   = 0
    last_tick    = 0
    first_chunk  = True
    stream_sec   = 0
    STREAM_LIMIT = 300
    slides_seen  = 0

    def _count_slides(txt):
        return len(re.findall(
            r'"type"\s*:\s*"(?:title|section|bullets|stats|comparison|recommendation|closing|chart)"',
            txt))

    done = False
    while not done:
        try:
            item = chunk_q.get(timeout=2)
        except _q.Empty:
            stream_sec += 2
            if stream_sec >= STREAM_LIMIT:
                yield ("error", f"Generation timed out after {STREAM_LIMIT}s.")
                return
            if first_chunk:
                yield ("spinner", f"⏳ Waiting for first token… {stream_sec}s")
            else:
                pct = min(int(char_count / (slide_count * 220) * 100), 95)
                yield ("spinner", f"  📝 Generating… {char_count:,} chars (~{pct}%)")
            continue

        if item is SENTINEL:
            if err_box[0]:
                raise err_box[0]
            done = True
            break

        first_chunk = False
        raw_chunks.append(item)
        char_count += len(item)

        cur = "".join(raw_chunks)
        n_sl = _count_slides(cur)
        if n_sl > slides_seen:
            for _ in range(n_sl - slides_seen):
                slides_seen += 1
                yield ("spinner", f"  ✏️ Writing slide {slides_seen} of {slide_count}…")

        if char_count - last_tick >= 600:
            last_tick = char_count
            pct = min(int(char_count / (slide_count * 220) * 100), 95)
            yield ("spinner", f"  📝 Generating… {char_count:,} chars (~{pct}%)")

    # ── Parse JSON ─────────────────────────────────────────────────────────────
    raw_json = "".join(raw_chunks).strip()
    if raw_json.startswith("```"):
        raw_json = raw_json.split("\n", 1)[1].rsplit("```", 1)[0]

    def _recover(s):
        try: return json.loads(s)
        except Exception: pass
        for pos in [s.rfind("}\n    ]"), s.rfind("}, {"), s.rfind('"}')]:
            if pos > 100:
                frag = s[:pos+1]
                opens  = frag.count("{") - frag.count("}")
                opens2 = frag.count("[") - frag.count("]")
                try: return json.loads(frag + "]"*opens2 + "}"*opens)
                except Exception: pass
        return None

    try:
        slide_data = json.loads(raw_json)
    except json.JSONDecodeError as e:
        recovered = _recover(raw_json)
        if recovered and recovered.get("slides"):
            n_rec = len(recovered["slides"])
            yield ("log", f"⚠️ Output truncated — recovered {n_rec} slides.")
            slide_data = recovered
        else:
            yield ("error", (
                f"JSON parse error: {e}<br>"
                f"First 400 chars: <code>{raw_json[:400]}</code>"
            ))
            return

    n_slides = len(slide_data.get("slides", []))
    yield ("step_done", "analyze")
    yield ("log", (
        f"✅ <b>Outline complete</b> — {n_slides} slides, "
        f"{char_count:,} chars<br>"
        f"<span class='log-detail'>Slide types: "
        f"{', '.join(sorted(set(s.get('type','?') for s in slide_data.get('slides',[]))))}"
        f"</span>"
    ))
    yield ("slide_data", slide_data)

    if plan_mode:
        yield ("plan_ready", slide_data)
        return

    # ── Stage 4: PPTX rendering ────────────────────────────────────────────────
    yield ("spinner", (
        f"🖥️ <b>Stage 4 of 4 — Building PPTX</b><br>"
        f"<span class='log-detail'>Rendering {n_slides} slides with python-pptx. "
        "Placing text, shapes, native charts, speaker notes.</span>"
    ))
    try:
        pptx_bytes = generate_pptx(slide_data, template_bytes=template_bytes)
    except Exception as e:
        yield ("error", f"PPTX render error: {e}")
        return

    # ── Post-process: embed images + source lines ──────────────────────────────
    slides_list = slide_data.get("slides", [])
    if any(s.get("image_search_query","").strip() for s in slides_list) or \
       any(s.get("source","").strip() for s in slides_list):
        yield ("spinner", "🖼️ <b>Fetching images and adding source lines…</b>")
    pptx_bytes = _postprocess_pptx(pptx_bytes, slides_list)

    yield ("step_done", "generate")
    yield ("log", (
        f"🎉 <b>Done!</b> — {len(pptx_bytes)//1024} KB, {n_slides} slides<br>"
        "<span class='log-detail'>Download below. Open in PowerPoint or Google Slides. "
        "Speaker notes included on every slide.</span>"
    ))
    yield ("pptx_bytes_out", pptx_bytes)


# ─────────────────────────────────────────────────────────────────────────────
# UI helpers
# ─────────────────────────────────────────────────────────────────────────────

def _render_log(log_lines: list) -> str:
    CSS = (
        "<style>"
        "@keyframes _sp{to{transform:rotate(360deg)}}"
        "._lw{background:#040A1C;border:1px solid #0A1832;border-radius:8px;"
        "padding:1rem 1.25rem;font-size:.82rem;"
        "font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;"
        "color:#8899BB;max-height:480px;overflow-y:auto;line-height:1.7}"
        "._lw b{color:#D0E4FF}._lw i{color:#00CCFF}"
        "._ld{color:#4A6A9A;font-size:.76rem;display:block;margin-top:.15rem}"
        "._le{border-left:2px solid #0A1A3A;padding-left:.6rem;margin-bottom:.5rem}"
        "._la{border-left:2px solid #0055AA;padding-left:.6rem;"
        "margin-bottom:.5rem;display:flex;align-items:flex-start;gap:.55rem}"
        "._lr{flex-shrink:0;width:14px;height:14px;margin-top:3px;"
        "border:2px solid rgba(0,85,170,0.25);border-top-color:#00CCFF;"
        "border-radius:50%;animation:_sp .8s linear infinite}"
        "._lt{flex:1}"
        "</style>"
    )
    parts = [CSS, '<div class="_lw">']
    for kind, text in log_lines[-14:]:
        t = text.replace("class='log-detail'", "class='_ld'")
        if kind == "spinner":
            parts.append(f'<div class="_la"><div class="_lr"></div><div class="_lt">{t}</div></div>')
        else:
            parts.append(f'<div class="_le">{t}</div>')
    parts.append("</div>")
    return "".join(parts)


# ─────────────────────────────────────────────────────────────────────────────
# Project management helpers
# ─────────────────────────────────────────────────────────────────────────────

def _load_project(owner: str, name: str):
    data = load_project(owner, name)
    if not data:
        return
    st.session_state["active_project"]       = name
    st.session_state["proj_topic"]           = data.get("business_question", "")
    st.session_state["proj_purpose"]         = data.get("audience", "General / Other")
    st.session_state["proj_industry"]        = data.get("industry", data.get("game_title", ""))
    st.session_state["proj_audience"]        = data.get("audience", "")
    st.session_state["project_doc_names"]    = data.get("doc_names", [])
    st.session_state["plan_slide_data"]      = data.get("slide_json") or {}
    st.session_state["plan_chat"]            = data.get("plan_chat", [])
    st.session_state["guided_messages"]      = data.get("guided_chat", [])
    st.session_state["research_sources"]     = data.get("sources", [])
    st.session_state["proj_web_research"]    = data.get("web_research", True)
    st.session_state["pptx_bytes"]           = data.get("pptx_bytes") or None
    st.session_state["saved_template_bytes"] = data.get("template_bytes") or None
    if st.session_state["pptx_bytes"] is None:
        st.session_state.pop("pptx_bytes", None)
    if st.session_state["saved_template_bytes"] is None:
        st.session_state.pop("saved_template_bytes", None)
    if st.session_state["plan_slide_data"]:
        st.session_state["plan_mode_active"] = True


def _save_project(owner: str, name: str):
    save_project(
        owner, name,
        business_question = st.session_state.get("proj_topic", ""),
        game_title        = st.session_state.get("proj_industry", ""),
        industry          = st.session_state.get("proj_industry", ""),
        audience          = st.session_state.get("proj_audience", ""),
        doc_names         = st.session_state.get("project_doc_names", []),
        slide_json        = st.session_state.get("plan_slide_data") or {},
        plan_chat         = st.session_state.get("plan_chat", []),
        guided_chat       = st.session_state.get("guided_messages", []),
        sources           = st.session_state.get("research_sources", []),
        web_research      = st.session_state.get("proj_web_research", True),
        pptx_bytes        = st.session_state.get("pptx_bytes"),
        template_bytes    = st.session_state.get("saved_template_bytes"),
    )


def _clear_project():
    for k in [
        "active_project", "proj_topic", "proj_purpose", "proj_industry",
        "proj_audience", "project_doc_names", "plan_slide_data", "plan_chat",
        "guided_messages", "guided_ready", "guided_params",
        "guided_pptx_bytes", "guided_pptx_filename",
        "research_sources", "proj_web_research",
        "pptx_bytes", "pptx_filename", "saved_template_bytes",
        "plan_mode_active", "plan_slide_history", "pipeline_steps",
    ]:
        st.session_state.pop(k, None)


# ─────────────────────────────────────────────────────────────────────────────
# Main app
# ─────────────────────────────────────────────────────────────────────────────

st.set_page_config(
    page_title="SEGA PowerPoint Creator",
    page_icon="🎮",
    layout="wide",
)

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Orbitron:wght@700;900&family=Rajdhani:wght@400;600&display=swap');
.section-label{font-size:.72rem;font-weight:700;letter-spacing:.1em;
    text-transform:uppercase;color:#00CCFF;margin:.9rem 0 .35rem;
    font-family:'Rajdhani',sans-serif;}
.sidebar-section{font-size:.65rem;font-weight:700;letter-spacing:.12em;
    text-transform:uppercase;color:#00CCFF;display:block;margin:.8rem 0 .3rem;
    font-family:'Rajdhani',sans-serif;}
.step-row{font-size:.78rem;padding:.18rem 0;color:#4A6A9A}
.step-done{color:#22DD88}
.step-pending{color:#6080A8}
.status-card{background:#0A1832;border:1px solid #0D2860;border-top:2px solid #0055AA;
    border-radius:8px;padding:1.25rem 1.5rem;margin-top:.5rem}
.status-card-label{font-size:.65rem;font-weight:700;letter-spacing:.12em;
    text-transform:uppercase;color:#00CCFF;margin-bottom:.5rem;
    font-family:'Rajdhani',sans-serif;}
</style>
""", unsafe_allow_html=True)

# Auth is handled inline above — st.stop() is called there if not verified.
OWNER = st.session_state.auth_email

# ── Sidebar ────────────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown(
        "<div style='font-family:\"Arial Black\",Arial,sans-serif;font-size:.65rem;"
        "letter-spacing:.2em;color:#00CCFF;text-transform:uppercase;"
        "margin-bottom:.1rem'>SEGA America</div>"
        f"<div style='font-size:.75rem;color:#8899BB;margin-bottom:1rem'>"
        f"{OWNER}</div>",
        unsafe_allow_html=True,
    )
    if st.button("Sign out", use_container_width=True):
        for k in ["auth_verified", "auth_email", "auth_token", "otp_sent", "otp_code"]:
            st.session_state.pop(k, None)
        st.query_params.clear()
        st.rerun()

    st.divider()

    # ── Project management ────────────────────────────────────────────────────
    st.markdown('<span class="sidebar-section">Project</span>', unsafe_allow_html=True)

    _projects = get_projects(OWNER)
    _proj_names = [p["name"] for p in _projects]
    _active = st.session_state.get("active_project", "")

    _select_opts = ["— select a project —"] + _proj_names
    _cur_idx = (_select_opts.index(_active) if _active in _select_opts else 0)
    _selected = st.selectbox("Project", _select_opts, index=_cur_idx,
                             label_visibility="hidden")

    if _selected != "— select a project —" and _selected != _active:
        _load_project(OWNER, _selected)
        st.query_params["project"] = _selected
        st.rerun()

    _new_name = st.text_input("New project name", placeholder="e.g. Q3 Investor Update",
                               label_visibility="hidden")
    if st.button("＋ Create project", use_container_width=True):
        if _new_name.strip():
            if project_exists(OWNER, _new_name.strip()):
                st.error("Name already exists.")
            else:
                create_project(OWNER, _new_name.strip())
                _clear_project()
                st.session_state["active_project"] = _new_name.strip()
                st.query_params["project"] = _new_name.strip()
                st.rerun()

    if _active:
        _c_save, _c_rename, _c_del = st.columns([3, 2, 1])
        with _c_save:
            if st.button("💾 Save", use_container_width=True):
                _save_project(OWNER, _active)
                st.toast(f'Saved "{_active}"', icon="✅")
        with _c_rename:
            if st.button("✏️", help="Rename", use_container_width=True):
                st.session_state["_renaming"] = True
        with _c_del:
            if st.button("🗑", help="Delete", use_container_width=True):
                st.session_state["_confirm_del"] = True

        if st.session_state.get("_renaming"):
            _rn = st.text_input("New name", value=_active, key="rename_input")
            if st.button("Confirm rename", key="confirm_rename"):
                if _rn.strip() and _rn.strip() != _active:
                    rename_project(OWNER, _active, _rn.strip())
                    st.session_state["active_project"] = _rn.strip()
                    st.session_state.pop("_renaming", None)
                    st.rerun()

        if st.session_state.get("_confirm_del"):
            st.warning(f'Delete "{_active}"?')
            _dc, _dk = st.columns(2)
            with _dc:
                if st.button("Yes, delete", key="del_yes"):
                    delete_project(OWNER, _active)
                    _clear_project()
                    st.session_state.pop("_confirm_del", None)
                    st.query_params.clear()
                    st.rerun()
            with _dk:
                if st.button("Cancel", key="del_no"):
                    st.session_state.pop("_confirm_del", None)
                    st.rerun()

        if st.session_state.get("pptx_bytes") or st.session_state.get("plan_slide_data"):
            if st.button("🔄 Reset output", use_container_width=True):
                for k in ["pptx_bytes","pptx_filename","plan_slide_data",
                          "plan_mode_active","plan_chat","plan_slide_history","pipeline_steps"]:
                    st.session_state.pop(k, None)
                save_project(OWNER, _active, slide_json={}, plan_chat=[], clear_pptx=True)
                st.rerun()

    st.divider()

    # ── Model ────────────────────────────────────────────────────────────────
    st.markdown('<span class="sidebar-section">Model</span>', unsafe_allow_html=True)
    model = st.selectbox(
        "Model", ["claude-sonnet-4-6", "claude-opus-4-6", "claude-haiku-4-5-20251001"],
        label_visibility="hidden",
    )
    st.session_state["_selected_model"] = model

    # ── Options ───────────────────────────────────────────────────────────────
    st.markdown('<span class="sidebar-section">Options</span>', unsafe_allow_html=True)
    _web_default = st.session_state.get("proj_web_research", True)
    web_search_enabled = st.checkbox("Web research", value=_web_default,
        help="Search the web for current data. Disable if your documents cover everything needed.",
        key="web_search_checkbox")
    # Persist preference immediately so it's saved with the project
    st.session_state["proj_web_research"] = web_search_enabled
    slide_count = st.slider("Target slides", 6, 25, 12)

    # ── Theme ────────────────────────────────────────────────────────────────
    st.markdown('<span class="sidebar-section">Theme</span>', unsafe_allow_html=True)
    theme_name = st.selectbox("Color theme", list(THEME_PRESETS.keys()),
                               label_visibility="hidden")
    theme = THEME_PRESETS[theme_name]

    # ── Template ────────────────────────────────────────────────────────────
    st.markdown('<span class="sidebar-section">Template (optional)</span>', unsafe_allow_html=True)
    template_file = st.file_uploader(
        "Upload .pptx template", type=["pptx"],
        label_visibility="hidden", key="template_upload",
        help="Upload a branded .pptx — layouts, fonts and colors will be preserved.",
    )

    # ── Pipeline status ─────────────────────────────────────────────────────
    st.divider()
    st.markdown('<span class="sidebar-section">Pipeline</span>', unsafe_allow_html=True)
    _pipeline = st.session_state.get("pipeline_steps", {
        "upload":False,"extract":False,"research":False,"analyze":False,"generate":False
    })
    for key, label in {
        "upload":"Document upload","extract":"Content extraction",
        "research":"Web research","analyze":"AI analysis","generate":"PPTX generation"
    }.items():
        done = _pipeline.get(key, False)
        st.markdown(
            f'<div class="step-row {"step-done" if done else "step-pending"}">'
            f'{"✓" if done else "○"}&nbsp; {label}</div>',
            unsafe_allow_html=True,
        )

# ── No-project gate ───────────────────────────────────────────────────────────
if not _active:
    # Check URL param
    _qp = st.query_params.get("project", "")
    if _qp and project_exists(OWNER, _qp):
        _load_project(OWNER, _qp)
        st.rerun()

    st.markdown("""
    <div style='max-width:520px;margin:4rem auto;text-align:center'>
      <div style='font-family:"Arial Black",Arial,sans-serif;font-size:2rem;font-weight:900;letter-spacing:.15em;color:#FFFFFF;text-shadow:0 0 15px #0055AA;margin-bottom:.5rem'>SEGA</div>
      <h2 style='color:#D0E4FF;margin-bottom:.5rem'>PowerPoint Creator</h2>
      <p style='color:#4A6A9A;margin-bottom:2rem'>
        AI-powered presentations for any topic — market analysis, project proposals,
        executive briefings, research summaries, and more.
      </p>
      <p style='color:#8899BB;font-size:.85rem'>
        Select an existing project or create a new one in the sidebar to get started.
      </p>
    </div>
    """, unsafe_allow_html=True)
    st.stop()

# ── Page header ────────────────────────────────────────────────────────────────
st.markdown(
    f"<h1>🎮 SEGA PowerPoint Creator</h1>"
    f"<div style='font-size:.78rem;color:#6080A8;margin-bottom:1.5rem'>"
    f"Project: <b>{_active}</b> &nbsp;·&nbsp; "
    f"Upload documents &nbsp;·&nbsp; Describe your goal &nbsp;·&nbsp; Generate a polished PPTX"
    f"</div>",
    unsafe_allow_html=True,
)

_tab_guided, _tab_main, _tab_pdf, _tab_transfer = st.tabs([
    "💬 Guided Build",
    "📊 Create Presentation",
    "📄 PDF → Editable PPTX",
    "🔄 Template Transfer",
])

with _tab_guided:
    st.markdown(
        "<div style='font-size:.85rem;color:#8899BB;margin-bottom:1rem'>"
        "Describe what you need and Claude will ask clarifying questions, "
        "then build your presentation when you're ready."
        "</div>",
        unsafe_allow_html=True,
    )

    # ── Session state init ────────────────────────────────────────────────────
    if "guided_messages" not in st.session_state:
        st.session_state["guided_messages"] = []
    if "guided_ready" not in st.session_state:
        st.session_state["guided_ready"] = False
    if "guided_params" not in st.session_state:
        st.session_state["guided_params"] = {}

    _gc_left, _gc_right = st.columns([1, 1], gap="large")

    with _gc_left:
        st.markdown('<div class="section-label">Conversation</div>', unsafe_allow_html=True)

        # Chat history
        _chat_box = st.container(height=420)
        with _chat_box:
            if not st.session_state["guided_messages"]:
                st.markdown(
                    "<div style='color:#4A6A9A;font-size:.85rem;padding:.5rem 0'>"
                    "👋 Hi! Tell me about the presentation you need — topic, "
                    "purpose, audience, anything you have in mind."
                    "</div>",
                    unsafe_allow_html=True,
                )
            for _msg in st.session_state["guided_messages"]:
                with st.chat_message(_msg["role"]):
                    st.markdown(_msg["content"])

        _guided_input = st.chat_input(
            "Tell me about your presentation…", key="guided_chat_input"
        )

        if _guided_input:
            st.session_state["guided_messages"].append(
                {"role": "user", "content": _guided_input}
            )

            _guide_system = """You are a presentation planning assistant for SEGA America.
Your role is to collect the minimum information needed to build a presentation brief, then confirm it.

Rules:
- Tone: professional, direct, concise. No filler. No enthusiasm markers. Do not say "Perfect", "Great", "Absolutely", "Sounds good", "Happy to help", or any similar phrase. Never compliment a choice.
- Ask one or two questions at a time. Never ask all questions at once.
- Do not volunteer opinions on the topic.

Information to collect (in order of priority):
1. Topic
2. Purpose (executive briefing, market analysis, project proposal, sales deck, etc.)
3. Audience (executive team, product team, external clients, etc.)
4. Slide count (assume 10 if not stated)
5. Specific sections, data points, or emphasis areas (optional)

Once you have topic + purpose + audience confirmed, output the plan summary THEN the signal line.
Never skip the plan summary. Never output the signal without the summary above it.

Format exactly:

**Topic:** [value]
**Purpose:** [value]
**Audience:** [value]
**Slides:** [value]
**Focus:** [the core question or goal this deck will answer]

READY_TO_GENERATE: {"topic": "...", "purpose": "...", "industry": "...", "audience": "...", "question": "...", "slide_count": 10}

Populate all JSON fields with real values from the conversation. Use empty string for anything not discussed. Do not wrap the JSON in markdown fences."""

            _guide_msgs = [
                {"role": m["role"], "content": m["content"]}
                for m in st.session_state["guided_messages"]
            ]

            try:
                _gc_client = _make_anthropic_client()
                _gc_resp = _gc_client.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=600,
                    system=_guide_system,
                    messages=_guide_msgs,
                )
                _reply = (
                    _gc_resp.content[0].text if _gc_resp.content else
                    "Sorry, I had trouble responding. Please try again."
                )
            except Exception as _ge:
                _reply = f"(Connection error: {_ge})"

            # Parse READY_TO_GENERATE signal
            _clean_reply = _reply
            if "READY_TO_GENERATE:" in _reply:
                try:
                    _json_line = (
                        _reply.split("READY_TO_GENERATE:")[1].strip().split("\n")[0]
                    )
                    _parsed_params = json.loads(_json_line)
                    st.session_state["guided_params"] = _parsed_params
                    st.session_state["guided_ready"] = True
                    # Strip the machine signal — keep everything before it
                    _clean_reply = _reply.split("READY_TO_GENERATE:")[0].strip()
                    # Always guarantee the plan summary is visible
                    _p = _parsed_params
                    _plan_block = (
                        f"**Topic:** {_p.get('topic', '')}\n"
                        f"**Purpose:** {_p.get('purpose', '')}\n"
                        f"**Audience:** {_p.get('audience', '')}\n"
                        f"**Slides:** {_p.get('slide_count', 10)}\n\n"
                        "Click **Generate presentation** on the right when ready."
                    )
                    if not _clean_reply:
                        _clean_reply = "Plan confirmed:\n\n" + _plan_block
                    elif not any(k in _clean_reply for k in ["**Topic:**", "**Purpose:**"]):
                        _clean_reply = _clean_reply + "\n\n" + _plan_block
                except Exception:
                    st.session_state["guided_ready"] = False


            st.session_state["guided_messages"].append(
                {"role": "assistant", "content": _clean_reply}
            )
            st.rerun()

        if st.button("🔄 Start over", key="guided_reset", use_container_width=True):
            for _k in [
                "guided_messages", "guided_ready", "guided_params",
                "guided_pptx_bytes", "guided_pptx_filename",
            ]:
                st.session_state.pop(_k, None)
            st.rerun()

    with _gc_right:
        st.markdown('<div class="section-label">Generate</div>', unsafe_allow_html=True)

        _gp      = st.session_state.get("guided_params", {})
        _g_ready = st.session_state.get("guided_ready", False)

        if _g_ready and _gp:
            st.success(
                f"**{_gp.get('topic', 'Your topic')}**  \n"
                f"{_gp.get('purpose', '')} · {_gp.get('audience', '')} · "
                f"{_gp.get('slide_count', 10)} slides"
            )
        else:
            st.markdown(
                "<div class='status-card'>"
                "<div class='status-card-label'>Waiting for plan</div>"
                "<div style='color:#6080A8;font-size:.82rem;line-height:1.9'>"
                "Chat with Claude on the left.<br>"
                "Once it has enough info it will confirm the plan<br>"
                "and the Generate button will unlock."
                "</div></div>",
                unsafe_allow_html=True,
            )

        st.markdown('<div class="section-label">Options</div>', unsafe_allow_html=True)
        _g_model = st.selectbox(
            "Model",
            ["claude-sonnet-4-6", "claude-opus-4-6", "claude-haiku-4-5-20251001"],
            key="guided_model",
        )
        _g_theme_name = st.selectbox(
            "Theme", list(THEME_PRESETS.keys()), key="guided_theme"
        )
        _g_theme = THEME_PRESETS[_g_theme_name]
        _g_web_default = st.session_state.get("proj_web_research", True)
        _g_web = st.checkbox("Web research", value=_g_web_default, key="guided_web",
            help="Disable if your uploaded documents cover everything needed.")
        st.session_state["proj_web_research"] = _g_web
        _g_template = st.file_uploader(
            "Custom template (.pptx)", type=["pptx"], key="guided_template"
        )
        _g_docs = st.file_uploader(
            "Supporting documents",
            type=["pdf", "docx", "txt", "csv", "xlsx"],
            accept_multiple_files=True,
            key="guided_docs",
        )

        _g_btn = st.button(
            "⚡ Generate presentation",
            use_container_width=True,
            type="primary",
            disabled=not _g_ready,
            key="guided_generate_btn",
        )

    # ── Full-width area below columns: pipeline log + plan modal + download ───
    # Must be outside _gc_left / _gc_right so _render_plan_modal gets full width
    _g_log_area    = st.empty()
    _g_plan_area   = st.container()
    _g_dl_area     = st.empty()

    # Persist plan modal across reruns
    if st.session_state.get("plan_mode_active") and not _g_btn:
        with _g_plan_area:
            _render_plan_modal(st.session_state.get("guided_template"), ns="guided")

    if st.session_state.get("guided_pptx_bytes") and not _g_btn:
        _g_dl_area.download_button(
            "⬇️ Download PPTX",
            data=st.session_state["guided_pptx_bytes"],
            file_name=st.session_state.get("guided_pptx_filename", "presentation.pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )

    if _g_btn and _g_ready and _gp:
        _g_template_bytes = _DEFAULT_TEMPLATE_BYTES
        if _g_template:
            _g_template.seek(0)
            _g_template_bytes = _g_template.read()

        _g_topic = _gp.get("topic", "")
        st.session_state["proj_topic"]        = _g_topic
        st.session_state["proj_purpose"]      = _gp.get("purpose", "")
        st.session_state["proj_industry"]     = _gp.get("industry", "")
        st.session_state["proj_audience"]     = _gp.get("audience", "")
        st.session_state["project_doc_names"] = [f.name for f in (_g_docs or [])]
        if _g_template_bytes:
            st.session_state["saved_template_bytes"] = _g_template_bytes

        _g_active = st.session_state.get("active_project", "")
        if not _g_active:
            _g_auto = re.sub(r"[^a-zA-Z0-9 ]+", "", _g_topic).strip()[:40] or "Untitled"
            _g_cand = _g_auto; _g_ctr = 2
            while project_exists(OWNER, _g_cand):
                _g_cand = f"{_g_auto} {_g_ctr}"; _g_ctr += 1
            create_project(OWNER, _g_cand)
            st.session_state["active_project"] = _g_cand
            _g_active = _g_cand

        _g_logs = []
        for _ev in run_pipeline(
            model=_g_model,
            uploaded_files=_g_docs or [],
            topic=_g_topic,
            purpose=_gp.get("purpose", "General / Other"),
            industry=_gp.get("industry", ""),
            audience=_gp.get("audience", "General audience"),
            question=_gp.get("question", ""),
            web_search_en=_g_web,
            slide_count=int(_gp.get("slide_count", 10)),
            theme=_g_theme,
            template_bytes=_g_template_bytes,
            plan_mode=True,
        ):
            _et = _ev[0]
            if _et in ("log", "spinner"):
                if _et == "spinner":
                    if _g_logs and _g_logs[-1][0] == "spinner":
                        _g_logs[-1] = ("log", _g_logs[-1][1])
                    _g_logs.append(("spinner", _ev[1]))
                else:
                    if _g_logs and _g_logs[-1][0] == "spinner":
                        _g_logs[-1] = ("log", _g_logs[-1][1])
                    _g_logs.append(("log", _ev[1]))
                _g_log_area.markdown(_render_log(_g_logs), unsafe_allow_html=True)
            elif _et == "sources":
                st.session_state["research_sources"] = _ev[1]
            elif _et == "plan_ready":
                st.session_state["plan_slide_data"]  = _ev[1]
                st.session_state["plan_mode_active"] = True
                _save_project(OWNER, _g_active)
                st.toast(f'Outline saved to "{_g_active}"', icon="📋")
            elif _et == "pptx_bytes_out":
                _slug = re.sub(r"[^a-zA-Z0-9]+", "_", _g_topic)[:50]
                st.session_state["guided_pptx_bytes"]    = _ev[1]
                st.session_state["guided_pptx_filename"] = f"Presentation_{_slug}.pptx"
                st.session_state["pptx_bytes"] = _ev[1]
                _save_project(OWNER, _g_active)
                st.toast(f'Auto-saved to "{_g_active}"', icon="💾")
            elif _et == "error":
                st.error(_ev[1], icon="🚨")
                break

        # Render plan modal full-width after pipeline
        if st.session_state.get("plan_mode_active"):
            with _g_plan_area:
                _render_plan_modal(st.session_state.get("guided_template"), ns="guided")

        _g_sources = st.session_state.get("research_sources", [])
        if _g_sources:
            with st.expander(f"🔗 {len(_g_sources)} verified source(s)", expanded=True):
                for _s in _g_sources:
                    st.markdown(f"- [{_s['title']}]({_s['url']})")


with _tab_pdf:
    st.markdown(
        "<div style='font-size:.85rem;color:#8899BB;margin-bottom:1rem'>"
        "Upload any PDF of slides and convert it to a fully editable PPTX — "
        "works even on image-only exports from Canva, Google Slides, or NotebookLM."
        "</div>",
        unsafe_allow_html=True,
    )
    _pdf_col, _pdf_out = st.columns([1, 1], gap="large")
    with _pdf_col:
        _pdf_file  = st.file_uploader("PDF", type=["pdf"],
                                       label_visibility="hidden", key="pdf_upload")
        _pdf_model = st.selectbox("Vision model",
                                   ["claude-opus-4-6", "claude-sonnet-4-6"],
                                   key="pdf_vision_model")
        _pdf_dpi   = st.select_slider("Quality (DPI)",
                                       options=[96,120,150,200], value=150, key="pdf_dpi")
        _pdf_btn   = st.button("⚡ Convert", use_container_width=True,
                                type="primary", disabled=not _pdf_file)
    with _pdf_out:
        _pdf_area = st.empty()
        if st.session_state.get("pdf_pptx_bytes"):
            _pdf_area.download_button(
                "⬇️ Download editable PPTX",
                data=st.session_state["pdf_pptx_bytes"],
                file_name=st.session_state.get("pdf_pptx_filename","converted.pptx"),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
    if _pdf_btn and _pdf_file:
        from pdf_to_pptx import pdf_to_editable_pptx as _pdf_convert
        _pdf_file.seek(0)
        _raw = _pdf_file.read()
        _pb  = _pdf_area.progress(0, text="Starting…")
        _la  = st.empty()
        _ll  = []
        def _plog(m):
            _ll.append(m)
            _la.markdown(
                "<div style='background:#040A1C;border-radius:6px;padding:.75rem;"
                "font-size:.8rem;color:#8899BB;font-family:monospace'>"
                + "".join(f"<div>{x}</div>" for x in _ll[-12:]) + "</div>",
                unsafe_allow_html=True)
        def _pprog(f):
            _pb.progress(min(f,0.99), text=f"Page {max(1,round(f*10))}…")
        _plog("🔍 Extracting and converting…")
        try:
            _out, _errs = _pdf_convert(
                _raw, api_key=st.secrets.get("ANTHROPIC_API_KEY",""),
                model=_pdf_model, dpi=_pdf_dpi, progress_cb=_pprog)
            for e in _errs: _plog(f"⚠️ {e}")
            _pb.progress(1.0, text="Done!")
            _fname = _pdf_file.name.rsplit(".",1)[0] + "_editable.pptx"
            st.session_state["pdf_pptx_bytes"]    = _out
            st.session_state["pdf_pptx_filename"] = _fname
            _pdf_area.download_button(
                "⬇️ Download editable PPTX", data=_out, file_name=_fname,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True)
            st.rerun()
        except Exception as _ex:
            st.error(f"Conversion failed: {_ex}")

with _tab_transfer:
    st.markdown(
        "<div style='font-size:.85rem;color:#8899BB;margin-bottom:1.25rem'>"
        "Re-skin an existing presentation into a new template. Upload your source "
        "<code>.pptx</code> and a target template — Claude will map every slide's "
        "content (text, bullets, speaker notes, chart titles) into the new layout, "
        "preserving the narrative while adopting the new brand."
        "</div>",
        unsafe_allow_html=True,
    )

    _tx_col1, _tx_col2 = st.columns(2, gap="large")
    with _tx_col1:
        st.markdown(
            "<div style='font-size:.72rem;font-weight:700;letter-spacing:.07em;"
            "text-transform:uppercase;color:#0055AA;margin-bottom:.4rem'>"
            "Source presentation</div>",
            unsafe_allow_html=True,
        )
        _tx_source = st.file_uploader(
            "Source PPTX", type=["pptx"],
            label_visibility="hidden", key="tx_source_upload",
        )
        st.markdown(
            "<div style='font-size:.72rem;font-weight:700;letter-spacing:.07em;"
            "text-transform:uppercase;color:#0055AA;margin:.9rem 0 .4rem'>"
            "Target template</div>",
            unsafe_allow_html=True,
        )
        _tx_template = st.file_uploader(
            "Target template PPTX", type=["pptx"],
            label_visibility="hidden", key="tx_template_upload",
        )
        st.markdown(
            "<div style='font-size:.72rem;font-weight:700;letter-spacing:.07em;"
            "text-transform:uppercase;color:#0055AA;margin:.9rem 0 .4rem'>"
            "Options</div>",
            unsafe_allow_html=True,
        )
        _tx_model = st.selectbox(
            "Model", ["claude-sonnet-4-6", "claude-opus-4-6", "claude-haiku-4-5-20251001"],
            key="tx_model",
        )
        _tx_preserve_charts = st.checkbox(
            "Preserve chart data", value=True,
            help="Keep chart categories and series values from the source.",
            key="tx_preserve_charts",
        )
        _tx_notes = st.checkbox(
            "Copy speaker notes", value=True,
            help="Transfer speaker notes to matching slides.",
            key="tx_notes",
        )
        _tx_btn = st.button(
            "⚡ Transfer to new template",
            use_container_width=True, type="primary",
            disabled=not (_tx_source and _tx_template),
            key="tx_run_btn",
        )

    with _tx_col2:
        st.markdown(
            "<div style='font-size:.72rem;font-weight:700;letter-spacing:.07em;"
            "text-transform:uppercase;color:#0055AA;margin-bottom:.4rem'>"
            "Output</div>",
            unsafe_allow_html=True,
        )
        _tx_out_area = st.empty()
        _tx_log_area = st.empty()

        if st.session_state.get("tx_pptx_bytes") and not _tx_btn:
            _tx_out_area.download_button(
                "⬇️ Download converted PPTX",
                data=st.session_state["tx_pptx_bytes"],
                file_name=st.session_state.get("tx_pptx_filename", "transferred.pptx"),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
        elif not _tx_btn:
            _tx_out_area.markdown(
                "<div class='status-card'>"
                "<div class='status-card-label'>How it works</div>"
                "<div style='color:#6080A8;font-size:.82rem;line-height:1.9'>"
                "1. Upload your existing presentation<br>"
                "2. Upload the target brand template<br>"
                "3. Claude extracts all content from the source<br>"
                "4. Remaps each slide into the best-matching layout<br>"
                "5. Delivers a fully styled PPTX in the new template"
                "</div></div>",
                unsafe_allow_html=True,
            )

    # ── Transfer handler ───────────────────────────────────────────────────────
    if _tx_btn and _tx_source and _tx_template:
        _tx_logs = []

        def _tx_log(msg: str, kind: str = "log"):
            _tx_logs.append((kind, msg))
            _tx_log_area.markdown(
                "<div style='background:#040A1C;border:1px solid #0A1832;border-radius:8px;"
                "padding:.85rem 1.1rem;font-size:.8rem;font-family:monospace;color:#8899BB;"
                "max-height:360px;overflow-y:auto;line-height:1.75'>"
                + "".join(
                    f"<div style='border-left:2px solid "
                    f"{'#0055AA' if k=='spin' else '#0A1A3A'};padding-left:.6rem;"
                    f"margin-bottom:.35rem'>{m}</div>"
                    for k, m in _tx_logs[-12:]
                )
                + "</div>",
                unsafe_allow_html=True,
            )

        try:
            import pptx as _pptx_lib
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.text import PP_ALIGN
            import copy, textwrap

            _tx_log("📂 Reading source presentation…", "spin")
            _tx_source.seek(0)
            _src_bytes = _tx_source.read()
            src_prs = _pptx_lib.Presentation(io.BytesIO(_src_bytes))

            _tx_log("🎨 Reading target template…", "spin")
            _tx_template.seek(0)
            _tmpl_bytes = _tx_template.read()

            # ── Extract content from source via Claude ─────────────────────────
            _tx_log("🤖 Extracting slide content with Claude…", "spin")

            # Build a text dump of the source
            # Patterns that identify footer/watermark/chrome text to skip
            _SKIP_PATTERNS = {
                "SEGA POWERPOINT CREATOR", "SEGA INTELLIGENCE ANALYZER",
                "POWERPOINT CREATOR", "INTELLIGENCE ANALYZER",
                "SEGA CONFIDENTIAL", "CONFIDENTIAL",
            }

            def _is_footer_shape(shape, slide_w, slide_h):
                """True if this shape looks like a footer/watermark (thin, near bottom or top edge)."""
                if not shape.has_text_frame:
                    return False
                txt = shape.text_frame.text.strip().upper()
                if not txt:
                    return True  # blank shape
                if txt in _SKIP_PATTERNS or any(p in txt for p in _SKIP_PATTERNS):
                    return True
                # Very thin shapes near slide bottom (bottom 10%) are usually footers
                try:
                    bottom = shape.top + shape.height
                    if shape.height < slide_h * 0.07 and bottom > slide_h * 0.88:
                        return True
                except Exception:
                    pass
                return False

            src_dump_parts = []
            for i, slide in enumerate(src_prs.slides, 1):
                lines = [f"SLIDE {i}"]
                sw = src_prs.slide_width
                sh = src_prs.slide_height
                for shape in slide.shapes:
                    if _is_footer_shape(shape, sw, sh):
                        continue
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                lines.append(f"  TEXT: {t}")
                    if hasattr(shape, "chart"):
                        try:
                            chart = shape.chart
                            lines.append(f"  CHART_TYPE: {chart.chart_type}")
                            lines.append(f"  CHART_TITLE: {chart.has_title and chart.chart_title.text_frame.text or ''}")
                            for series in chart.series:
                                vals = [str(v) for v in series.values]
                                lines.append(f"  SERIES: {series.name} = {', '.join(vals[:8])}")
                        except Exception:
                            pass
                if slide.has_notes_slide and _tx_notes:
                    notes_tf = slide.notes_slide.notes_text_frame
                    notes_t = notes_tf.text.strip() if notes_tf else ""
                    if notes_t:
                        lines.append(f"  NOTES: {notes_t}")
                src_dump_parts.append("\n".join(lines))

            src_dump = "\n\n".join(src_dump_parts)

            n_src_slides = len(src_prs.slides)

            xfer_prompt = f"""You are a presentation editor. Extract and rewrite the content of each source slide for a new context.

SOURCE SLIDES (total {n_src_slides}):
{src_dump}

For each source slide produce a JSON object with:
- "slide": slide number (1-based)
- "placeholders": dict where key "0" = slide title, key "1" = all body content as newline-separated lines
- "speaker_notes": verbatim speaker notes if present, else ""

Return ONLY a valid JSON array, no markdown, no explanation:
[
  {{
    "slide": 1,
    "placeholders": {{
      "0": "Title text",
      "1": "Line 1\\nLine 2\\nLine 3"
    }},
    "speaker_notes": ""
  }}
]

Rules:
- Keep ALL factual content — do not invent, summarise or drop data
- "0" must be the slide title only (short, no bullets)
- "1" contains all body text as newline-separated lines (strip leading bullet chars)
- speaker_notes: copy verbatim from source NOTES lines if present
"""
            client = _make_anthropic_client()
            xfer_resp = client.messages.create(
                model=_tx_model, max_tokens=8000,
                system="Return only valid JSON. No markdown fences.",
                messages=[{"role": "user", "content": xfer_prompt}],
            )
            raw_xfer = xfer_resp.content[0].text.strip()
            if raw_xfer.startswith("```"):
                raw_xfer = raw_xfer.split("\n", 1)[1].rsplit("```", 1)[0]

            slide_map = json.loads(raw_xfer)
            _tx_log(f"✅ Content mapped — {len(slide_map)} slides", "log")

            # ── Build output PPTX by cloning source slides + replacing text ────
            # Strategy: copy the source PPTX as a ZIP, then for each slide
            # do an XML text-node replacement of old→new content.
            # This preserves ALL backgrounds, images, shapes, and master graphics.
            _tx_log("🖥️ Building new PPTX…", "spin")

            import zipfile, copy, re as _re
            import lxml.etree as _etree

            # First build a per-slide old→new text replacement map from slide_map
            # We need to know what text was on each source slide so we can replace it
            # Re-extract source slide texts (same logic as above, already done in src_dump)
            # Collect ordered <a:t> text values per slide (keep bullets/prefix as-is)
            # These are used only for Claude's context — XML replacement is purely positional
            src_slide_texts = []
            _A_NS  = "http://schemas.openxmlformats.org/drawingml/2006/main"
            _T_TAG = f"{{{_A_NS}}}t"
            import zipfile as _zf, re as _re
            import lxml.etree as _etree

            src_zip_buf = io.BytesIO(_src_bytes)
            with _zf.ZipFile(src_zip_buf, "r") as _z:
                _all_names = _z.namelist()
                _slide_xmls = sorted(
                    [n for n in _all_names if _re.match(r"ppt/slides/slide\d+\.xml$", n)],
                    key=lambda x: int(_re.search(r"\d+", x.split("/")[-1]).group())
                )
                for sxml in _slide_xmls:
                    root = _etree.fromstring(_z.read(sxml))
                    nodes_on_slide = []
                    for node in root.iter(_T_TAG):
                        txt = (node.text or "").strip()
                        if not txt:
                            continue
                        up = txt.upper().lstrip("▸►•·▶➤➜–—- ")
                        if any(pat in up for pat in _SKIP_PATTERNS):
                            continue
                        nodes_on_slide.append(txt)
                    src_slide_texts.append(nodes_on_slide)

            # Open source as ZIP and clone it
            src_zip_buf  = io.BytesIO(_src_bytes)
            out_zip_buf  = io.BytesIO()

            # Helper: flatten Claude's placeholders into ordered list of new strings
            def _new_texts_for_entry(entry):
                ph = entry.get("placeholders", {})
                result = []
                for k in sorted(ph.keys(), key=lambda x: int(x) if x.isdigit() else 999):
                    for ln in str(ph[k]).split("\n"):
                        up = ln.strip().upper().lstrip("▸►•·▶➤➜–—- ")
                        if any(pat in up for pat in _SKIP_PATTERNS):
                            continue
                        if ln.strip():
                            result.append(ln.strip())
                return result

            def _replace_texts_in_xml(xml_bytes, new_texts):
                """
                Positional replacement: collect all non-empty, non-footer <a:t> nodes
                in document order, replace them 1-to-1 with new_texts.
                Extra old nodes beyond len(new_texts) are blanked.
                """
                try:
                    root = _etree.fromstring(xml_bytes)
                except Exception:
                    return xml_bytes

                t_nodes = []
                for node in root.iter(_T_TAG):
                    txt = (node.text or "").strip()
                    if not txt:
                        continue
                    up = txt.upper().lstrip("▸►•·▶➤➜–—- ")
                    if any(pat in up for pat in _SKIP_PATTERNS):
                        node.text = ""  # wipe footer
                        continue
                    t_nodes.append(node)

                for i, node in enumerate(t_nodes):
                    if i < len(new_texts):
                        node.text = new_texts[i]
                    else:
                        node.text = ""  # blank any extra old nodes

                return _etree.tostring(root, xml_declaration=True,
                                       encoding="UTF-8", standalone=True)

            with _zf.ZipFile(src_zip_buf, "r") as zin, \
                 _zf.ZipFile(out_zip_buf, "w", _zf.ZIP_DEFLATED) as zout:

                all_names = zin.namelist()
                slide_xmls = sorted(
                    [n for n in all_names if _re.match(r"ppt/slides/slide\d+\.xml$", n)],
                    key=lambda x: int(_re.search(r"\d+", x.split("/")[-1]).group())
                )

                for name in all_names:
                    data = zin.read(name)

                    if name in slide_xmls:
                        slide_i = slide_xmls.index(name)
                        if slide_i < len(slide_map):
                            entry     = slide_map[slide_i]
                            new_texts = _new_texts_for_entry(entry)
                            data = _replace_texts_in_xml(data, new_texts)

                    zout.writestr(name, data)

            # Now re-open the output ZIP to patch notes slides
            if any(entry.get("speaker_notes") for entry in slide_map):
                _tx_log("📝 Writing speaker notes…", "spin")
                out_zip_buf.seek(0)
                tmp_buf = io.BytesIO()
                with _zf.ZipFile(out_zip_buf, "r") as zin2, \
                     _zf.ZipFile(tmp_buf, "w", _zf.ZIP_DEFLATED) as zout2:
                    for name in zin2.namelist():
                        data = zin2.read(name)
                        m = _re.match(r"ppt/notesSlides/notesSlide(\d+)\.xml$", name)
                        if m:
                            idx = int(m.group(1)) - 1
                            if idx < len(slide_map):
                                nt = slide_map[idx].get("speaker_notes", "")
                                if nt and _tx_notes:
                                    try:
                                        root = _etree.fromstring(data)
                                        for node in root.iter(_T_TAG):
                                            if node.text and node.text.strip():
                                                node.text = nt
                                                nt = ""  # only replace first text node
                                        data = _etree.tostring(
                                            root, xml_declaration=True,
                                            encoding="UTF-8", standalone=True)
                                    except Exception:
                                        pass
                        zout2.writestr(name, data)
                out_zip_buf = tmp_buf

            out_bytes = out_zip_buf.getvalue()

            # ── Save output ────────────────────────────────────────────────────
            _tx_fname = _tx_source.name.rsplit(".", 1)[0] + "_transferred.pptx"
            st.session_state["tx_pptx_bytes"]    = out_bytes
            st.session_state["tx_pptx_filename"] = _tx_fname

            _tx_log(f"✅ Done — {len(slide_map)} slides transferred", "log")
            _tx_out_area.download_button(
                "⬇️ Download converted PPTX",
                data=out_bytes,
                file_name=_tx_fname,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )

        except json.JSONDecodeError as _je:
            st.error(f"Failed to parse Claude's layout mapping: {_je}")
        except Exception as _te:
            import traceback
            st.error(f"Template transfer failed: {_te}")
            st.code(traceback.format_exc())


with _tab_main:
    col_left, col_right = st.columns([1.1, 0.9], gap="large")

    with col_left:
        st.markdown('<div class="section-label">Documents (optional)</div>', unsafe_allow_html=True)
        uploaded_files = st.file_uploader(
            "Upload supporting documents",
            type=["pdf","xlsx","xls","csv","txt","docx"],
            accept_multiple_files=True, label_visibility="hidden",
        )
        if uploaded_files:
            st.caption(f"{len(uploaded_files)} file(s): " + ", ".join(f.name for f in uploaded_files))
            st.session_state["project_doc_names"] = [f.name for f in uploaded_files]

        st.markdown('<div class="section-label">Data for charts (optional)</div>', unsafe_allow_html=True)
        data_files = st.file_uploader(
            "Upload data files", type=["xlsx","xls","csv"],
            accept_multiple_files=True, label_visibility="hidden", key="data_upload",
            help="Excel or CSV files — mention charts in the goal field below.",
        )

        st.markdown('<div class="section-label">Presentation details</div>', unsafe_allow_html=True)

        topic = st.text_input(
            "Topic / subject",
            value=st.session_state.get("proj_topic", ""),
            placeholder="e.g. EV market landscape, Q3 earnings review, office relocation proposal",
            key="topic_input",
        )
        st.session_state["proj_topic"] = topic

        purpose = st.selectbox(
            "Purpose",
            list(PURPOSE_PRESETS.keys()),
            index=list(PURPOSE_PRESETS.keys()).index(
                st.session_state.get("proj_purpose", "General / Other")
                if st.session_state.get("proj_purpose", "General / Other")
                   in PURPOSE_PRESETS else "General / Other"
            ),
            help="Shapes the structure, depth, and tone of the presentation.",
        )
        st.session_state["proj_purpose"] = purpose

        industry = st.text_input(
            "Industry / context (optional)",
            value=st.session_state.get("proj_industry", ""),
            placeholder="e.g. SaaS, healthcare, real estate, fintech",
            key="industry_input",
        )
        st.session_state["proj_industry"] = industry

        audience = st.text_input(
            "Audience",
            value=st.session_state.get("proj_audience", "Executive team"),
            placeholder="e.g. Board of directors, Sales team, Engineering leads",
            key="audience_input",
        )
        st.session_state["proj_audience"] = audience

        question = st.text_area(
            "Goal / business question",
            value=st.session_state.get("proj_question", ""),
            placeholder=(
                "e.g. Make the case for expanding into the EU market by Q2, "
                "highlighting opportunity size, competitive gaps, and required investment…"
            ),
            height=120,
            key="question_input",
        )
        st.session_state["proj_question"] = question

        run_btn = st.button("⚡ Generate presentation", use_container_width=True, type="primary")

    with col_right:
        st.markdown('<div class="section-label">Output</div>', unsafe_allow_html=True)
        output_area   = st.empty()
        download_area = st.empty()

        if st.session_state.get("plan_mode_active") and not run_btn:
            with output_area.container():
                _render_plan_modal(st.session_state.get("template_upload"), ns="main")
        elif not run_btn and "pptx_bytes" not in st.session_state:
            output_area.markdown("""
<div class="status-card">
<div class="status-card-label">Ready</div>
<div class="status-card-value" style="color:#6080A8;font-size:.82rem;line-height:1.9">
Fill in the details on the left and click <strong style="color:#D0E4FF">Generate presentation</strong>.<br><br>
The pipeline will:<br>
&nbsp;1. Extract your uploaded documents<br>
&nbsp;2. Search the web for current data on your topic<br>
&nbsp;3. Generate a structured slide outline with Claude<br>
&nbsp;4. Let you review and edit the outline<br>
&nbsp;5. Export to a polished, editable PPTX
</div>
</div>
""", unsafe_allow_html=True)

        if st.session_state.get("pptx_bytes") and not run_btn:
            download_area.download_button(
                "⬇️ Download previous PPTX",
                data=st.session_state["pptx_bytes"],
                file_name=st.session_state.get("pptx_filename", "presentation.pptx"),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )

    # ── Run button handler ─────────────────────────────────────────────────────
    if run_btn:
        if not topic.strip():
            st.error("Please enter a topic.")
        else:
            # ── Sync all form fields into session state so _save_project captures them ──
            st.session_state["proj_topic"]    = topic
            st.session_state["proj_purpose"]  = purpose
            st.session_state["proj_industry"] = industry
            st.session_state["proj_audience"] = audience
            st.session_state["project_doc_names"] = [f.name for f in (uploaded_files or [])]

            # ── Auto-create a project if none is selected ──────────────────
            if not _active:
                _auto_name = re.sub(r"[^a-zA-Z0-9 ]+", "", topic).strip()[:40] or "Untitled"
                # Deduplicate if name exists
                _candidate = _auto_name
                _counter = 2
                while project_exists(OWNER, _candidate):
                    _candidate = f"{_auto_name} {_counter}"
                    _counter += 1
                create_project(OWNER, _candidate)
                st.session_state["active_project"] = _candidate
                _active = _candidate

            data_files_list = st.session_state.get("data_upload") or []
            if hasattr(data_files_list, "read"):
                data_files_list = [data_files_list]

            _template_bytes = _DEFAULT_TEMPLATE_BYTES  # SOA-HD Blue by default
            _tf = st.session_state.get("template_upload")
            if _tf:
                try:
                    _tf.seek(0); _template_bytes = _tf.read()
                    st.session_state["saved_template_bytes"] = _template_bytes
                except Exception:
                    pass

            _pipeline_steps = {
                "upload": bool(uploaded_files), "extract": False,
                "research": False, "analyze": False, "generate": False,
            }
            st.session_state["pipeline_steps"] = _pipeline_steps
            log_lines = []

            with output_area.container():
                st.markdown('<div class="section-label">Pipeline log</div>', unsafe_allow_html=True)
                log_area = st.empty()

                try:
                    for event in run_pipeline(
                        model=model,
                        uploaded_files=uploaded_files,
                        topic=topic,
                        purpose=purpose,
                        industry=industry,
                        audience=audience,
                        question=question,
                        web_search_en=web_search_enabled,
                        slide_count=slide_count,
                        theme=theme,
                        template_bytes=_template_bytes,
                        data_files=data_files_list,
                        plan_mode=True,
                    ):
                        etype = event[0]

                        if etype in ("log", "spinner"):
                            if etype == "spinner":
                                if log_lines and log_lines[-1][0] == "spinner":
                                    log_lines[-1] = ("log", log_lines[-1][1])
                                log_lines.append(("spinner", event[1]))
                            else:
                                if log_lines and log_lines[-1][0] == "spinner":
                                    log_lines[-1] = ("log", log_lines[-1][1])
                                log_lines.append(("log", event[1]))
                            log_area.markdown(_render_log(log_lines), unsafe_allow_html=True)

                        elif etype == "step_done":
                            _pipeline_steps[event[1]] = True
                            st.session_state["pipeline_steps"] = _pipeline_steps

                        elif etype == "sources":
                            st.session_state["research_sources"] = event[1]

                        elif etype == "slide_data":
                            st.session_state["slide_data"] = event[1]

                        elif etype == "plan_ready":
                            st.session_state["plan_slide_data"]  = event[1]
                            st.session_state["plan_mode_active"] = True
                            # Auto-save outline immediately — even before PPTX is built
                            if _active:
                                _save_project(OWNER, _active)
                                st.toast(f'Outline saved to "{_active}"', icon="📋")

                        elif etype == "pptx_bytes_out":
                            _slug = re.sub(r"[^a-zA-Z0-9]+", "_", topic)[:50]
                            fname = f"Presentation_{_slug}.pptx"
                            st.session_state["pptx_bytes"]    = event[1]
                            st.session_state["pptx_filename"] = fname
                            if _active:
                                _save_project(OWNER, _active)
                                st.toast(f'Auto-saved to "{_active}"', icon="💾")

                        elif etype == "error":
                            st.error(event[1], icon="🚨")
                            break

                except Exception as ex:
                    st.error(f"Unexpected error: {ex}")
                    import traceback; st.code(traceback.format_exc())

            if st.session_state.get("plan_mode_active"):
                with output_area.container():
                    _render_plan_modal(st.session_state.get("template_upload"), ns="main")

            if st.session_state.get("pptx_bytes"):
                with download_area.container():
                    st.success("Presentation ready!")
                    st.download_button(
                        "⬇️ Download PPTX",
                        data=st.session_state["pptx_bytes"],
                        file_name=st.session_state.get("pptx_filename","presentation.pptx"),
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                    )
                    _main_sources = st.session_state.get("research_sources", [])
                    if _main_sources:
                        with st.expander(f"🔗 {len(_main_sources)} verified source(s)", expanded=False):
                            for _s in _main_sources:
                                st.markdown(f"- [{_s['title']}]({_s['url']})")