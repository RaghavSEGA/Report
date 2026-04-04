"""
pdf_to_pptx.py — Convert PDF slides to fully editable PPTX.

Pipeline per page:
  1. Rasterise via PyMuPDF (no system deps)
  2. Detect if page is a full-page raster (NotebookLM/Canva/etc export)
  3. Send page image to Claude vision → JSON element spec
  4. Reconstruct as native pptx shapes/text/tables + image crops
"""
from __future__ import annotations
import base64, io, json, re
from typing import Callable
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

SLIDE_W_IN = 13.33
SLIDE_H_IN = 7.5
_ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
              "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}

EXTRACTION_PROMPT = """\
Analyze this presentation slide and extract ALL content as JSON.
Return ONLY valid JSON — no markdown fences, no explanation.
Slide is 13.33 inches wide x 7.5 inches tall. All x/y/w/h in inches from top-left.

{
  "background_color": "hex6",
  "elements": [
    {"type":"shape","x":0.0,"y":6.9,"w":13.33,"h":0.6,"fill":"1565C0","border_color":null,"border_pt":0},
    {"type":"text","x":0.35,"y":0.1,"w":10.5,"h":0.9,"text":"Title text — \\n for line breaks",
     "font_size":26.0,"bold":true,"italic":false,"color":"1A3A6B","bg_color":null,"align":"left","font":"Calibri"},
    {"type":"table","x":0.35,"y":1.6,"w":12.6,"h":4.1,
     "col_widths_in":[3.5,2.3,2.7,2.7],"row_heights_in":[0.72,1.09,1.09,1.09],
     "header_row":{"cells":["C1","C2","C3","C4"],"bg":"BDD7EE","fg":"1A3A6B","bold":true,"font_size":12.0,
       "cell_overrides":{"3":{"bg":"1565C0","fg":"FFFFFF"}}},
     "data_rows":[
       {"cells":["A","B","C","D"],"bg":"FFFFFF","fg":"1A1A2E","bold":false,"font_size":12.0,
        "cell_overrides":{"3":{"bg":"1565C0","fg":"FFFFFF","bold":true}}}]},
    {"type":"image","x":0.38,"y":2.35,"w":0.9,"h":0.75,"description":"SEGA Genesis controller logo"}
  ]
}

Rules:
- Emit back-to-front (backgrounds first, text/tables on top).
- shapes: colored bars, accent panels, bordered boxes, footer bars.
- text: every title, subtitle, body, bullet, label, footnote, page number.
  Preserve exact capitalisation and punctuation. Use \\n for line breaks.
- table: every cell. cell_overrides for cells differing from row defaults.
  col_widths_in must sum to ~w.
- image: one entry per logo, photo, icon, illustration visible on the slide.
  x/y/w/h = where it sits on the slide. DO NOT emit image for the overall
  slide background — only for distinct inline images (logos, photos, icons).
- Colors: 6-digit hex, no # prefix.
"""


def _rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _page_to_b64(pil_img) -> str:
    buf = io.BytesIO()
    pil_img.save(buf, format="PNG")
    return base64.b64encode(buf.getvalue()).decode()

def _is_full_page_raster(pdf_imgs: list, threshold: float = 0.65) -> bool:
    for im in pdf_imgs:
        frac = (im["w_in"] * im["h_in"]) / (SLIDE_W_IN * SLIDE_H_IN)
        if frac > threshold:
            return True
    return False

def _extract_pdf_images(fitz_doc, page_idx: int) -> list[dict]:
    page = fitz_doc[page_idx]
    pw, ph = page.rect.width, page.rect.height
    out = []
    for meta in page.get_images(full=True):
        xref = meta[0]
        try:
            base = fitz_doc.extract_image(xref)
            for r in page.get_image_rects(xref):
                out.append({
                    "bytes": base["image"], "ext": base["ext"],
                    "x_in": max(0.0, r.x0/pw*SLIDE_W_IN),
                    "y_in": max(0.0, r.y0/ph*SLIDE_H_IN),
                    "w_in": (r.x1-r.x0)/pw*SLIDE_W_IN,
                    "h_in": (r.y1-r.y0)/ph*SLIDE_H_IN,
                })
        except Exception:
            pass
    return out

def _crop_from_raster(pil_img, el: dict) -> bytes:
    iw, ih = pil_img.size
    sx, sy = iw/SLIDE_W_IN, ih/SLIDE_H_IN
    l = max(0, int(el["x"]*sx));       t = max(0, int(el["y"]*sy))
    r = min(iw, int((el["x"]+el["w"])*sx)); b = min(ih, int((el["y"]+el["h"])*sy))
    if r <= l or b <= t:
        raise ValueError("degenerate crop")
    buf = io.BytesIO()
    pil_img.crop((l, t, r, b)).save(buf, format="PNG")
    return buf.getvalue()

def _place_image_bytes(slide, img_bytes: bytes, x, y, w, h) -> None:
    from PIL import Image as _PIL
    try:
        pil = _PIL.open(io.BytesIO(img_bytes))
        buf = io.BytesIO(); pil.save(buf, format="PNG"); buf.seek(0)
    except Exception:
        buf = io.BytesIO(img_bytes); buf.seek(0)
    slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))

def _extract_slide_json(image_b64: str, api_key: str, model: str) -> dict:
    import anthropic
    client = anthropic.Anthropic(api_key=api_key)
    msg = client.messages.create(
        model=model, max_tokens=4096,
        messages=[{"role":"user","content":[
            {"type":"image","source":{"type":"base64","media_type":"image/png","data":image_b64}},
            {"type":"text","text":EXTRACTION_PROMPT},
        ]}],
    )
    raw = msg.content[0].text.strip()
    raw = re.sub(r"^```[a-z]*\n?", "", raw)
    raw = re.sub(r"\n?```$", "", raw)
    return json.loads(raw.strip())

# ── Renderers ─────────────────────────────────────────────────────────────────

def _render_shape(slide, el):
    s = slide.shapes.add_shape(1, Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"]))
    fill = el.get("fill")
    if fill: s.fill.solid(); s.fill.fore_color.rgb = _rgb(fill)
    else: s.fill.background()
    border = el.get("border_color")
    if border: s.line.color.rgb = _rgb(border); s.line.width = Pt(float(el.get("border_pt",1)))
    else: s.line.fill.background()

def _render_text(slide, el):
    txb = slide.shapes.add_textbox(Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"]))
    tf = txb.text_frame; tf.word_wrap = True
    bg = el.get("bg_color")
    if bg: txb.fill.solid(); txb.fill.fore_color.rgb = _rgb(bg)
    else: txb.fill.background()
    for li, line in enumerate(str(el.get("text","")).split("\n")):
        p = tf.paragraphs[0] if li==0 else tf.add_paragraph()
        p.alignment = _ALIGN_MAP.get(el.get("align","left"), PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = line
        run.font.name   = el.get("font","Calibri")
        run.font.bold   = bool(el.get("bold",False))
        run.font.italic = bool(el.get("italic",False))
        run.font.size   = Pt(float(el.get("font_size",12)))
        run.font.color.rgb = _rgb(el.get("color","000000"))

def _set_cell_fill(cell, h):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    for x in tcPr.findall(qn("a:solidFill")): tcPr.remove(x)
    sf = etree.SubElement(tcPr, qn("a:solidFill"))
    sc = etree.SubElement(sf,   qn("a:srgbClr")); sc.set("val", h.lstrip("#"))

def _render_table(slide, el):
    header = el.get("header_row", {}); data_rows = el.get("data_rows", [])
    hcells = header.get("cells", [])
    n_cols = len(hcells) or (len(data_rows[0]["cells"]) if data_rows else 1)
    n_rows = (1 if header else 0) + len(data_rows)
    tbl = slide.shapes.add_table(n_rows, n_cols,
        Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"])).table
    for i, cw in enumerate(el.get("col_widths_in",[])[:n_cols]):
        tbl.columns[i].width = Inches(float(cw))
    for i, rh in enumerate(el.get("row_heights_in",[])[:n_rows]):
        tbl.rows[i].height = Inches(float(rh))

    def fc(r, c, text, bg=None, fg="000000", bold=False, fs=12.0, align=PP_ALIGN.CENTER):
        cell = tbl.cell(r,c); cell.text = ""
        tf = cell.text_frame; tf.word_wrap = True
        for li, line in enumerate(str(text).split("\n")):
            p = tf.paragraphs[0] if li==0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run(); run.text = line
            run.font.name="Calibri"; run.font.bold=bold
            run.font.size=Pt(float(fs)); run.font.color.rgb=_rgb(fg)
        if bg: _set_cell_fill(cell, bg)

    row_off = 0
    if header:
        hbg=header.get("bg","BDD7EE"); hfg=header.get("fg","000000")
        hbold=header.get("bold",True); hsize=float(header.get("font_size",12))
        hov=header.get("cell_overrides",{})
        for c, txt in enumerate(hcells[:n_cols]):
            ov=hov.get(str(c),{})
            fc(0,c,txt,bg=ov.get("bg",hbg),fg=ov.get("fg",hfg),
               bold=ov.get("bold",hbold),fs=ov.get("font_size",hsize))
        row_off=1

    for ri, rd in enumerate(data_rows):
        rbg=rd.get("bg"); rfg=rd.get("fg","000000")
        rbold=rd.get("bold",False); rsize=float(rd.get("font_size",12))
        rov=rd.get("cell_overrides",{})
        for c, txt in enumerate(rd.get("cells",[])[:n_cols]):
            ov=rov.get(str(c),{})
            fc(row_off+ri,c,txt,bg=ov.get("bg",rbg),fg=ov.get("fg",rfg),
               bold=ov.get("bold",rbold),fs=ov.get("font_size",rsize),
               align=_ALIGN_MAP.get(ov.get("align","center"),PP_ALIGN.CENTER))

def _render_element(slide, el: dict, page_raster=None) -> None:
    t = el.get("type","shape")
    if t == "shape":  _render_shape(slide, el)
    elif t == "text": _render_text(slide, el)
    elif t == "table":_render_table(slide, el)
    elif t == "image" and page_raster is not None:
        try:
            img_bytes = _crop_from_raster(page_raster, el)
            _place_image_bytes(slide, img_bytes, el["x"], el["y"], el["w"], el["h"])
        except Exception:
            pass

# ── Main ──────────────────────────────────────────────────────────────────────

def pdf_to_editable_pptx(
    pdf_bytes: bytes,
    api_key: str,
    model: str = "claude-opus-4-5",
    dpi: int = 150,
    progress_cb: Callable | None = None,
) -> tuple[bytes, list[str]]:
    """
    Convert PDF slides to editable PPTX. Returns (pptx_bytes, errors).

    For full-page raster PDFs (NotebookLM, Canva exports):
      - Claude reconstructs all shapes/text/tables as native elements
      - Inline images (logos, icons) are cropped from the raster by coordinate
      - The full-page background image is NOT placed (everything is editable)

    For vector PDFs with discrete embedded images:
      - Embedded images are placed directly from the PDF XObjects
      - Text/shapes/tables reconstructed as native elements
    """
    import fitz
    from PIL import Image as _PIL

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    mat = fitz.Matrix(dpi/72, dpi/72)
    page_rasters = []
    for pg in doc:
        pix = pg.get_pixmap(matrix=mat)
        page_rasters.append(_PIL.frombytes("RGB", [pix.width, pix.height], pix.samples))

    n = len(page_rasters)
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    errors: list[str] = []

    for i, raster in enumerate(page_rasters):
        if progress_cb:
            progress_cb(i / n)

        pdf_imgs      = _extract_pdf_images(doc, i)
        has_bg_raster = _is_full_page_raster(pdf_imgs)
        b64           = _page_to_b64(raster)

        try:
            spec = _extract_slide_json(b64, api_key, model)
        except Exception as e:
            errors.append(f"Page {i+1}: extraction failed — {e}")
            slide = prs.slides.add_slide(blank)
            txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            txb.text_frame.paragraphs[0].add_run().text = f"[Page {i+1}: {e}]"
            continue

        slide = prs.slides.add_slide(blank)

        # White/colored base
        bg = spec.get("background_color", "FFFFFF")
        bgs = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bgs.fill.solid(); bgs.fill.fore_color.rgb = _rgb(bg); bgs.line.fill.background()

        # For vector PDFs: place discrete embedded images before overlays
        if not has_bg_raster:
            for pi in pdf_imgs:
                try:
                    _place_image_bytes(slide, pi["bytes"],
                                       pi["x_in"], pi["y_in"], pi["w_in"], pi["h_in"])
                except Exception as e:
                    errors.append(f"Page {i+1} image: {e}")

        # Render all elements — pass raster so "image" crops work
        for el in spec.get("elements", []):
            try:
                _render_element(slide, el, page_raster=raster)
            except Exception as e:
                errors.append(f"Page {i+1} {el.get('type','?')}: {e}")

    if progress_cb:
        progress_cb(1.0)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), errors