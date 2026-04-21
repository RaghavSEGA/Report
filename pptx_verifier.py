"""
pptx_verifier.py — Action-First PPTX Verification Tool
Run with:  streamlit run pptx_verifier.py
"""

import streamlit as st
import anthropic
import hashlib
import hmac
import base64
import time
import random
import string
import sqlite3
import os
import json
import tempfile
from datetime import datetime

# ── Try to import pptx reader ─────────────────────────────────
try:
    from markitdown import MarkItDown
    HAS_MARKITDOWN = True
except ImportError:
    HAS_MARKITDOWN = False

try:
    from pptx import Presentation
    from pptx.util import Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ═══════════════════════════════════════════════════════════════
# CONFIG
# ═══════════════════════════════════════════════════════════════
st.set_page_config(
    page_title="Slide Verifier",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="collapsed",
)

ALLOWED_DOMAIN    = st.secrets.get("ALLOWED_DOMAIN", "")
TOKEN_EXPIRY_DAYS = int(st.secrets.get("TOKEN_EXPIRY_DAYS", 7))
DB_PATH           = st.secrets.get("DB_PATH", "verifier_users.db")
OTP_EXPIRY_SECS   = 600
MAX_OTP_ATTEMPTS  = 5

# ═══════════════════════════════════════════════════════════════
# REFERENCE GALLERY DATA
# Sourced from BCG Melbourne Cultural Destination &
# McKinsey/IACPM GenAI in Credit decks
# ═══════════════════════════════════════════════════════════════

REFERENCE_SLIDES = [
    {
        "deck": "BCG — Melbourne as a Global Cultural Destination",
        "firm": "BCG",
        "slide_num": 4,
        "title": "Melbourne has a compelling creative and cultural offer; the city attracted >10m visitors in 2015",
        "what_works": [
            "Full-sentence action title leads with the conclusion — reader knows the 'so what' before reading a word of body content.",
            "Three distinct claim blocks (strength / economic case / challenge) each with its own bold lead-in, creating scannable Z-pattern reading flow.",
            "Data is specific and cited inline (>10m, 8% GVA, 40% by 2025) — evidence serves the headline rather than replacing it.",
            "Deliberate contrast: strengths stated first, competitive threat introduced last — classic 'complication' structure.",
        ],
        "annotation": "PASS — Strong Action-First title. Proof is organized to support one claim per paragraph. No orphan bullets.",
        "badge": "pass",
        "pattern": "Stacked claim blocks with bold lead-ins",
    },
    {
        "deck": "BCG — Melbourne as a Global Cultural Destination",
        "firm": "BCG",
        "slide_num": 5,
        "title": "Five strategic priorities to improve Melbourne's position as a global cultural and creative destination",
        "what_works": [
            "Numbered list turns abstract strategy into a concrete, countable deliverable.",
            "Italic emphasis on the priority name separates label from description without needing extra headers.",
            "The slide title is a full verb phrase ('to improve Melbourne's position') — an imperative, not a category label.",
            "Scope exclusions were stated upfront on slide 2, so this summary carries no ambiguity.",
        ],
        "annotation": "PASS — Title states a direction, not a topic. Numbered structure gives the audience an instant mental model.",
        "badge": "pass",
        "pattern": "Numbered priority list with inline italic labels",
    },
    {
        "deck": "BCG — Melbourne as a Global Cultural Destination",
        "firm": "BCG",
        "slide_num": 9,
        "title": "Creative industries increasingly important to the economy",
        "what_works": [
            "Sub-title ('Workers generate slightly more GVA than the average Victorian worker') acts as an Action Subtitle — the charts exist to prove it.",
            "Two side-by-side charts answer one compound question: is the creative economy growing in employment AND GVA?",
            "CAGR callout (5.0% vs 1.8% average) is the punchline, placed at far right so the eye lands there last.",
            "Footnotes define abbreviations (GVA, CAGR) inline — no ambiguity left for the reader.",
        ],
        "annotation": "PASS — Action Subtitle functions as a second claim. Dual-chart parallel layout elegantly handles a two-part argument.",
        "badge": "pass",
        "pattern": "Main claim title + Action Subtitle + parallel evidence charts",
    },
    {
        "deck": "BCG — Melbourne as a Global Cultural Destination",
        "firm": "BCG",
        "slide_num": 22,
        "title": "Melbourne has strengths in offer and governance, underperforms in brand and marketing",
        "what_works": [
            "Title is a direct contrast statement — two competing facts. Audience has the full verdict before reading the matrix.",
            "Color coding (green/amber/red) maps directly to the three-word legend — no guessing required.",
            "Matrix structure (6 dimensions × 3 sub-factors) creates visual density without clutter.",
            "Framework introduced earlier (slide 20) means this slide goes straight to assessment without re-explaining.",
        ],
        "annotation": "PASS — Title is a judgment, not a description. Color-coded assessment matrix is rare but highly effective.",
        "badge": "pass",
        "pattern": "Contrast-statement title + color-coded assessment matrix",
    },
    {
        "deck": "BCG — Melbourne as a Global Cultural Destination",
        "firm": "BCG",
        "slide_num": 35,
        "title": "Travellers more likely to recommend Melbourne after visiting — suggests marketing could boost visitation",
        "what_works": [
            "Long title earns its length: contains both the finding ('more likely after visiting') AND the implication ('marketing could boost visitation').",
            "Single bar chart with a +24 callout does one job cleanly — chart title repeats the claim in plain language.",
            "Pull-quote boxes serve as qualitative evidence — they complement, not replace, the quantitative bar chart.",
            "Shaded background on pull-quotes makes them visually distinct — classic Takeaway Block execution.",
        ],
        "annotation": "PASS — Finding + implication in title. Pull-quote boxes = Takeaway Blocks well executed.",
        "badge": "pass",
        "pattern": "Finding + implication title; pull-quote boxes as takeaway blocks",
    },
    {
        "deck": "McKinsey/IACPM — Emerging GenAI Use Cases in Credit",
        "firm": "McKinsey",
        "slide_num": 3,
        "title": "Executive Summary",
        "what_works": [
            "Each bullet starts with a quantified fact — no vague language ('Most institutions…' is immediately followed by a number).",
            "Functions as a standalone document — a reader who only sees this slide has the full story.",
            "Each paragraph addresses one distinct theme — no paragraph mixes topics.",
            "Bold emphasis within bullets highlights the key data point in each sentence, enabling rapid scanning.",
        ],
        "annotation": "PASS — Executive Summary done right: every bullet is a claim with evidence, not a topic label.",
        "badge": "pass",
        "pattern": "Claim-first bullets with inline data; bold for scannability",
    },
    {
        "deck": "McKinsey/IACPM — Emerging GenAI Use Cases in Credit",
        "firm": "McKinsey",
        "slide_num": 6,
        "title": "Institutions are prioritizing use cases like supporting underwriting (synthesizing, drafting memo) and portfolio monitoring (early warning)",
        "what_works": [
            "Title is an active observation with specific examples in parentheses — reader knows the conclusion AND the top two use cases.",
            "Stacked bar by institution size lets each audience member find their own segment without re-reading.",
            "Development stage column (Ideation → PoC → Pilot → Deployment) shows not just *what* but *how far along*.",
            "Percentage totals alongside absolute counts support both reading modes simultaneously.",
        ],
        "annotation": "PASS — Title names the conclusion and examples. Chart adds two dimensions (what + maturity) elegantly.",
        "badge": "pass",
        "pattern": "Active observation with examples; multi-dimension bar chart",
    },
    {
        "deck": "McKinsey/IACPM — Emerging GenAI Use Cases in Credit",
        "firm": "McKinsey",
        "slide_num": 12,
        "title": "Institutions prioritize productivity improvement as the most important factor when initiating or developing GenAI use cases",
        "what_works": [
            "Title is a full declarative sentence — the ranked #1 finding is in the headline, not buried in bullets.",
            "Four summary callout boxes at the top give the ranked answer before the supporting data appears.",
            "Callout boxes use consistent structure: rank statement + percentage + label — identical format reduces reading effort.",
            "Detailed ranking data below provides nuance — top boxes give the verdict, chart gives the breakdown.",
        ],
        "annotation": "PASS — Summary callout boxes above the chart function perfectly as pre-chart Takeaway Blocks.",
        "badge": "pass",
        "pattern": "Ranked finding in title; pre-chart callout boxes as takeaways",
    },
    {
        "deck": "McKinsey/IACPM — Emerging GenAI Use Cases in Credit",
        "firm": "McKinsey",
        "slide_num": 13,
        "title": "Leadership at majority of the institutions are positioning GenAI as a priority",
        "what_works": [
            "'Majority' is precise without being falsely exact — better than 'many' (too vague) or '52%' (too precise for a headline).",
            "Three-tier classification maps to three distinct descriptions — the reader can self-classify their institution.",
            "Descriptions use consistent voice ('Senior leadership promotes…' / 'The organization is encouraged…') making comparison easy.",
            "Bar chart breakdown by institution type appears to the right — answer first, segmentation second.",
        ],
        "annotation": "PASS — Classification with self-assessment language is a McKinsey signature pattern. Verdict in title.",
        "badge": "pass",
        "pattern": "Three-tier classification with self-assessment descriptions",
    },
    {
        "deck": "McKinsey/IACPM — Emerging GenAI Use Cases in Credit",
        "firm": "McKinsey",
        "slide_num": 16,
        "title": "Insufficient performance and complexity challenge are the top reasons noted for slowing down GenAI use cases",
        "what_works": [
            "Title names both top findings specifically — 'insufficient performance' AND 'complexity challenge' are in the headline.",
            "Three-column structure (Top / Relevant / Not Relevant) forces every reason to be ranked, not just listed.",
            "Percentages shown only for 'top reasons' — prevents over-indexing on less important columns.",
            "Reasons list uses parallel grammar ('Insufficient X', 'Marginal X') — each item consistently framed as a problem noun.",
        ],
        "annotation": "PASS — Ranking taxonomy turns a list into a hierarchy. Title names the conclusion, not the topic.",
        "badge": "pass",
        "pattern": "Ranked taxonomy (Top/Relevant/Not) with percentages for top column only",
    },
]

# Text version injected into the AI system prompt
REFERENCE_PATTERNS_TEXT = """
## Patterns from Best-in-Class Consulting Decks (BCG & McKinsey)

### TITLE PATTERNS
- **Contrast title**: "Melbourne has strengths in X, underperforms in Y" — verdict before body.
- **Finding + implication**: "Travellers more likely to recommend after visiting — suggests marketing could boost visitation" — data point THEN business implication.
- **Active observation with examples**: "Institutions are prioritizing use cases like underwriting and portfolio monitoring" — conclusion + specific examples in parentheses.
- **Ranked finding**: "Institutions prioritize productivity improvement as the most important factor" — #1 result in the headline.
- **Majority qualifier**: Use 'majority' instead of 'many' (vague) or '52%' (too precise for headlines).

### TAKEAWAY BLOCK PATTERNS
- **Pre-chart callout boxes**: Summary boxes above the chart. Each box: rank statement + percentage + label. Verdict before data.
- **Pull-quote box**: Stakeholder quote in shaded box at bottom-right. Shading = Takeaway Block signal.
- **Action Subtitle**: Second line below main title that states a sub-claim. Body charts prove the sub-claim.
- **Three-tier classification**: Classify findings into 3 tiers with self-assessment descriptions.

### BENCHMARK TABLE
| Pattern | Weak title | Strong title |
|---|---|---|
| Contrast | "SWOT Analysis" | "Melbourne has strengths in offer and governance, underperforms in brand and marketing" |
| Ranked finding | "Factor Prioritization" | "Institutions prioritize productivity improvement as the most important factor" |
| Finding + implication | "Visitor Data" | "Travellers more likely to recommend Melbourne after visiting — suggests marketing could boost visitation" |
| Active observation | "GenAI Use Case Update" | "Institutions are prioritizing use cases like underwriting (synthesizing, drafting memo) and portfolio monitoring" |
"""


# ═══════════════════════════════════════════════════════════════
# DATABASE
# ═══════════════════════════════════════════════════════════════

def _get_db():
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.row_factory = sqlite3.Row
    conn.execute("""CREATE TABLE IF NOT EXISTS users (
        email TEXT PRIMARY KEY, display_name TEXT DEFAULT '',
        created_at REAL DEFAULT (unixepoch()), last_login REAL)""")
    conn.execute("""CREATE TABLE IF NOT EXISTS analyses (
        id INTEGER PRIMARY KEY AUTOINCREMENT, email TEXT NOT NULL,
        filename TEXT NOT NULL, created_at REAL DEFAULT (unixepoch()), result_json TEXT)""")
    conn.commit()
    return conn

def _upsert_user(email):
    with _get_db() as conn:
        conn.execute("INSERT INTO users(email) VALUES(?) ON CONFLICT(email) DO UPDATE SET last_login=unixepoch()", (email,))

def _save_analysis(email, filename, result):
    with _get_db() as conn:
        conn.execute("INSERT INTO analyses(email, filename, result_json) VALUES(?,?,?)",
                     (email, filename, json.dumps(result)))

def _get_history(email, limit=10):
    with _get_db() as conn:
        return conn.execute(
            "SELECT id, filename, created_at, result_json FROM analyses WHERE email=? ORDER BY created_at DESC LIMIT ?",
            (email, limit)).fetchall()


# ═══════════════════════════════════════════════════════════════
# AUTH
# ═══════════════════════════════════════════════════════════════

def _make_token(email):
    secret  = st.secrets.get("COOKIE_SIGNING_KEY", "change-me-please")
    expiry  = int(time.time()) + TOKEN_EXPIRY_DAYS * 86400
    payload = f"{email}|{expiry}"
    sig     = hmac.new(secret.encode(), payload.encode(), hashlib.sha256).hexdigest()
    return base64.urlsafe_b64encode(f"{payload}|{sig}".encode()).decode()

def _verify_token(token):
    try:
        secret  = st.secrets.get("COOKIE_SIGNING_KEY", "change-me-please")
        decoded = base64.urlsafe_b64decode(token.encode()).decode()
        email, expiry_str, sig = decoded.rsplit("|", 2)
        payload  = f"{email}|{expiry_str}"
        expected = hmac.new(secret.encode(), payload.encode(), hashlib.sha256).hexdigest()
        if not hmac.compare_digest(sig, expected): return None
        if int(time.time()) > int(expiry_str): return None
        return email
    except Exception:
        return None

def _send_otp(email, code):
    import boto3
    client = boto3.client("ses",
        region_name=st.secrets.get("AWS_SES_REGION", "us-east-1"),
        aws_access_key_id=st.secrets.get("AWS_ACCESS_KEY_ID"),
        aws_secret_access_key=st.secrets.get("AWS_SECRET_ACCESS_KEY"))
    client.send_email(
        Source=st.secrets.get("EMAIL_FROM", "noreply@example.com"),
        Destination={"ToAddresses": [email]},
        Message={"Subject": {"Data": "Your Slide Verifier sign-in code"},
                 "Body": {"Html": {"Data": f"""<div style="font-family:sans-serif;max-width:480px">
                   <h2>📊 Slide Verifier</h2><p>Your one-time sign-in code:</p>
                   <div style="font-size:36px;font-weight:700;letter-spacing:8px;color:#4f46e5;
                   padding:16px;background:#f1f5f9;border-radius:8px;text-align:center">{code}</div>
                   <p style="color:#64748b;font-size:13px">Expires in 10 minutes.</p></div>"""}}})


# ═══════════════════════════════════════════════════════════════
# PPTX EXTRACTION
# ═══════════════════════════════════════════════════════════════

def extract_slide_content(uploaded_file):
    slides = []
    if HAS_PPTX:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp.write(uploaded_file.getvalue()); tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            for i, slide in enumerate(prs.slides, 1):
                sd = {"slide_num": i, "title": "", "body_texts": [], "font_sizes": [], "notes": ""}
                if slide.has_notes_slide:
                    sd["notes"] = slide.notes_slide.notes_text_frame.text.strip()
                for shape in slide.shapes:
                    if not shape.has_text_frame: continue
                    tf = shape.text_frame
                    full_text = tf.text.strip()
                    if not full_text: continue
                    is_title = shape.shape_type == 13 or (
                        hasattr(shape, "placeholder_format") and
                        shape.placeholder_format is not None and
                        shape.placeholder_format.idx in (0, 1))
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if run.font.size: sd["font_sizes"].append(run.font.size.pt)
                    if is_title: sd["title"] = full_text
                    else: sd["body_texts"].append(full_text)
                slides.append(sd)
        finally:
            os.unlink(tmp_path)
    elif HAS_MARKITDOWN:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp.write(uploaded_file.getvalue()); tmp_path = tmp.name
        try:
            md = MarkItDown(); result = md.convert(tmp_path)
            for chunk in result.text_content.split("<!-- Slide ")[1:]:
                num = int(chunk[:chunk.index(" ")])
                text = chunk[chunk.index("-->")+3:].strip()
                lines = [l.strip() for l in text.split("\n") if l.strip()]
                slides.append({"slide_num": num, "title": lines[0] if lines else "",
                                "body_texts": lines[1:], "font_sizes": [], "notes": ""})
        finally:
            os.unlink(tmp_path)
    return slides


# ═══════════════════════════════════════════════════════════════
# AI ANALYSIS
# ═══════════════════════════════════════════════════════════════

ACTION_FIRST_SYSTEM = """You are an expert presentation coach specializing in the "Action-First" communication framework used by top-tier consulting firms. Your role is to review PowerPoint presentations slide by slide and provide specific, actionable feedback — benchmarked against real BCG and McKinsey slides.

## The Action-First Framework

### Core Logic: The Action-Takeaway Sandwich
1. **Action-Title Mandate** — Title MUST be a complete sentence summarizing the conclusion.
   - REJECT: "Q4 Revenue Analysis" | ENFORCE: "Q4 Revenue exceeded targets by 15% due to optimized checkout flows"
2. **Takeaway Block** — Slide must END with a "Key Takeaway" or "So What?" block connecting data to business objective.
3. **The ArtCenter "Read"** — Action-Title must be heaviest weight (Bold) and largest scale.

### The Action-First Mastery Checklist
**I. The Header**
- Action-Title Test: Is the title a full sentence with a verb?
- The "Flip" Check: If you delete everything except the title, does the audience still know the key message?
- No Category Labels: Does the title avoid dead words like "Overview," "Update," "Status," or "Summary"?

**II. The Footer**
- Takeaway Lead-in: Does the slide end with a bolded "Key Takeaway" or "Executive Insight"?
- WIIFY Validation: Does the takeaway name the benefit to the stakeholder?
- Supportive Linkage: Does the takeaway reference the data in the body?

**III. The Visual Hierarchy**
- Title Dominance: Is the Action-Title at least 24pt–36pt Bold?
- Weight Contrast: Is there a clear visual drop from title to body?
- The Z-Pattern: Does the eye flow Action-Title → Evidence → Takeaway?

""" + REFERENCE_PATTERNS_TEXT + """

## Your Output Format

For each slide:
1. **Slide [N]: [Current Title]** — PASS / NEEDS WORK / FAIL
2. Brief summary of what works (reference the consulting pattern it resembles, if applicable)
3. **Issues Found** — numbered list with checklist item cited
4. **Revised Action-Title** — using one of the named consulting patterns (contrast / finding+implication / ranked finding / active observation with examples)
5. **Suggested Takeaway Block** — using one of the named patterns (pre-chart callout / pull-quote box / action subtitle / three-tier classification)

End with **Overall Deck Score** (0–100) and top 3 highest-impact changes."""


def analyze_slides_with_claude(slides, filename):
    client = anthropic.Anthropic(api_key=st.secrets["ANTHROPIC_API_KEY"])
    parts = []
    for s in slides:
        p = [f"## Slide {s['slide_num']}", f"**Title:** {s['title'] or '(no title)'}"]
        if s["body_texts"]:
            p.append("**Body:**")
            p.extend(f"- {bt[:500]}" for bt in s["body_texts"])
        if s["font_sizes"]:
            p.append(f"**Fonts:** min={min(s['font_sizes']):.0f}pt, max={max(s['font_sizes']):.0f}pt")
        if s["notes"]:
            p.append(f"**Notes:** {s['notes'][:200]}")
        parts.append("\n".join(p))

    prompt = f"""Review this presentation using the Action-First framework, benchmarking against BCG and McKinsey standards.

**Filename:** {filename} | **Total Slides:** {len(slides)}

---

{chr(10).join(parts)}

---

Apply the full checklist. Quote actual slide text when citing issues. Use the named consulting patterns when rewriting titles and takeaways."""

    with st.spinner("Analyzing your presentation with Claude…"):
        msg = client.messages.create(
            model="claude-opus-4-5", max_tokens=4096,
            system=ACTION_FIRST_SYSTEM,
            messages=[{"role": "user", "content": prompt}])
    return msg.content[0].text


# ═══════════════════════════════════════════════════════════════
# STYLES
# ═══════════════════════════════════════════════════════════════

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');
html,body,[class*="css"]{font-family:'Inter',sans-serif}
.stApp{background:#0f172a;color:#e2e8f0}
.auth-card{max-width:420px;margin:4rem auto 2rem;background:#1e293b;border:1px solid #334155;border-radius:16px;padding:2.5rem 2rem}
.auth-logo{font-size:2.5rem;text-align:center;margin-bottom:.5rem}
.auth-title{font-size:1.4rem;font-weight:700;color:#f1f5f9;text-align:center;margin-bottom:.25rem}
.auth-sub{font-size:.85rem;color:#94a3b8;text-align:center;margin-bottom:2rem}
.ref-card{background:#1e293b;border:1px solid #334155;border-radius:12px;padding:1.25rem 1.5rem;margin-bottom:.75rem;height:100%}
.ref-deck{font-size:.68rem;color:#64748b;text-transform:uppercase;letter-spacing:.05em;margin-bottom:.3rem}
.ref-title{font-size:.9rem;font-weight:700;color:#f1f5f9;margin-bottom:.6rem;line-height:1.4}
.ref-point{display:flex;gap:.5rem;margin-bottom:.35rem;font-size:.8rem;color:#cbd5e1;line-height:1.45}
.ref-icon{flex-shrink:0;color:#22c55e}
.ref-badge{display:inline-block;background:#064e3b;color:#6ee7b7;padding:2px 10px;border-radius:99px;font-size:.68rem;font-weight:700;margin-bottom:.6rem}
.ref-annotation{font-size:.78rem;color:#94a3b8;border-left:3px solid #4f46e5;padding-left:.6rem;margin-top:.6rem}
.ref-pattern{font-size:.73rem;color:#818cf8;font-style:italic;margin-top:.5rem}
.chip-bcg{background:#064e3b;color:#6ee7b7;padding:1px 7px;border-radius:5px;font-size:.68rem;font-weight:700;margin-right:.4rem}
.chip-mckinsey{background:#1e3a5f;color:#93c5fd;padding:1px 7px;border-radius:5px;font-size:.68rem;font-weight:700;margin-right:.4rem}
.stButton>button{background:#4f46e5!important;color:white!important;border:none!important;border-radius:8px!important;font-weight:600!important}
.stButton>button:hover{background:#4338ca!important}
.stTextInput>div>div>input{background:#0f172a!important;border:1px solid #334155!important;color:#e2e8f0!important;border-radius:8px!important}
.stMarkdown h2,.stMarkdown h3{color:#818cf8}
.stExpander{background:#1e293b!important;border:1px solid #334155!important;border-radius:10px!important}
</style>
""", unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════
# SESSION STATE
# ═══════════════════════════════════════════════════════════════

_url_token   = st.query_params.get("t", "")
_token_email = _verify_token(_url_token) if _url_token else None

for k, v in {
    "auth_verified": False, "auth_email": "", "auth_token": "",
    "otp_code": "", "otp_email": "", "otp_expiry": 0,
    "otp_sent": False, "otp_attempts": 0,
    "analysis_result": None, "analysis_filename": None,
    "show_history": False, "gallery_filter": "All",
}.items():
    if k not in st.session_state:
        st.session_state[k] = v

if _token_email and not st.session_state.auth_verified:
    st.session_state.auth_verified = True
    st.session_state.auth_email    = _token_email
    _upsert_user(_token_email)


# ═══════════════════════════════════════════════════════════════
# LOGIN GATE
# ═══════════════════════════════════════════════════════════════

if not st.session_state.auth_verified:
    st.markdown("<style>section[data-testid='stSidebar']{display:none!important}</style>",
                unsafe_allow_html=True)
    _, col, _ = st.columns([1, 2, 1])
    with col:
        st.markdown('<div class="auth-card"><div class="auth-logo">📊</div>'
                    '<div class="auth-title">Slide Verifier</div>'
                    '<div class="auth-sub">Action-First presentation checker</div></div>',
                    unsafe_allow_html=True)
        if not st.session_state.otp_sent:
            email_in = st.text_input("Email address", placeholder="you@yourcompany.com",
                                     label_visibility="collapsed", key="auth_email_input")
            if st.button("Send verification code", use_container_width=True):
                addr = email_in.strip().lower()
                if not addr or "@" not in addr:
                    st.error("Please enter a valid email address.")
                elif ALLOWED_DOMAIN and not addr.endswith(ALLOWED_DOMAIN):
                    st.error(f"Access restricted to {ALLOWED_DOMAIN} addresses.")
                else:
                    code = "".join(random.choices(string.digits, k=6))
                    try:
                        _send_otp(addr, code)
                        st.session_state.update(otp_code=code, otp_email=addr,
                            otp_expiry=int(time.time())+OTP_EXPIRY_SECS,
                            otp_sent=True, otp_attempts=0)
                        st.rerun()
                    except Exception as e:
                        st.error(f"Failed to send email: {e}")
        else:
            st.info(f"Code sent to **{st.session_state.otp_email}**")
            otp_in = st.text_input("6-digit code", placeholder="123456",
                                   max_chars=6, label_visibility="collapsed")
            c1, c2 = st.columns(2)
            if c1.button("Verify", use_container_width=True):
                st.session_state.otp_attempts += 1
                if st.session_state.otp_attempts > MAX_OTP_ATTEMPTS:
                    st.error("Too many attempts."); st.session_state.otp_sent = False
                elif int(time.time()) > st.session_state.otp_expiry:
                    st.error("Code expired."); st.session_state.otp_sent = False
                elif otp_in.strip() == st.session_state.otp_code:
                    token = _make_token(st.session_state.otp_email)
                    st.session_state.update(auth_verified=True,
                        auth_email=st.session_state.otp_email, auth_token=token)
                    _upsert_user(st.session_state.otp_email)
                    st.query_params["t"] = token; st.rerun()
                else:
                    st.error(f"Incorrect code. {MAX_OTP_ATTEMPTS - st.session_state.otp_attempts} attempts remaining.")
            if c2.button("Back", use_container_width=True):
                st.session_state.otp_sent = False; st.rerun()
    st.stop()


# ═══════════════════════════════════════════════════════════════
# MAIN APP
# ═══════════════════════════════════════════════════════════════

# Header
c1, c2, c3, c4 = st.columns([.07, .6, .15, .18])
c1.markdown("## 📊")
c2.markdown(f"**Slide Verifier** — Action-First Checker")
c2.caption(f"Signed in as {st.session_state.auth_email}")
if c3.button("📋 History", use_container_width=True):
    st.session_state.show_history = not st.session_state.show_history
if c4.button("Sign out", use_container_width=True):
    for k in ["auth_verified", "auth_email", "auth_token", "otp_sent", "analysis_result", "analysis_filename"]:
        st.session_state[k] = False if isinstance(st.session_state[k], bool) else ""
    st.query_params.clear(); st.rerun()

st.divider()

# History panel
if st.session_state.show_history:
    with st.expander("📋 Past Analyses", expanded=True):
        rows = _get_history(st.session_state.auth_email)
        if not rows:
            st.caption("No analyses yet.")
        else:
            for row in rows:
                dt = datetime.fromtimestamp(row["created_at"]).strftime("%b %d %H:%M")
                result = json.loads(row["result_json"]) if row["result_json"] else {}
                hc1, hc2, hc3 = st.columns([3, 1, 1])
                hc1.markdown(f"**{row['filename']}**"); hc2.caption(dt)
                if hc3.button("View", key=f"h_{row['id']}"):
                    st.session_state.analysis_result   = result.get("text", "")
                    st.session_state.analysis_filename = row["filename"]
                    st.session_state.show_history      = False; st.rerun()

# Main tabs
t_verify, t_gallery, t_patterns, t_create = st.tabs([
    "🔍 Verify a Deck",
    "📚 Reference Gallery",
    "💡 Pattern Cheatsheet",
    "🎮 Create PPTX",
])


# ════════════════════════════════════════════════════════════════
# TAB 1 — VERIFY
# ════════════════════════════════════════════════════════════════
with t_verify:
    st.subheader("Upload your PPTX")
    uploaded = st.file_uploader("Drop a .pptx file", type=["pptx"],
                                label_visibility="collapsed")

    if uploaded:
        if not HAS_PPTX and not HAS_MARKITDOWN:
            st.error("Run: `pip install python-pptx markitdown`")
        else:
            st.success(f"📁 **{uploaded.name}** ({uploaded.size/1024:.0f} KB)")
            if st.button("🔍 Analyze with Action-First Framework", use_container_width=True):
                try:
                    with st.spinner("Extracting slide content…"):
                        slides = extract_slide_content(uploaded)
                    if not slides:
                        st.warning("No slides detected.")
                    else:
                        st.info(f"Found **{len(slides)} slides**. Analyzing…")
                        result_text = analyze_slides_with_claude(slides, uploaded.name)
                        st.session_state.analysis_result   = result_text
                        st.session_state.analysis_filename = uploaded.name
                        _save_analysis(st.session_state.auth_email, uploaded.name,
                                       {"text": result_text, "score": "—", "slides": len(slides)})
                        st.rerun()
                except Exception as e:
                    st.error(f"Analysis failed: {e}")

    if st.session_state.analysis_result:
        st.divider()
        st.subheader(f"📝 Analysis: {st.session_state.analysis_filename}")
        r_tab1, r_tab2 = st.tabs(["📊 Full Report", "📄 Raw Text"])
        with r_tab1:
            st.markdown(st.session_state.analysis_result)
        with r_tab2:
            st.code(st.session_state.analysis_result, language=None)
            st.download_button("⬇️ Download (.txt)", data=st.session_state.analysis_result,
                file_name=f"review_{st.session_state.analysis_filename}.txt",
                mime="text/plain", use_container_width=True)
        if st.button("🔄 Analyze a new file", use_container_width=True):
            st.session_state.analysis_result = None
            st.session_state.analysis_filename = None; st.rerun()

    elif not uploaded:
        st.markdown("""<div style="text-align:center;padding:3rem 1rem;color:#475569">
          <div style="font-size:3rem;margin-bottom:1rem">📂</div>
          <div style="font-size:1.1rem;font-weight:600;color:#94a3b8">Upload a PPTX to get started</div>
          <div style="font-size:.85rem;margin-top:.5rem">Slides checked against Action-First framework and<br>
          benchmarked against BCG & McKinsey reference examples.</div></div>""",
          unsafe_allow_html=True)


# ════════════════════════════════════════════════════════════════
# TAB 2 — REFERENCE GALLERY
# ════════════════════════════════════════════════════════════════
with t_gallery:
    st.subheader("Reference Gallery — What Good Looks Like")
    st.caption(
        "Annotated slides from two real consulting decks: **BCG — Melbourne as a Global Cultural Destination** "
        "and **McKinsey/IACPM — Emerging GenAI Use Cases in Credit**. "
        "These same patterns are embedded in Claude's analysis prompt."
    )

    # Filter buttons
    fc1, fc2, fc3, _ = st.columns([1, 1, 1, 5])
    for btn, firm in [(fc1, "All"), (fc2, "BCG"), (fc3, "McKinsey")]:
        active = st.session_state.gallery_filter == firm
        label  = f"{'✓ ' if active else ''}{firm}"
        if btn.button(label, key=f"gf_{firm}", use_container_width=True):
            st.session_state.gallery_filter = firm; st.rerun()

    st.markdown("")

    filtered = [s for s in REFERENCE_SLIDES
                if st.session_state.gallery_filter == "All"
                or s["firm"] == st.session_state.gallery_filter]

    for i in range(0, len(filtered), 2):
        cols = st.columns(2, gap="medium")
        for j, col in enumerate(cols):
            if i + j >= len(filtered): break
            s = filtered[i + j]
            chip  = "chip-bcg" if s["firm"] == "BCG" else "chip-mckinsey"
            pts   = "".join(
                f'<div class="ref-point"><span class="ref-icon">✓</span><span>{p}</span></div>'
                for p in s["what_works"])
            col.markdown(f"""<div class="ref-card">
              <div class="ref-deck"><span class="{chip}">{s['firm']}</span>Slide {s['slide_num']}</div>
              <div class="ref-title">"{s['title']}"</div>
              <span class="ref-badge">✓ PASS</span>
              {pts}
              <div class="ref-annotation">{s['annotation']}</div>
              <div class="ref-pattern">Pattern: {s['pattern']}</div>
            </div>""", unsafe_allow_html=True)

    st.markdown("---")
    st.caption(f"Showing {len(filtered)} of {len(REFERENCE_SLIDES)} reference slides.")


# ════════════════════════════════════════════════════════════════
# TAB 3 — PATTERN CHEATSHEET
# ════════════════════════════════════════════════════════════════
with t_patterns:
    st.subheader("Pattern Cheatsheet")
    st.caption("Quick reference for the consulting patterns extracted from the gallery. "
               "Claude uses these same pattern names when rewriting your titles and takeaways.")

    st.markdown("### 🏷️ Title Patterns")

    title_patterns = [
        {
            "name": "Contrast title",
            "source": "BCG Slide 22",
            "desc": "State two competing facts separated by a comma. The audience has the verdict before reading the body.",
            "bad":  "SWOT Analysis",
            "good": "Melbourne has strengths in offer and governance, underperforms in brand and marketing",
        },
        {
            "name": "Finding + implication",
            "source": "BCG Slide 35",
            "desc": "Data point THEN business implication — separated by a dash. Both halves belong in the title.",
            "bad":  "Visitor Recommendation Data",
            "good": "Travellers more likely to recommend Melbourne after visiting — suggests marketing could boost visitation",
        },
        {
            "name": "Active observation with examples",
            "source": "McKinsey Slide 6",
            "desc": "Observation verb + specific examples in parentheses. Conclusion AND supporting specifics in one line.",
            "bad":  "GenAI Use Case Update",
            "good": "Institutions are prioritizing use cases like underwriting (synthesizing, drafting memo) and portfolio monitoring (early warning)",
        },
        {
            "name": "Ranked finding",
            "source": "McKinsey Slide 12",
            "desc": "The #1 result belongs in the headline. Don't bury the ranking in bullets.",
            "bad":  "Factor Prioritization Overview",
            "good": "Institutions prioritize productivity improvement as the most important factor when initiating GenAI use cases",
        },
        {
            "name": "Majority qualifier",
            "source": "McKinsey Slide 13",
            "desc": "'Majority' is more credible than 'many' (too vague) and less awkward than '52%' (too precise for a headline).",
            "bad":  "Leadership Support for GenAI",
            "good": "Leadership at majority of the institutions are positioning GenAI as a priority",
        },
    ]

    for p in title_patterns:
        with st.expander(f"**{p['name']}** — *{p['source']}*"):
            st.markdown(f"**What it does:** {p['desc']}")
            bc, gc = st.columns(2)
            bc.markdown(
                f"<div style='background:#7f1d1d;padding:.7rem;border-radius:8px;font-size:.82rem;color:#fca5a5'>"
                f"❌ <strong>Weak:</strong> {p['bad']}</div>", unsafe_allow_html=True)
            gc.markdown(
                f"<div style='background:#064e3b;padding:.7rem;border-radius:8px;font-size:.82rem;color:#6ee7b7'>"
                f"✓ <strong>Strong:</strong> {p['good']}</div>", unsafe_allow_html=True)

    st.markdown("### 📦 Takeaway Block Patterns")

    takeaway_patterns = [
        {
            "name": "Pre-chart callout boxes",
            "source": "McKinsey Slide 12",
            "desc": "Summary boxes ABOVE the chart. Each box: rank statement + percentage + label. Verdict before data.",
            "example": "47% ranked productivity #1 | 44% ranked business needs #2 | 25% ranked regulatory #3 | 50% ranked ROI last",
        },
        {
            "name": "Pull-quote box",
            "source": "BCG Slide 35",
            "desc": "Stakeholder quote in shaded box at bottom-right. The shading IS the Takeaway Block signal.",
            "example": "'Melbourne lacks a clear value proposition… the proposition should be holistic' — Thought Leader",
        },
        {
            "name": "Action Subtitle",
            "source": "BCG Slide 9",
            "desc": "Second line directly below the main title that states a sub-claim. Body charts prove the sub-claim.",
            "example": "Title: 'Creative industries increasingly important to the economy' | Sub: 'Workers generate slightly more GVA than the average Victorian worker'",
        },
        {
            "name": "Three-tier classification",
            "source": "McKinsey Slide 13",
            "desc": "Classify findings into 3 tiers. Write a self-assessment description for each so readers can place themselves.",
            "example": "Priority (52%) — 'Senior leadership promotes GenAI as a priority through investments' | Interested (39%) | Not a priority (9%)",
        },
    ]

    for p in takeaway_patterns:
        with st.expander(f"**{p['name']}** — *{p['source']}*"):
            st.markdown(f"**What it does:** {p['desc']}")
            st.markdown(
                f"<div style='background:#1e293b;border-left:3px solid #4f46e5;padding:.7rem;"
                f"border-radius:0 8px 8px 0;font-size:.8rem;color:#94a3b8;margin-top:.5rem'>"
                f"<strong>Example:</strong> {p['example']}</div>", unsafe_allow_html=True)

    st.markdown("---")
    st.caption("All patterns above are injected into Claude's analysis prompt. "
               "When your slide resembles one of these patterns, Claude will name it.")

# Footer
st.markdown("---")
st.caption("Slide Verifier · Action-First Framework · Benchmarked against BCG & McKinsey · Internal use only")

# ════════════════════════════════════════════════════════════════
# TAB 4 — CREATE PPTX
# Imports the creator engine from pptx_creator.py + pptxruns.py
# which must live in the same directory as pptx_verifier.py.
# Auth is already handled above — OWNER is the signed-in email.
# ════════════════════════════════════════════════════════════════
with t_create:
    _creator_ok = False
    try:
        import sys as _sys, os as _os
        _here = _os.path.dirname(_os.path.abspath(__file__))
        if _here not in _sys.path:
            _sys.path.insert(0, _here)

        from pptx_creator import (
            run_pipeline        as _run_pipeline,
            _render_plan_modal  as _render_plan_modal,
            _render_log         as _render_log,
            _load_project       as _load_project,
            _save_project       as _save_project,
            _clear_project      as _clear_project,
            _DEFAULT_TEMPLATE_BYTES,
            PURPOSE_PRESETS,
            THEME_PRESETS,
        )
        from storage_pptx import (
            get_projects, project_exists, create_project,
            rename_project, delete_project,
        )
        _creator_ok = True
    except ImportError:
        pass

    if not _creator_ok:
        st.markdown("""
        <div style='text-align:center;padding:3rem 1rem;color:#475569'>
          <div style='font-size:3rem;margin-bottom:1rem'>🔧</div>
          <div style='font-size:1.1rem;font-weight:600;color:#94a3b8'>Creator not available</div>
          <div style='font-size:.85rem;margin-top:.5rem'>
            Place <code>pptx_creator.py</code>, <code>pptxruns.py</code>, and
            <code>storage_pptx.py</code> in the same directory as this file,
            then restart the app.
          </div>
        </div>
        """, unsafe_allow_html=True)
    else:
        import re as _re
        OWNER = st.session_state.auth_email

        st.markdown("""
        <style>
        .cr-section{font-size:.68rem;font-weight:700;letter-spacing:.1em;
            text-transform:uppercase;color:#00CCFF;margin:.9rem 0 .3rem}
        .cr-step-done{font-size:.78rem;color:#22DD88;padding:.12rem 0}
        .cr-step-pend{font-size:.78rem;color:#6080A8;padding:.12rem 0}
        .cr-status{background:#0A1832;border:1px solid #0D2860;border-top:2px solid #0055AA;
            border-radius:8px;padding:1.25rem 1.5rem;margin-top:.5rem}
        .cr-status-lbl{font-size:.65rem;font-weight:700;letter-spacing:.12em;
            text-transform:uppercase;color:#00CCFF;margin-bottom:.5rem}
        </style>
        """, unsafe_allow_html=True)

        # ── Project management ──────────────────────────────────────────────────
        _projects   = get_projects(OWNER)
        _proj_names = [p["name"] for p in _projects]
        _active     = st.session_state.get("active_project", "")

        with st.expander("📁 Project", expanded=not bool(_active)):
            _pcol1, _pcol2 = st.columns([3, 1])
            with _pcol1:
                _opts    = ["— select —"] + _proj_names
                _cur_idx = (_opts.index(_active) if _active in _opts else 0)
                _sel     = st.selectbox("Project", _opts, index=_cur_idx,
                                        label_visibility="hidden", key="cr_proj_sel")
                if _sel != "— select —" and _sel != _active:
                    _load_project(OWNER, _sel)
                    st.rerun()

            with _pcol2:
                _new_proj = st.text_input("New name", placeholder="e.g. Q3 Deck",
                                          label_visibility="hidden", key="cr_new_proj")
                if st.button("＋ Create", key="cr_create_proj", use_container_width=True):
                    if _new_proj.strip():
                        if project_exists(OWNER, _new_proj.strip()):
                            st.error("Name already exists.")
                        else:
                            create_project(OWNER, _new_proj.strip())
                            _clear_project()
                            st.session_state["active_project"] = _new_proj.strip()
                            st.rerun()

            if _active:
                _pa, _pb, _pc = st.columns([3, 2, 1])
                with _pa:
                    if st.button("💾 Save", key="cr_save", use_container_width=True):
                        _save_project(OWNER, _active)
                        st.toast(f'Saved "{_active}"', icon="✅")
                with _pb:
                    if st.button("🔄 Reset output", key="cr_reset", use_container_width=True):
                        for _k in ["pptx_bytes","pptx_filename","plan_slide_data",
                                   "plan_mode_active","plan_chat","plan_slide_history","pipeline_steps"]:
                            st.session_state.pop(_k, None)
                        st.rerun()
                with _pc:
                    if st.button("🗑", key="cr_del", help="Delete", use_container_width=True):
                        st.session_state["cr_confirm_del"] = True

                if st.session_state.get("cr_confirm_del"):
                    st.warning(f'Delete "{_active}"?')
                    _dy, _dn = st.columns(2)
                    if _dy.button("Yes, delete", key="cr_del_yes"):
                        delete_project(OWNER, _active)
                        _clear_project()
                        st.session_state.pop("cr_confirm_del", None)
                        st.rerun()
                    if _dn.button("Cancel", key="cr_del_no"):
                        st.session_state.pop("cr_confirm_del", None)
                        st.rerun()

        if not _active:
            st.markdown("""
            <div style='text-align:center;padding:2rem 1rem'>
              <div style='font-size:2rem;margin-bottom:.5rem'>🎮</div>
              <div style='font-size:1rem;font-weight:600;color:#94a3b8'>
                Select or create a project above to get started
              </div>
            </div>
            """, unsafe_allow_html=True)
        else:
            st.markdown(f"<div style='font-size:.75rem;color:#4A6A9A;margin-bottom:.75rem'>"
                        f"Project: <b style='color:#D0E4FF'>{_active}</b></div>",
                        unsafe_allow_html=True)

            # ── Options ─────────────────────────────────────────────────────────
            with st.expander("⚙️ Options", expanded=False):
                _cr_oc1, _cr_oc2, _cr_oc3 = st.columns(3)
                with _cr_oc1:
                    _cr_model = st.selectbox("Model",
                        ["claude-sonnet-4-6","claude-opus-4-6","claude-haiku-4-5-20251001"],
                        key="cr_model")
                    _cr_web = st.checkbox("Web research",
                                          value=st.session_state.get("proj_web_research", True),
                                          key="cr_web")
                    st.session_state["proj_web_research"] = _cr_web
                with _cr_oc2:
                    _cr_slides = st.slider("Target slides", 6, 25, 12, key="cr_slides")
                    _cr_theme_name = st.selectbox("Theme", list(THEME_PRESETS.keys()), key="cr_theme")
                    _cr_theme = THEME_PRESETS[_cr_theme_name]
                with _cr_oc3:
                    _cr_template = st.file_uploader("Custom .pptx template",
                                                    type=["pptx"], key="cr_template")
                    _cr_pipeline = st.session_state.get("pipeline_steps",{})
                    for _k, _lbl in {"upload":"Docs","extract":"Extract",
                                     "research":"Research","analyze":"Analyze","generate":"Generate"}.items():
                        _d = _cr_pipeline.get(_k, False)
                        st.markdown(f'<div class="{"cr-step-done" if _d else "cr-step-pend"}">'
                                    f'{"✓" if _d else "○"} {_lbl}</div>', unsafe_allow_html=True)

            def _get_tpl():
                if _cr_template:
                    _cr_template.seek(0)
                    return _cr_template.read()
                return st.session_state.get("saved_template_bytes") or _DEFAULT_TEMPLATE_BYTES

            # ── Inner tabs ───────────────────────────────────────────────────────
            _cr_t1, _cr_t2, _cr_t3, _cr_t4 = st.tabs([
                "💬 Guided Build",
                "📊 Create Presentation",
                "✏️ Edit Outline",
                "📄 PDF → PPTX",
            ])

            # ══ GUIDED BUILD ═══════════════════════════════════════════════════
            with _cr_t1:
                if "guided_messages" not in st.session_state:
                    st.session_state["guided_messages"] = []
                if "guided_ready" not in st.session_state:
                    st.session_state["guided_ready"] = False
                if "guided_params" not in st.session_state:
                    st.session_state["guided_params"] = {}

                _gl, _gr = st.columns([1, 1], gap="large")
                with _gl:
                    _gchat = st.container(height=340)
                    with _gchat:
                        if not st.session_state["guided_messages"]:
                            st.caption("👋 Tell me about the presentation you need.")
                        for _m in st.session_state["guided_messages"]:
                            with st.chat_message(_m["role"]):
                                st.markdown(_m["content"])

                    _gi = st.chat_input("Tell me about your presentation…", key="cr_guided_input")
                    if _gi:
                        st.session_state["guided_messages"].append({"role":"user","content":_gi})
                        _gsys = """You are a presentation planning assistant.
Rules: professional, direct, concise. No filler. Ask 1-2 questions at a time.
Collect: topic, purpose, audience, slide count (default 10), focus areas.
Once you have topic+purpose+audience, output the plan then the signal line.

Format:
**Topic:** [value]
**Purpose:** [value]
**Audience:** [value]
**Slides:** [value]
**Focus:** [core goal]

READY_TO_GENERATE: {"topic":"...","purpose":"...","industry":"...","audience":"...","question":"...","slide_count":10}

Populate all fields from the conversation. Empty string for anything not discussed. No markdown fences."""
                        try:
                            from pptxruns import _make_anthropic_client as _mac
                            _grc = _mac().messages.create(
                                model="claude-sonnet-4-6", max_tokens=600,
                                system=_gsys,
                                messages=[{"role":m["role"],"content":m["content"]}
                                          for m in st.session_state["guided_messages"]],
                            )
                            _rep = _grc.content[0].text if _grc.content else "Try again."
                        except Exception as _ge2:
                            _rep = f"(Error: {_ge2})"

                        _clean = _rep
                        if "READY_TO_GENERATE:" in _rep:
                            try:
                                _jl = _rep.split("READY_TO_GENERATE:")[1].strip().split("\n")[0]
                                _pp = json.loads(_jl)
                                st.session_state.update(guided_params=_pp, guided_ready=True)
                                _clean = _rep.split("READY_TO_GENERATE:")[0].strip() or (
                                    f"**Topic:** {_pp.get('topic','')}\n"
                                    f"**Purpose:** {_pp.get('purpose','')}\n"
                                    f"**Audience:** {_pp.get('audience','')}\n"
                                    f"**Slides:** {_pp.get('slide_count',10)}\n\n"
                                    "Click **Generate** on the right when ready.")
                            except Exception:
                                st.session_state["guided_ready"] = False

                        st.session_state["guided_messages"].append({"role":"assistant","content":_clean})
                        st.rerun()

                    if st.button("🔄 Start over", key="cr_guided_reset", use_container_width=True):
                        for _kk in ["guided_messages","guided_ready","guided_params",
                                    "guided_pptx_bytes","guided_pptx_filename"]:
                            st.session_state.pop(_kk, None)
                        st.rerun()

                with _gr:
                    _gp   = st.session_state.get("guided_params",{})
                    _grdy = st.session_state.get("guided_ready", False)
                    if _grdy and _gp:
                        st.success(f"**{_gp.get('topic','Topic')}**  \n"
                                   f"{_gp.get('purpose','')} · {_gp.get('audience','')} · "
                                   f"{_gp.get('slide_count',10)} slides")
                    else:
                        st.markdown('<div class="cr-status"><div class="cr-status-lbl">Waiting for plan</div>'
                                    '<div style="color:#6080A8;font-size:.82rem;line-height:1.9">'
                                    'Chat on the left. Once Claude has enough info the Generate button unlocks.'
                                    '</div></div>', unsafe_allow_html=True)

                    _gbtn = st.button("⚡ Generate presentation", use_container_width=True,
                                      type="primary", disabled=not _grdy, key="cr_guided_gen")

                _g_log  = st.empty()
                _g_plan = st.container()
                _g_dl   = st.empty()

                if st.session_state.get("plan_mode_active") and not _gbtn:
                    with _g_plan:
                        _render_plan_modal(_get_tpl(), ns="cr_guided")

                if st.session_state.get("guided_pptx_bytes") and not _gbtn:
                    _g_dl.download_button("⬇️ Download PPTX",
                        data=st.session_state["guided_pptx_bytes"],
                        file_name=st.session_state.get("guided_pptx_filename","presentation.pptx"),
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True)

                if _gbtn and _grdy and _gp:
                    _gt = _get_tpl()
                    _gtopic = _gp.get("topic","")
                    st.session_state.update(proj_topic=_gtopic,
                        proj_purpose=_gp.get("purpose",""),
                        proj_industry=_gp.get("industry",""),
                        proj_audience=_gp.get("audience",""))
                    _glogs = []
                    for _ev in _run_pipeline(
                        model=_cr_model, uploaded_files=[],
                        topic=_gtopic, purpose=_gp.get("purpose","General / Other"),
                        industry=_gp.get("industry",""), audience=_gp.get("audience","General audience"),
                        question=_gp.get("question",""), web_search_en=_cr_web,
                        slide_count=int(_gp.get("slide_count",10)),
                        theme=_cr_theme, template_bytes=_gt, plan_mode=True,
                    ):
                        _et = _ev[0]
                        if _et in ("log","spinner"):
                            if _et=="spinner":
                                if _glogs and _glogs[-1][0]=="spinner": _glogs[-1]=("log",_glogs[-1][1])
                                _glogs.append(("spinner",_ev[1]))
                            else:
                                if _glogs and _glogs[-1][0]=="spinner": _glogs[-1]=("log",_glogs[-1][1])
                                _glogs.append(("log",_ev[1]))
                            _g_log.markdown(_render_log(_glogs), unsafe_allow_html=True)
                        elif _et=="sources": st.session_state["research_sources"]=_ev[1]
                        elif _et=="plan_ready":
                            st.session_state["plan_slide_data"]=_ev[1]
                            st.session_state["plan_mode_active"]=True
                            _save_project(OWNER,_active)
                        elif _et=="pptx_bytes_out":
                            _slg=_re.sub(r"[^a-zA-Z0-9]+","_",_gtopic)[:50]
                            st.session_state["guided_pptx_bytes"]=_ev[1]
                            st.session_state["guided_pptx_filename"]=f"Presentation_{_slg}.pptx"
                            st.session_state["pptx_bytes"]=_ev[1]
                            _save_project(OWNER,_active)
                        elif _et=="error": st.error(_ev[1],icon="🚨"); break

                    if st.session_state.get("plan_mode_active"):
                        with _g_plan:
                            _render_plan_modal(_get_tpl(), ns="cr_guided")

            # ══ CREATE PRESENTATION ════════════════════════════════════════════
            with _cr_t2:
                _fc1, _fc2 = st.columns([1,1], gap="large")
                with _fc1:
                    st.markdown('<div class="cr-section">Topic & documents</div>', unsafe_allow_html=True)
                    _topic    = st.text_input("Topic / title",
                        value=st.session_state.get("proj_topic",""),
                        placeholder="e.g. Q3 Market Analysis, Product Launch…", key="cr_topic")
                    _purpose  = st.selectbox("Purpose", list(PURPOSE_PRESETS.keys()), key="cr_purpose")
                    _industry = st.text_input("Industry / context (optional)",
                        value=st.session_state.get("proj_industry",""),
                        placeholder="e.g. Healthcare, Gaming…", key="cr_industry")
                    _audience = st.text_input("Audience",
                        value=st.session_state.get("proj_audience",""),
                        placeholder="e.g. Executive team, Board…", key="cr_audience")
                    _question = st.text_area("Business question / goal (optional)",
                        placeholder="What question should this deck answer?",
                        height=72, key="cr_question")
                    _uploads  = st.file_uploader("Supporting documents",
                        type=["pdf","docx","txt","csv","xlsx"],
                        accept_multiple_files=True, key="cr_uploads")

                with _fc2:
                    st.markdown('<div class="cr-section">Output</div>', unsafe_allow_html=True)
                    _out_area = st.empty()
                    _dl_area  = st.empty()

                    if st.session_state.get("plan_mode_active") and not st.session_state.get("cr_run_clicked"):
                        with _out_area.container():
                            _render_plan_modal(_get_tpl(), ns="cr_main")
                    if st.session_state.get("pptx_bytes") and not st.session_state.get("cr_run_clicked"):
                        _dl_area.download_button("⬇️ Download previous PPTX",
                            data=st.session_state["pptx_bytes"],
                            file_name=st.session_state.get("pptx_filename","presentation.pptx"),
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True)

                _run_btn = st.button("⚡ Generate presentation", use_container_width=True,
                                     type="primary", key="cr_run_btn")
                st.session_state["cr_run_clicked"] = _run_btn

                if _run_btn:
                    if not _topic.strip():
                        st.error("Please enter a topic.")
                    else:
                        st.session_state.update(
                            proj_topic=_topic, proj_purpose=_purpose,
                            proj_industry=_industry, proj_audience=_audience,
                            project_doc_names=[f.name for f in (_uploads or [])])
                        _tpl2 = _get_tpl()
                        _pipe = {"upload":bool(_uploads),"extract":False,
                                 "research":False,"analyze":False,"generate":False}
                        st.session_state["pipeline_steps"] = _pipe
                        _logs2 = []
                        with _out_area.container():
                            _la2 = st.empty()
                            try:
                                for _ev2 in _run_pipeline(
                                    model=_cr_model, uploaded_files=_uploads or [],
                                    topic=_topic, purpose=_purpose,
                                    industry=_industry, audience=_audience,
                                    question=_question, web_search_en=_cr_web,
                                    slide_count=_cr_slides, theme=_cr_theme,
                                    template_bytes=_tpl2, plan_mode=True,
                                ):
                                    _et2 = _ev2[0]
                                    if _et2 in ("log","spinner"):
                                        if _et2=="spinner":
                                            if _logs2 and _logs2[-1][0]=="spinner":
                                                _logs2[-1]=("log",_logs2[-1][1])
                                            _logs2.append(("spinner",_ev2[1]))
                                        else:
                                            if _logs2 and _logs2[-1][0]=="spinner":
                                                _logs2[-1]=("log",_logs2[-1][1])
                                            _logs2.append(("log",_ev2[1]))
                                        _la2.markdown(_render_log(_logs2),unsafe_allow_html=True)
                                    elif _et2=="step_done":
                                        _pipe[_ev2[1]]=True
                                        st.session_state["pipeline_steps"]=_pipe
                                    elif _et2=="sources":
                                        st.session_state["research_sources"]=_ev2[1]
                                    elif _et2=="plan_ready":
                                        st.session_state["plan_slide_data"]=_ev2[1]
                                        st.session_state["plan_mode_active"]=True
                                        _save_project(OWNER,_active)
                                        st.toast(f'Outline saved to "{_active}"',icon="📋")
                                    elif _et2=="pptx_bytes_out":
                                        _slg2=_re.sub(r"[^a-zA-Z0-9]+","_",_topic)[:50]
                                        st.session_state["pptx_bytes"]=_ev2[1]
                                        st.session_state["pptx_filename"]=f"Presentation_{_slg2}.pptx"
                                        _save_project(OWNER,_active)
                                        st.toast(f'Saved "{_active}"',icon="💾")
                                    elif _et2=="error":
                                        st.error(_ev2[1],icon="🚨"); break
                            except Exception as _ex2:
                                st.error(f"Error: {_ex2}")

                        if st.session_state.get("plan_mode_active"):
                            with _out_area.container():
                                _render_plan_modal(_get_tpl(), ns="cr_main")
                        if st.session_state.get("pptx_bytes"):
                            st.success("Presentation ready!")
                            _dl_area.download_button("⬇️ Download PPTX",
                                data=st.session_state["pptx_bytes"],
                                file_name=st.session_state.get("pptx_filename","presentation.pptx"),
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True)
                            _src2=st.session_state.get("research_sources",[])
                            if _src2:
                                with st.expander(f"🔗 {len(_src2)} source(s)"):
                                    for _s2 in _src2:
                                        st.markdown(f"- [{_s2['title']}]({_s2['url']})")

            # ══ EDIT OUTLINE ══════════════════════════════════════════════════
            with _cr_t3:
                if st.session_state.get("plan_slide_data"):
                    _render_plan_modal(_get_tpl(), ns="cr_edit")
                elif st.session_state.get("pptx_bytes"):
                    st.info("Outline not available — the PPTX was generated without an editable outline. "
                            "Generate a new presentation to access outline editing.")
                    st.download_button("⬇️ Download current PPTX",
                        data=st.session_state["pptx_bytes"],
                        file_name=st.session_state.get("pptx_filename","presentation.pptx"),
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True)
                else:
                    st.info("Generate a presentation first — then come back here to edit the outline.")

            # ══ PDF → PPTX ════════════════════════════════════════════════════
            with _cr_t4:
                st.caption("Upload any PDF of slides — even rasterised exports from NotebookLM, "
                           "Canva, etc. Claude reads each page with vision and reconstructs it "
                           "as fully editable PowerPoint shapes.")

                try:
                    from pdf_to_pptx import pdf_to_editable_pptx as _p2p
                    _HAS_P2P = True
                except ImportError:
                    _HAS_P2P = False

                _pd1, _pd2 = st.columns([1,1], gap="large")
                with _pd1:
                    _pup = st.file_uploader("Upload PDF", type=["pdf"],
                                            label_visibility="hidden", key="cr_pdf_up")
                    if _pup:
                        try:
                            import pypdf as _pp4
                            _np4 = len(_pp4.PdfReader(io.BytesIO(_pup.read())).pages)
                            _pup.seek(0)
                            st.caption(f"📄 {_np4} page{'s' if _np4!=1 else ''} detected")
                        except Exception:
                            pass
                    _pmd = st.selectbox("Vision model",
                        ["claude-opus-4-5","claude-sonnet-4-5"], key="cr_pdf_model",
                        help="Opus = more accurate; Sonnet = faster.")
                    _pdpi = st.select_slider("Quality",
                        options=[96,120,150,200], value=150, key="cr_pdf_dpi")
                    _pbtn = st.button("⚡ Convert to editable PPTX", use_container_width=True,
                                     type="primary", key="cr_pdf_btn", disabled=not bool(_pup))

                with _pd2:
                    _pout = st.empty()
                    if st.session_state.get("cr_pdf_out") and not _pbtn:
                        _pout.download_button("⬇️ Download editable PPTX",
                            data=st.session_state["cr_pdf_out"],
                            file_name=st.session_state.get("cr_pdf_fname","converted.pptx"),
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True)

                if _pbtn and _pup:
                    if not _HAS_P2P:
                        st.error("pdf_to_pptx.py not found — place it in the same directory.")
                    else:
                        _akey = st.secrets.get("ANTHROPIC_API_KEY","")
                        if not _akey:
                            st.error("ANTHROPIC_API_KEY not set.")
                        else:
                            _pup.seek(0); _praw = _pup.read()
                            _pprog = _pout.progress(0,"Starting…")
                            _plog2 = st.empty(); _pll = []
                            def _plg(m):
                                _pll.append(m)
                                _plog2.markdown(
                                    "<div style='background:#0f172a;border-radius:6px;"
                                    "padding:.6rem 1rem;font-size:.78rem;color:#94a3b8;"
                                    "font-family:monospace;max-height:200px;overflow-y:auto'>"
                                    +"".join(f"<div>{l}</div>" for l in _pll[-15:])+"</div>",
                                    unsafe_allow_html=True)
                            try:
                                import pypdf as _pp5
                                _npg2=len(_pp5.PdfReader(io.BytesIO(_praw)).pages)
                                _plg(f"📄 {_npg2} page(s) · sending to {_pmd}…")
                                _pout2,_perrs=_p2p(_praw,api_key=_akey,model=_pmd,dpi=_pdpi,
                                    progress_cb=lambda f:_pprog.progress(min(f,0.99),
                                        text=f"Page {max(1,round(f*_npg2))} of {_npg2}…"))
                                _pfn=(_pup.name.rsplit(".",1)[0] if "." in _pup.name else _pup.name)+"_editable.pptx"
                                st.session_state["cr_pdf_out"]=_pout2
                                st.session_state["cr_pdf_fname"]=_pfn
                                for _pe in (_perrs or []): _plg(f"⚠️ {_pe}")
                                _pprog.progress(1.0,text="Done!")
                                _plg(f"🎉 Complete — {len(_pout2)//1024} KB")
                                _pout.download_button("⬇️ Download editable PPTX",
                                    data=_pout2, file_name=_pfn,
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    use_container_width=True)
                                st.rerun()
                            except Exception as _pe2:
                                st.error(f"Conversion failed: {_pe2}")

# Footer
st.markdown("---")
st.caption("Slide Verifier · Action-First Framework · Benchmarked against BCG & McKinsey · Internal use only")